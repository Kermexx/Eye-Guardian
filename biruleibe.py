import os
import re
import fitz#PymuPDF
from google.cloud import vision
from docx import Document  # Para lidar com arquivos DOCX
from pptx import Presentation
import openpyxl  # Para lidar com arquivos XLSX
import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox
from tkinter import ttk
from tkinter import Canvas, Entry, Text, Button, PhotoImage
from pathlib import Path
import sys
from customtkinter import *
from PIL import Image, ImageTk

religioes = [
    'Cristianismo',
    'Cristão',
    'Cristã',
    'Islamismo',
    'Islâmico',
    'Islâmica',
    'Hinduísmo',
    'Hindu',
    'Hinduísta',
    'Budismo',
    'Budista',
    'Sikhismo',
    'Sikh',
    'Judaísmo',
    'Judeu',
    'Judaica',
    'Bahaí',
    'Bahaísta',
    'Jainismo',
    'Jainista',
    'Espiritismo',
    'Espírita',
    'Ateísmo',
    'Ateu',
    'Ateia'
]

cores_etnias = [
    'Branco',
    'Negro',
    'Pardo',
    'Indígena',
    'Amarelo',
    'Asiático',
    'Outro / Não Declarado',
]

def extract_info_by_pattern(pattern, text, info_type, results):
    matches = re.findall(pattern, text, flags=re.IGNORECASE)
    if matches:
        results.extend([(info_type, match) for match in matches])
    return results

def format_rg(rg):
    # Remove pontos e hífens do RG e adiciona o hífen no formato desejado
    rg_limpo = re.sub(r'[^\d]', '', rg)
    return f'{rg_limpo[:-1]}.{rg_limpo[-1]}'

# Modifica a função extract_sensitive_info_from_xlsx para incluir formatação de RG
def extract_sensitive_info_from_xlsx(xlsx_path, results):
    wb = openpyxl.load_workbook(xlsx_path)

    sensitive_info = []

    # Itera sobre todas as folhas no arquivo Excel
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Itera sobre todas as células na folha
        for row in sheet.iter_rows(min_row=1, max_col=sheet.max_column, max_row=sheet.max_row, values_only=True):
            for cell_value in row:
                if cell_value:
                    # Aplica expressões regulares para encontrar informações sensíveis
                    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})|\d{8,9}-\d{1,2}', cell_value)
                    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', str(cell_value))
                    matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', str(cell_value))
                    matches_email = re.findall(r'\S+@\S+', str(cell_value))
                    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', str(cell_value))
                    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', cell_value)
                    valid_telefones = []


                    for telefone in matches_telefone:
                        numero_limpo = re.sub(r'[^\d]', '', telefone)
                        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                            valid_telefones.append(telefone)


                    for rg in matches_rg:
                        rg_formatado = format_rg(rg)
                        rg_in_cpf = any(rg_formatado in cpf for cpf in matches_cpf)
                        if not rg_in_cpf:
                            sensitive_info.append(('RG', rg_formatado))

                    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
                    sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
                    sensitive_info.extend([('Email', email) for email in matches_email])
                    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
                    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
                    # Extrai informações sobre religiões
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', cell_value, 'Religião',sensitive_info)
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', cell_value,'Cor/Etnia', sensitive_info)
    if sensitive_info:
        results[xlsx_path] = results.get(xlsx_path, [])
        results[xlsx_path].extend(sensitive_info)

    return results

def process_directory_with_xlsx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos XLSX
            if filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)

def extract_sensitive_info_from_pptx(pptx_path_or_text, results):
    if os.path.isfile(pptx_path_or_text):  # Verifica se é um caminho de arquivo
        presentation = Presentation(pptx_path_or_text)
        text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

    else:
        text = pptx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)
    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pptx_path_or_text] = results.get(pptx_path_or_text, [])
        results[pptx_path_or_text].extend(sensitive_info)

    return results

# Modifica a função process_directory para incluir arquivos PPTX
def process_directory_with_pptx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos PPTX
            if filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)

# Função para extrair informações sensíveis de um arquivo PDF
def extract_sensitive_info_from_pdf(pdf_path, results):
    # Inicializa o documento PDF
    doc = fitz.open(pdf_path)

    sensitive_info = []

    # Itera sobre as páginas do documento PDF
    for page_number in range(doc.page_count):
        page = doc[page_number]

        # Extrai texto da página
        text = page.get_text()

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
        matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', text)
        matches_email = re.findall(r'\S+@\S+', text)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

        valid_telefones = []

        # Filtra e formata números de telefone válidos
        for telefone in matches_telefone:
            numero_limpo = re.sub(r'[^\d]', '', telefone)
            if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                valid_telefones.append(telefone)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
            if not rg_in_cpf:
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
        # Extrai informações sobre religiões
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião',sensitive_info)
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pdf_path] = results.get(pdf_path, [])
        results[pdf_path].extend(sensitive_info)

    return results

# Função para extrair informações sensíveis de uma imagem
def extract_sensitive_info_from_image(image_path, results):
    # Inicializa o cliente Google Cloud Vision
    client = vision.ImageAnnotatorClient()

    # Lê o conteúdo da imagem
    with open(image_path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    # Envia a imagem para análise de texto
    response = client.text_detection(image=image)

    # Extrai texto identificado na imagem
    texts = response.text_annotations

    sensitive_info = []

    # Itera sobre os textos identificados
    for text in texts:
        text = text.description

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2})', text)
        matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', text)
        matches_email = re.findall(r'\S+@\S+', text)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)
        valid_telefones = []

        # Filtra e formata números de telefone válidos
        for telefone in matches_telefone:
            numero_limpo = re.sub(r'[^\d]', '', telefone)
            if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                valid_telefones.append(telefone)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
            if not rg_in_cpf:
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
        # Extrai informações sobre religiões
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião',sensitive_info)
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[image_path] = results.get(image_path, [])
        results[image_path].extend(sensitive_info)

    return results


# Função para extrair informações sensíveis de um arquivo TXT

def extract_sensitive_info_from_txt(txt_path, results):
    with open(txt_path, 'rb') as txt_file:
        text = txt_file.read().decode('utf-8', errors='ignore')  # Decodificar explicitamente como UTF-8

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-\d{1,2}|\d{8}-\d{1,2}|\d{7,9}-\d{1,2}', text)
    matches_cpf = re.findall(r'\b(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})\b', text)
    matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in matches_telefone])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[txt_path] = results.get(txt_path, [])
        results[txt_path].extend(sensitive_info)

    return results

# Função para processar um diretório e seus subdiretórios, incluindo arquivos TXT
def process_directory_with_txt(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos TXT
            if filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)

# Função para extrair informações sensíveis de um arquivo DOCX
def extract_sensitive_info_from_docx(docx_path_or_text, results):
    if os.path.isfile(docx_path_or_text):  # Verifica se é um caminho de arquivo
        with open(docx_path_or_text, 'rb') as docx_file:
            doc = Document(docx_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    else:
        text = docx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_cnpj = re.findall(r'\d{2}.\d{3}.\d{3}/\d{4}-\d{2}', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('CNPJ', cnpj) for cnpj in matches_cnpj])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[docx_path_or_text] = results.get(docx_path_or_text, [])
        results[docx_path_or_text].extend(sensitive_info)

    return results

# Função para processar um diretório e seus subdiretórios, incluindo arquivos DOCX
def process_directory_with_docx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos DOCX
            if filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)


def process_directory(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Verifica se o arquivo é uma imagem, PDF, TXT ou DOCX e chama a função correspondente
            if filename.endswith(('.jpg', '.png', '.bmp')):
                image_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_image(image_path, results)
            elif filename.endswith('.pdf'):
                pdf_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pdf(pdf_path, results)
            elif filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)
            elif filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)
            elif filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)
            elif filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


# Define o caminho do diretório a ser processado
caminho_diretorio = ''
results = {}

# Chama a função para processar o diretório e seus subdiretórios, incluindo arquivos PDF, TXT e DOCX
process_directory(caminho_diretorio, results)

# Remove duplicatas nas informações sensíveis
for path, data in results.items():
    results[path] = list(set(data))

# Exibe as informações sensíveis encontradas
for path, data in results.items():
    path = os.path.normpath(path)
    print(f"Informações sensíveis encontradas em: {path}")


import os
import tkinter as tk
import customtkinter as ctk
from tkinter import filedialog, messagebox
from PIL import Image, ImageTk


class MyApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Minha Aplicação")
        self.geometry("1200x700")

        # Variáveis
        self.directory_path = tk.StringVar()
        self.key_path = tk.StringVar()
        self.sensitive_files = []

        # Configuração da aparência
        ctk.set_appearance_mode("dark")

        # Criar e exibir widgets
        self.create_widgets()

    def create_widgets(self):
        frame = ctk.CTkScrollableFrame(master=self, fg_color="transparent", border_color="#962CCA", border_width=2, height=600)
        frame.grid(row=0, column=0, rowspan=3, padx=10, pady=10)

        ctk.CTkButton(master=frame, text="INFO", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=0, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Escanear", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9", command=self.start_scan).grid(row=1, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Escolher Diretório", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9", command=self.choose_directory).grid(row=2, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Escolher Chave", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9", command=self.choose_key_file).grid(row=3, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Excluir Arquivos", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9", command=self.delete_files).grid(row=4, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Mover Arquivos", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=5, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Diretório Blacklist", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=6, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Relatório Blacklist", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=7, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Salvar", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=8, column=0, padx=30, pady=20, sticky="ew")
        ctk.CTkButton(master=frame, text="Sair", corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=9, column=0, padx=30, pady=20, sticky="ew")

        #Quadrado Vazio
        quadrado_vazio = ctk.CTkFrame(master=self, width=900, height=500, border_color="#962CCA", border_width=2)
        quadrado_vazio.grid(row=2, column=1, padx=10, pady=(0,60))
        quadrado_vazio.grid_rowconfigure(0, weight=1)
        quadrado_vazio.grid_columnconfigure(0, weight=1)


        self.output_text = ctk.CTkTextbox(master=quadrado_vazio, wrap=tk.WORD, border_color="#962CCA", border_width=1, height=500, width=900)
        self.output_text.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")

        ctk.CTkButton(master=self, text="", width=300, height=50, corner_radius=32, fg_color="#0f0913", hover_color="#53DEC9").grid(row=1, column=1, pady=10)

        # Carregar e exibir a imagem
        image = Image.open("C:\\Users\\lucas\\OneDrive\\Área de Trabalho\\logo_grupo\\logo.png")
        image = ImageTk.PhotoImage(image)
        image_label = ctk.CTkLabel(master=self, image=image, text="")
        image_label.grid(row=1, column=1, padx=10, pady=10)

    def choose_key_file(self):
        key_file = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if key_file:
            self.key_path.set(key_file)
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = key_file

    def choose_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.directory_path.set(directory)

    def start_scan(self):
        directory_path = self.directory_path.get()
        if directory_path:
            results = {}
            process_directory(directory_path, results)

            sensitive_files = []

            for path, data in results.items():
                results[path] = list(set(data))
                sensitive_files.append(path)

            for path, data in results.items():
                path = os.path.normpath(path)
                self.output_text.insert(tk.END, f"Informações sensíveis encontradas em: {path}\n")

                types_found = set(info[0] for info in data)
                for info_type in types_found:
                    self.output_text.insert(tk.END, f"{info_type} encontrado\n")

                self.output_text.insert(tk.END, "\n")

            self.sensitive_files = sensitive_files

    def delete_files(self):
        directory_path = self.directory_path.get()
        if directory_path:
            confirmation = messagebox.askyesno("Confirmação",
                                               "Tem certeza de que deseja excluir todos os arquivos sensíveis no diretório?")
            if confirmation:
                for sensitive_file in self.sensitive_files:
                    if os.path.isfile(sensitive_file):
                        os.remove(sensitive_file)

                messagebox.showinfo("Concluído", "Todos os arquivos sensíveis foram excluídos com sucesso.")


if __name__ == "__main__":
    app = MyApp()
    app.mainloop()
