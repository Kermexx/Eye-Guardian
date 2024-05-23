# Antes de tudo instale no terminall > pip install PyMuPDF
# Antes de tudo instale no terminal > pip install google-cloud-vision
# Antes de tudo instale no terminal > pip install python-docx
# Antes de tudo instale no terminal > pip install python-pptx
# Antes de tudo instale no terminal > pip install openpyxl
# Antes de tudo instale no terminal > pip install Pillow
# Antes de tudo instale no terminal > pip install schedule
# Antes de tudo instale no terminal > pip install customtkinter
# Antes de tudo instale no terminal > pip install pandas
# Antes de tudo instale no terminal > pip install matplotlib

import textwrap
import io
import imaplib #biblioteca para se conectar na caixa de e-mail
import email #decodificar partes do e-mail
from tkinter import Scrollbar, Toplevel, Listbox
import threading
import os  # funções para manipular caminhos de arquivos
import re  # ajuda a usar padrões de busca do scan
import shutil  # Copia e/ou move os arquivos
import fitz  # PymuPDF
from google.cloud import vision
from docx import Document  # Para lidar com arquivos DOCX
from pptx import Presentation  # powerpoint
import openpyxl  # Para lidar com arquivos XLSX
import tkinter as tk
from tkinter import Tk, Label, Entry, Button, messagebox
# imports de design
from tkinter import filedialog
from tkinter import messagebox
from tkinter import ttk
from tkinter import Canvas, Entry, Text, Button, PhotoImage
from pathlib import Path
import sys
from customtkinter import *
from PIL import Image, ImageTk
import time  # contador
import schedule  # contador
import customtkinter as ctk
from tkinter import simpledialog
import json  # transforamar o save em um arquivo JSON
import base64  # para a imagem do design virar base64
import io  # load do save
import csv  # relatório
from datetime import datetime
from openpyxl import Workbook
import pandas as pd
import matplotlib.pyplot as plt

#--------------------------------------------------------------------------#

religioes = [
    'Cristianismo',
    'Cristão',
    'Cristã',
    'Islamismo',
    'Islâmico',
    'Islâmica',
    'Hinduísmo',
    'Hindu',
    'Hinduísta',
    'Budismo',
    'Budista',
    'Sikhismo',
    'Sikh',
    'Judaísmo',
    'Judeu',
    'Judaica',
    'Bahaí',
    'Bahaísta',
    'Jainismo',
    'Jainista',
    'Espiritismo',
    'Espírita',
    'Ateísmo',
    'Ateu',
    'Ateia'
]

cores_etnias = [
    'Branco',
    'Negro',
    'Pardo',
    'Indígena',
    'Amarelo',
    'Asiático',
    'Outro / Não Declarado',
]


def extract_info_by_pattern(pattern, text, info_type, results):
    matches = re.findall(pattern, text, flags=re.IGNORECASE)
    if matches:
        results.extend([(info_type, match) for match in matches])
    return results


def format_rg(rg):
    # Remove pontos e hífens do RG e adiciona o hífen no formato desejado
    rg_limpo = re.sub(r'[^\d]', '', rg)
    return f'{rg_limpo[:-1]}.{rg_limpo[-1]}'


# Modifica a função extract_sensitive_info_from_xlsx para incluir formatação de RG
def extract_sensitive_info_from_xlsx(xlsx_path, results):
    wb = openpyxl.load_workbook(xlsx_path)

    sensitive_info = []

    # Itera sobre todas as folhas no arquivo Excel
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Itera sobre todas as células na folha
        for row in sheet.iter_rows(min_row=1, max_col=sheet.max_column, max_row=sheet.max_row, values_only=True):
            for cell_value in row:
                if cell_value:
                    # Aplica expressões regulares para encontrar informações sensíveis
                    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})|\d{8,9}-\d{1,2}', cell_value)
                    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', str(cell_value))
                    matches_email = re.findall(r'\S+@\S+', str(cell_value))
                    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', str(cell_value))
                    matches_genero = re.findall(
                        r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', cell_value)
                    valid_telefones = []

                    for telefone in matches_telefone:
                        numero_limpo = re.sub(r'[^\d]', '', telefone)
                        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                            valid_telefones.append(telefone)

                    for rg in matches_rg:
                        rg_formatado = format_rg(rg)
                        rg_in_cpf = any(rg_formatado in cpf for cpf in matches_cpf)
                        if not rg_in_cpf:
                            sensitive_info.append(('RG', rg_formatado))

                    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
                    sensitive_info.extend([('Email', email) for email in matches_email])
                    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
                    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
                    # Extrai informações sobre religiões
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', cell_value,
                                                             'Religião', sensitive_info)
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', cell_value,
                                                             'Cor/Etnia', sensitive_info)
    if sensitive_info:
        results[xlsx_path] = results.get(xlsx_path, [])
        results[xlsx_path].extend(sensitive_info)

    return results


def process_directory_with_xlsx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos XLSX
            if filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


def extract_sensitive_info_from_pptx(pptx_path_or_text, results):
    if os.path.isfile(pptx_path_or_text):  # Verifica se é um caminho de arquivo
        presentation = Presentation(pptx_path_or_text)
        text = "\n".join(
            [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

    else:
        text = pptx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)
    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pptx_path_or_text] = results.get(pptx_path_or_text, [])
        results[pptx_path_or_text].extend(sensitive_info)

    return results


# Modifica a função process_directory para incluir arquivos PPTX
def process_directory_with_pptx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos PPTX
            if filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)


# Função para extrair informações sensíveis de um arquivo PDF
def extract_sensitive_info_from_pdf(pdf_path, results):
    # Inicializa o documento PDF
    doc = fitz.open(pdf_path)

    sensitive_info = []

    # Itera sobre as páginas do documento PDF
    for page_number in range(doc.page_count):
        page = doc[page_number]

        # Extrai texto da página
        text = page.get_text()

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
        matches_email = re.findall(r'\S+@\S+', text)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

        valid_telefones = []

        # Filtra e formata números de telefone válidos
        for telefone in matches_telefone:
            numero_limpo = re.sub(r'[^\d]', '', telefone)
            if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                valid_telefones.append(telefone)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
            if not rg_in_cpf:
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
        # Extrai informações sobre religiões
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião',
                                                 sensitive_info)
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                                 sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pdf_path] = results.get(pdf_path, [])
        results[pdf_path].extend(sensitive_info)

    return results


# Função para extrair informações sensíveis de uma imagem
def extract_sensitive_info_from_image(image_path, results):
    # Inicializa o cliente Google Cloud Vision
    client = vision.ImageAnnotatorClient()

    # Lê o conteúdo da imagem
    with open(image_path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    # Primeiro, realiza a detecção de rosto, mas não adiciona ainda à lista de informações sensíveis
    face_response = client.face_detection(image=image)
    faces_detected = len(face_response.face_annotations) > 0

    # Realiza a detecção de texto
    text_response = client.text_detection(image=image)
    texts = text_response.text_annotations

    sensitive_info = []

    # Itera sobre os textos identificados para encontrar informações sensíveis
    for text in texts:
        text_content = text.description

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text_content)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2})', text_content)
        matches_email = re.findall(r'\S+@\S+', text_content)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text_content)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text_content)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            if not any(rg in cpf for cpf in matches_cpf):
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in matches_telefone])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])

    # Se um rosto e outras informações sensíveis foram detectados, adiciona "Rosto" à lista
    if faces_detected and sensitive_info:
        sensitive_info.insert(0, ('Rosto', 'Rosto encontrado'))

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[image_path] = results.get(image_path, [])
        results[image_path].extend(sensitive_info)

    return results


# Função para extrair informações sensíveis de um arquivo TXT
def extract_sensitive_info_from_txt(txt_path, results):
    with open(txt_path, 'rb') as txt_file:
        text = txt_file.read().decode('utf-8', errors='ignore')  # Decodificar explicitamente como UTF-8

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-\d{1,2}|\d{8}-\d{1,2}|\d{7,9}-\d{1,2}', text)
    matches_cpf = re.findall(r'\b(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})\b', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in matches_telefone])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[txt_path] = results.get(txt_path, [])
        results[txt_path].extend(sensitive_info)

    return results


# Função para processar um diretório e seus subdiretórios, incluindo arquivos TXT
def process_directory_with_txt(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos TXT
            if filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)


# Função para extrair informações sensíveis de um arquivo DOCX
def extract_sensitive_info_from_docx(docx_path_or_text, results):
    if os.path.isfile(docx_path_or_text):  # Verifica se é um caminho de arquivo
        with open(docx_path_or_text, 'rb') as docx_file:
            doc = Document(docx_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    else:
        text = docx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[docx_path_or_text] = results.get(docx_path_or_text, [])
        results[docx_path_or_text].extend(sensitive_info)

    return results


# Função para processar um diretório e seus subdiretórios, incluindo arquivos DOCX
def process_directory_with_docx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos DOCX
            if filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)


def process_directory(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Verifica se o arquivo é uma imagem, PDF, TXT ou DOCX e chama a função correspondente
            if filename.endswith(('.jpg', '.png', '.bmp')):
                image_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_image(image_path, results)
            elif filename.endswith('.pdf'):
                pdf_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pdf(pdf_path, results)
            elif filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)
            elif filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)
            elif filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)
            elif filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


# Define o caminho do diretório a ser processado
caminho_diretorio = ''
results = {}

# Chama a função para processar o diretório e seus subdiretórios, incluindo arquivos PDF, TXT e DOCX
process_directory(caminho_diretorio, results)

# Remove duplicatas nas informações sensíveis
for path, data in results.items():
    results[path] = list(set(data))

# Exibe as informações sensíveis encontradas
for path, data in results.items():
    path = os.path.normpath(path)
    print(f"Informações sensíveis encontradas em: {path}")

def get_base_dir():
    if getattr(sys, 'frozen', False):
        # Estamos rodando em um executável PyInstaller
        return sys._MEIPASS
    else:
        # Estamos rodando em um ambiente Python normal
        return os.path.dirname(os.path.abspath(__file__))

class MeuApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Eye Guardian")
        self.geometry("1200x700")
        self.resizable(False, False)

        canto_preto_base64 = (
            """/9j/4QAiRXhpZgAATU0AKgAAAAgAAQESAAMAAAABAAEAAAAAAAD/7RuyUGhvdG9zaG9wIDMuMAA4QklNBCUAAAAAABAAAAAAAAAAAAAAAAAAAAAAOEJJTQQ6AAAAAAD5AAAAEAAAAAEAAAAAAAtwcmludE91dHB1dAAAAAUAAAAAUHN0U2Jvb2wBAAAAAEludGVlbnVtAAAAAEludGUAAAAAQ2xybQAAAA9wcmludFNpeHRlZW5CaXRib29sAAAAAAtwcmludGVyTmFtZVRFWFQAAAABAAAAAAAPcHJpbnRQcm9vZlNldHVwT2JqYwAAABYAQwBvAG4AZgBpAGcAdQByAGEA5wDjAG8AIABkAGUAIABQAHIAbwB2AGEAAAAAAApwcm9vZlNldHVwAAAAAQAAAABCbHRuZW51bQAAAAxidWlsdGluUHJvb2YAAAAJcHJvb2ZDTVlLADhCSU0EOwAAAAACLQAAABAAAAABAAAAAAAScHJpbnRPdXRwdXRPcHRpb25zAAAAFwAAAABDcHRuYm9vbAAAAAAAQ2xicmJvb2wAAAAAAFJnc01ib29sAAAAAABDcm5DYm9vbAAAAAAAQ250Q2Jvb2wAAAAAAExibHNib29sAAAAAABOZ3R2Ym9vbAAAAAAARW1sRGJvb2wAAAAAAEludHJib29sAAAAAABCY2tnT2JqYwAAAAEAAAAAAABSR0JDAAAAAwAAAABSZCAgZG91YkBv4AAAAAAAAAAAAEdybiBkb3ViQG/gAAAAAAAAAAAAQmwgIGRvdWJAb+AAAAAAAAAAAABCcmRUVW50RiNSbHQAAAAAAAAAAAAAAABCbGQgVW50RiNSbHQAAAAAAAAAAAAAAABSc2x0VW50RiNQeGxAUgAAAAAAAAAAAAp2ZWN0b3JEYXRhYm9vbAEAAAAAUGdQc2VudW0AAAAAUGdQcwAAAABQZ1BDAAAAAExlZnRVbnRGI1JsdAAAAAAAAAAAAAAAAFRvcCBVbnRGI1JsdAAAAAAAAAAAAAAAAFNjbCBVbnRGI1ByY0BZAAAAAAAAAAAAEGNyb3BXaGVuUHJpbnRpbmdib29sAAAAAA5jcm9wUmVjdEJvdHRvbWxvbmcAAAAAAAAADGNyb3BSZWN0TGVmdGxvbmcAAAAAAAAADWNyb3BSZWN0UmlnaHRsb25nAAAAAAAAAAtjcm9wUmVjdFRvcGxvbmcAAAAAADhCSU0D7QAAAAAAEABIAAAAAQACAEgAAAABAAI4QklNBCYAAAAAAA4AAAAAAAAAAAAAP4AAADhCSU0EDQAAAAAABAAAAHg4QklNBBkAAAAAAAQAAAAeOEJJTQPzAAAAAAAJAAAAAAAAAAABADhCSU0nEAAAAAAACgABAAAAAAAAAAI4QklNA/UAAAAAAEgAL2ZmAAEAbGZmAAYAAAAAAAEAL2ZmAAEAoZmaAAYAAAAAAAEAMgAAAAEAWgAAAAYAAAAAAAEANQAAAAEALQAAAAYAAAAAAAE4QklNA/gAAAAAAHAAAP////////////////////////////8D6AAAAAD/////////////////////////////A+gAAAAA/////////////////////////////wPoAAAAAP////////////////////////////8D6AAAOEJJTQQAAAAAAAACAAI4QklNBAIAAAAAAAYAAAAAAAA4QklNBDAAAAAAAAMBAQEAOEJJTQQtAAAAAAAGAAEAAAAGOEJJTQQIAAAAAAAQAAAAAQAAAkAAAAJAAAAAADhCSU0EHgAAAAAABAAAAAA4QklNBBoAAAAAA00AAAAGAAAAAAAAAAAAAAK8AAAEsAAAAAwAUwBlAG0AIABUAO0AdAB1AGwAbwAtADEAAAABAAAAAAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAAABLAAAAK8AAAAAAAAAAAAAAAAAAAAAAEAAAAAAAAAAAAAAAAAAAAAAAAAEAAAAAEAAAAAAABudWxsAAAAAgAAAAZib3VuZHNPYmpjAAAAAQAAAAAAAFJjdDEAAAAEAAAAAFRvcCBsb25nAAAAAAAAAABMZWZ0bG9uZwAAAAAAAAAAQnRvbWxvbmcAAAK8AAAAAFJnaHRsb25nAAAEsAAAAAZzbGljZXNWbExzAAAAAU9iamMAAAABAAAAAAAFc2xpY2UAAAASAAAAB3NsaWNlSURsb25nAAAAAAAAAAdncm91cElEbG9uZwAAAAAAAAAGb3JpZ2luZW51bQAAAAxFU2xpY2VPcmlnaW4AAAANYXV0b0dlbmVyYXRlZAAAAABUeXBlZW51bQAAAApFU2xpY2VUeXBlAAAAAEltZyAAAAAGYm91bmRzT2JqYwAAAAEAAAAAAABSY3QxAAAABAAAAABUb3AgbG9uZwAAAAAAAAAATGVmdGxvbmcAAAAAAAAAAEJ0b21sb25nAAACvAAAAABSZ2h0bG9uZwAABLAAAAADdXJsVEVYVAAAAAEAAAAAAABudWxsVEVYVAAAAAEAAAAAAABNc2dlVEVYVAAAAAEAAAAAAAZhbHRUYWdURVhUAAAAAQAAAAAADmNlbGxUZXh0SXNIVE1MYm9vbAEAAAAIY2VsbFRleHRURVhUAAAAAQAAAAAACWhvcnpBbGlnbmVudW0AAAAPRVNsaWNlSG9yekFsaWduAAAAB2RlZmF1bHQAAAAJdmVydEFsaWduZW51bQAAAA9FU2xpY2VWZXJ0QWxpZ24AAAAHZGVmYXVsdAAAAAtiZ0NvbG9yVHlwZWVudW0AAAARRVNsaWNlQkdDb2xvclR5cGUAAAAATm9uZQAAAAl0b3BPdXRzZXRsb25nAAAAAAAAAApsZWZ0T3V0c2V0bG9uZwAAAAAAAAAMYm90dG9tT3V0c2V0bG9uZwAAAAAAAAALcmlnaHRPdXRzZXRsb25nAAAAAAA4QklNBCgAAAAAAAwAAAACP/AAAAAAAAA4QklNBBEAAAAAAAEBADhCSU0EFAAAAAAABAAAAAY4QklNBAwAAAAAEnQAAAABAAAAoAAAAF0AAAHgAACuYAAAElgAGAAB/9j/4gxYSUNDX1BST0ZJTEUAAQEAAAxITGlubwIQAABtbnRyUkdCIFhZWiAHzgACAAkABgAxAABhY3NwTVNGVAAAAABJRUMgc1JHQgAAAAAAAAAAAAAAAQAA9tYAAQAAAADTLUhQICAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABFjcHJ0AAABUAAAADNkZXNjAAABhAAAAGx3dHB0AAAB8AAAABRia3B0AAACBAAAABRyWFlaAAACGAAAABRnWFlaAAACLAAAABRiWFlaAAACQAAAABRkbW5kAAACVAAAAHBkbWRkAAACxAAAAIh2dWVkAAADTAAAAIZ2aWV3AAAD1AAAACRsdW1pAAAD+AAAABRtZWFzAAAEDAAAACR0ZWNoAAAEMAAAAAxyVFJDAAAEPAAACAxnVFJDAAAEPAAACAxiVFJDAAAEPAAACAx0ZXh0AAAAAENvcHlyaWdodCAoYykgMTk5OCBIZXdsZXR0LVBhY2thcmQgQ29tcGFueQAAZGVzYwAAAAAAAAASc1JHQiBJRUM2MTk2Ni0yLjEAAAAAAAAAAAAAABJzUkdCIElFQzYxOTY2LTIuMQAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAWFlaIAAAAAAAAPNRAAEAAAABFsxYWVogAAAAAAAAAAAAAAAAAAAAAFhZWiAAAAAAAABvogAAOPUAAAOQWFlaIAAAAAAAAGKZAAC3hQAAGNpYWVogAAAAAAAAJKAAAA+EAAC2z2Rlc2MAAAAAAAAAFklFQyBodHRwOi8vd3d3LmllYy5jaAAAAAAAAAAAAAAAFklFQyBodHRwOi8vd3d3LmllYy5jaAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABkZXNjAAAAAAAAAC5JRUMgNjE5NjYtMi4xIERlZmF1bHQgUkdCIGNvbG91ciBzcGFjZSAtIHNSR0IAAAAAAAAAAAAAAC5JRUMgNjE5NjYtMi4xIERlZmF1bHQgUkdCIGNvbG91ciBzcGFjZSAtIHNSR0IAAAAAAAAAAAAAAAAAAAAAAAAAAAAAZGVzYwAAAAAAAAAsUmVmZXJlbmNlIFZpZXdpbmcgQ29uZGl0aW9uIGluIElFQzYxOTY2LTIuMQAAAAAAAAAAAAAALFJlZmVyZW5jZSBWaWV3aW5nIENvbmRpdGlvbiBpbiBJRUM2MTk2Ni0yLjEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAHZpZXcAAAAAABOk/gAUXy4AEM8UAAPtzAAEEwsAA1yeAAAAAVhZWiAAAAAAAEwJVgBQAAAAVx/nbWVhcwAAAAAAAAABAAAAAAAAAAAAAAAAAAAAAAAAAo8AAAACc2lnIAAAAABDUlQgY3VydgAAAAAAAAQAAAAABQAKAA8AFAAZAB4AIwAoAC0AMgA3ADsAQABFAEoATwBUAFkAXgBjAGgAbQByAHcAfACBAIYAiwCQAJUAmgCfAKQAqQCuALIAtwC8AMEAxgDLANAA1QDbAOAA5QDrAPAA9gD7AQEBBwENARMBGQEfASUBKwEyATgBPgFFAUwBUgFZAWABZwFuAXUBfAGDAYsBkgGaAaEBqQGxAbkBwQHJAdEB2QHhAekB8gH6AgMCDAIUAh0CJgIvAjgCQQJLAlQCXQJnAnECegKEAo4CmAKiAqwCtgLBAssC1QLgAusC9QMAAwsDFgMhAy0DOANDA08DWgNmA3IDfgOKA5YDogOuA7oDxwPTA+AD7AP5BAYEEwQgBC0EOwRIBFUEYwRxBH4EjASaBKgEtgTEBNME4QTwBP4FDQUcBSsFOgVJBVgFZwV3BYYFlgWmBbUFxQXVBeUF9gYGBhYGJwY3BkgGWQZqBnsGjAadBq8GwAbRBuMG9QcHBxkHKwc9B08HYQd0B4YHmQesB78H0gflB/gICwgfCDIIRghaCG4IggiWCKoIvgjSCOcI+wkQCSUJOglPCWQJeQmPCaQJugnPCeUJ+woRCicKPQpUCmoKgQqYCq4KxQrcCvMLCwsiCzkLUQtpC4ALmAuwC8gL4Qv5DBIMKgxDDFwMdQyODKcMwAzZDPMNDQ0mDUANWg10DY4NqQ3DDd4N+A4TDi4OSQ5kDn8Omw62DtIO7g8JDyUPQQ9eD3oPlg+zD88P7BAJECYQQxBhEH4QmxC5ENcQ9RETETERTxFtEYwRqhHJEegSBxImEkUSZBKEEqMSwxLjEwMTIxNDE2MTgxOkE8UT5RQGFCcUSRRqFIsUrRTOFPAVEhU0FVYVeBWbFb0V4BYDFiYWSRZsFo8WshbWFvoXHRdBF2UXiReuF9IX9xgbGEAYZRiKGK8Y1Rj6GSAZRRlrGZEZtxndGgQaKhpRGncanhrFGuwbFBs7G2MbihuyG9ocAhwqHFIcexyjHMwc9R0eHUcdcB2ZHcMd7B4WHkAeah6UHr4e6R8THz4faR+UH78f6iAVIEEgbCCYIMQg8CEcIUghdSGhIc4h+yInIlUigiKvIt0jCiM4I2YjlCPCI/AkHyRNJHwkqyTaJQklOCVoJZclxyX3JicmVyaHJrcm6CcYJ0kneierJ9woDSg/KHEooijUKQYpOClrKZ0p0CoCKjUqaCqbKs8rAis2K2krnSvRLAUsOSxuLKIs1y0MLUEtdi2rLeEuFi5MLoIuty7uLyQvWi+RL8cv/jA1MGwwpDDbMRIxSjGCMbox8jIqMmMymzLUMw0zRjN/M7gz8TQrNGU0njTYNRM1TTWHNcI1/TY3NnI2rjbpNyQ3YDecN9c4FDhQOIw4yDkFOUI5fzm8Ofk6Njp0OrI67zstO2s7qjvoPCc8ZTykPOM9Ij1hPaE94D4gPmA+oD7gPyE/YT+iP+JAI0BkQKZA50EpQWpBrEHuQjBCckK1QvdDOkN9Q8BEA0RHRIpEzkUSRVVFmkXeRiJGZ0arRvBHNUd7R8BIBUhLSJFI10kdSWNJqUnwSjdKfUrESwxLU0uaS+JMKkxyTLpNAk1KTZNN3E4lTm5Ot08AT0lPk0/dUCdQcVC7UQZRUFGbUeZSMVJ8UsdTE1NfU6pT9lRCVI9U21UoVXVVwlYPVlxWqVb3V0RXklfgWC9YfVjLWRpZaVm4WgdaVlqmWvVbRVuVW+VcNVyGXNZdJ114XcleGl5sXr1fD19hX7NgBWBXYKpg/GFPYaJh9WJJYpxi8GNDY5dj62RAZJRk6WU9ZZJl52Y9ZpJm6Gc9Z5Nn6Wg/aJZo7GlDaZpp8WpIap9q92tPa6dr/2xXbK9tCG1gbbluEm5rbsRvHm94b9FwK3CGcOBxOnGVcfByS3KmcwFzXXO4dBR0cHTMdSh1hXXhdj52m3b4d1Z3s3gReG54zHkqeYl553pGeqV7BHtje8J8IXyBfOF9QX2hfgF+Yn7CfyN/hH/lgEeAqIEKgWuBzYIwgpKC9INXg7qEHYSAhOOFR4Wrhg6GcobXhzuHn4gEiGmIzokziZmJ/opkisqLMIuWi/yMY4zKjTGNmI3/jmaOzo82j56QBpBukNaRP5GokhGSepLjk02TtpQglIqU9JVflcmWNJaflwqXdZfgmEyYuJkkmZCZ/JpomtWbQpuvnByciZz3nWSd0p5Anq6fHZ+Ln/qgaaDYoUehtqImopajBqN2o+akVqTHpTilqaYapoum/adup+CoUqjEqTepqaocqo+rAqt1q+msXKzQrUStuK4trqGvFq+LsACwdbDqsWCx1rJLssKzOLOutCW0nLUTtYq2AbZ5tvC3aLfguFm40blKucK6O7q1uy67p7whvJu9Fb2Pvgq+hL7/v3q/9cBwwOzBZ8Hjwl/C28NYw9TEUcTOxUvFyMZGxsPHQce/yD3IvMk6ybnKOMq3yzbLtsw1zLXNNc21zjbOts83z7jQOdC60TzRvtI/0sHTRNPG1EnUy9VO1dHWVdbY11zX4Nhk2OjZbNnx2nba+9uA3AXcit0Q3ZbeHN6i3ynfr+A24L3hROHM4lPi2+Nj4+vkc+T85YTmDeaW5x/nqegy6LzpRunQ6lvq5etw6/vshu0R7ZzuKO6070DvzPBY8OXxcvH/8ozzGfOn9DT0wvVQ9d72bfb794r4Gfio+Tj5x/pX+uf7d/wH/Jj9Kf26/kv+3P9t////7QAMQWRvYmVfQ00AAv/uAA5BZG9iZQBkgAAAAAH/2wCEAAwICAgJCAwJCQwRCwoLERUPDAwPFRgTExUTExgRDAwMDAwMEQwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwBDQsLDQ4NEA4OEBQODg4UFA4ODg4UEQwMDAwMEREMDAwMDAwRDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDP/AABEIAF0AoAMBIgACEQEDEQH/3QAEAAr/xAE/AAABBQEBAQEBAQAAAAAAAAADAAECBAUGBwgJCgsBAAEFAQEBAQEBAAAAAAAAAAEAAgMEBQYHCAkKCxAAAQQBAwIEAgUHBggFAwwzAQACEQMEIRIxBUFRYRMicYEyBhSRobFCIyQVUsFiMzRygtFDByWSU/Dh8WNzNRaisoMmRJNUZEXCo3Q2F9JV4mXys4TD03Xj80YnlKSFtJXE1OT0pbXF1eX1VmZ2hpamtsbW5vY3R1dnd4eXp7fH1+f3EQACAgECBAQDBAUGBwcGBTUBAAIRAyExEgRBUWFxIhMFMoGRFKGxQiPBUtHwMyRi4XKCkkNTFWNzNPElBhaisoMHJjXC0kSTVKMXZEVVNnRl4vKzhMPTdePzRpSkhbSVxNTk9KW1xdXl9VZmdoaWprbG1ub2JzdHV2d3h5ent8f/2gAMAwEAAhEDEQA/AOT9Kr9xv3JelV+437lNJJTD0qv3G/co211it5DGghpgx5Iqhb/NP/qn8iSn1zB+rP1Zfg4z39IwXPdTW5zjj1kkljS5zvajf81vqv8A+U2D/wCw9f8A5FXOn/8AJ+J/xFX/AFDEdJLmf81vqv8A+U2D/wCw9f8A5FL/AJrfVf8A8psH/wBh6/8AyK00klOZ/wA1vqv/AOU2D/7D1/8AkUv+a31X/wDKbB/9h6//ACK00klOZ/zW+q//AJTYP/sPX/5FL/mt9V//ACmwf/Yev/yK00klOZ/zW+q//lNg/wDsPX/5FL/mt9V//KbB/wDYev8A8itNJJTmf81vqv8A+U2D/wCw9f8A5FL/AJrfVf8A8psH/wBh6/8AyK00klOZ/wA1vqv/AOU2D/7D1/8AkUv+a31X/wDKbB/9h6//ACK00klOZ/zW+q//AJTYP/sPX/5FL/mt9V//ACmwf/Yev/yK00klP//Q5ZJJJJSlC3+af/VP5FNQt/mn/wBU/kSU+3dP/wCT8T/iKv8AqGI6B0//AJPxP+Iq/wCoYjoJUkkkkpSSSSSlJJJJKUkkkkpSSSSSlJJJJKUkkkkp/9HlkkkklKULf5p/9U/kU1C3+af/AFT+RJT7d0//AJPxP+Iq/wCoYjoHT/8Ak/E/4ir/AKhiOglSSSSSlJJJJKUkkkkpSSSSSlJJJJKUkkkkpSSSSSn/0uWSSSSUpQt/mn/1T+RTULf5p/8AVP5ElPt3T/8Ak/E/4ir/AKhiOgdP/wCT8T/iKv8AqGI6CVJJJJKUkkkkpSSSSSlJJJJKUkkkkpSSSSSlJJJJKf/T5ZJJJJSlC3+af/VP5FNQt/mn/wBU/kSU+3dP/wCT8T/iKv8AqGI6B0//AJPxP+Iq/wCoYjoJUkkkkpSSSSSlJJJJKUkkkkpSSSSSlJJJJKUkkkkp/9TGq6Fm24rb2ur3vaH14xJ9RzXB9jO2xtlldVllVTvpsWe1r3Dc1jiI3SGk6QHbuP3XLscN3W/2VUG1y/bVDy+wU7/Ts+zuya/SNf2ltX02er6L7fs/q7LPsy5+r9vfZqdm/wBLa30N2yZ2j0/T3e719v0Pz08iGlE/Rq45c0TPijHcVxy4P7/B7ccnFj/m+D3PW0BXaeGPMCTDTwPzuPoodzXCl5LXAbXakHwWnd+2PRPqf0f3RHp+nGyz1Nm327PQ9X1Nv/pJNn/t37HlfaZ9LY77T/NxwPp+n/hPo7f8J/4Ih6e5+z+1lBz2Ljjrr653/wCkn1vp/wDyfif8RV/1DEdA6f8A8n4n/EVf9QxHTGZSSSSSlJJJJKUkkkkpSSSSSlJJJJKUkkkkpSSSSSn/2ThCSU0EIQAAAAAAVQAAAAEBAAAADwBBAGQAbwBiAGUAIABQAGgAbwB0AG8AcwBoAG8AcAAAABMAQQBkAG8AYgBlACAAUABoAG8AdABvAHMAaABvAHAAIABDAFMANgAAAAEAOEJJTQQGAAAAAAAHAAgAAAABAQD/4Q2taHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLwA8P3hwYWNrZXQgYmVnaW49Iu+7vyIgaWQ9Ilc1TTBNcENlaGlIenJlU3pOVGN6a2M5ZCI/PiA8eDp4bXBtZXRhIHhtbG5zOng9ImFkb2JlOm5zOm1ldGEvIiB4OnhtcHRrPSJBZG9iZSBYTVAgQ29yZSA1LjMtYzAxMSA2Ni4xNDU2NjEsIDIwMTIvMDIvMDYtMTQ6NTY6MjcgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0RXZ0PSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VFdmVudCMiIHhtbG5zOmRjPSJodHRwOi8vcHVybC5vcmcvZGMvZWxlbWVudHMvMS4xLyIgeG1sbnM6cGhvdG9zaG9wPSJodHRwOi8vbnMuYWRvYmUuY29tL3Bob3Rvc2hvcC8xLjAvIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCBDUzYgKFdpbmRvd3MpIiB4bXA6Q3JlYXRlRGF0ZT0iMjAyNC0wNC0yOVQyMjowMDozOS0wMzowMCIgeG1wOk1ldGFkYXRhRGF0ZT0iMjAyNC0wNC0yOVQyMjowMDozOS0wMzowMCIgeG1wOk1vZGlmeURhdGU9IjIwMjQtMDQtMjlUMjI6MDA6MzktMDM6MDAiIHhtcE1NOkluc3RhbmNlSUQ9InhtcC5paWQ6NTQ4NTg1MEY4RDA2RUYxMTgzQ0NDMjdEOUZDOEU5QUQiIHhtcE1NOkRvY3VtZW50SUQ9InhtcC5kaWQ6QURFMThFMkM4OTA2RUYxMTgzQ0NDMjdEOUZDOEU5QUQiIHhtcE1NOk9yaWdpbmFsRG9jdW1lbnRJRD0ieG1wLmRpZDpBREUxOEUyQzg5MDZFRjExODNDQ0MyN0Q5RkM4RTlBRCIgZGM6Zm9ybWF0PSJpbWFnZS9qcGVnIiBwaG90b3Nob3A6Q29sb3JNb2RlPSIzIj4gPHhtcE1NOkhpc3Rvcnk+IDxyZGY6U2VxPiA8cmRmOmxpIHN0RXZ0OmFjdGlvbj0iY3JlYXRlZCIgc3RFdnQ6aW5zdGFuY2VJRD0ieG1wLmlpZDpBREUxOEUyQzg5MDZFRjExODNDQ0MyN0Q5RkM4RTlBRCIgc3RFdnQ6d2hlbj0iMjAyNC0wNC0yOVQyMjowMDozOS0wMzowMCIgc3RFdnQ6c29mdHdhcmVBZ2VudD0iQWRvYmUgUGhvdG9zaG9wIENTNiAoV2luZG93cykiLz4gPHJkZjpsaSBzdEV2dDphY3Rpb249InNhdmVkIiBzdEV2dDppbnN0YW5jZUlEPSJ4bXAuaWlkOjU0ODU4NTBGOEQwNkVGMTE4M0NDQzI3RDlGQzhFOUFEIiBzdEV2dDp3aGVuPSIyMDI0LTA0LTI5VDIyOjAwOjM5LTAzOjAwIiBzdEV2dDpzb2Z0d2FyZUFnZW50PSJBZG9iZSBQaG90b3Nob3AgQ1M2IChXaW5kb3dzKSIgc3RFdnQ6Y2hhbmdlZD0iLyIvPiA8L3JkZjpTZXE+IDwveG1wTU06SGlzdG9yeT4gPC9yZGY6RGVzY3JpcHRpb24+IDwvcmRmOlJERj4gPC94OnhtcG1ldGE+ICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgICAgPD94cGFja2V0IGVuZD0idyI/Pv/uAA5BZG9iZQBkQAAAAAH/2wCEAAEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQEBAQECAgICAgICAgICAgMDAwMDAwMDAwMBAQEBAQEBAQEBAQICAQICAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDAwMDA//AABEIArwEsAMBEQACEQEDEQH/3QAEAJb/xAGiAAAABgIDAQAAAAAAAAAAAAAHCAYFBAkDCgIBAAsBAAAGAwEBAQAAAAAAAAAAAAYFBAMHAggBCQAKCxAAAgEDBAEDAwIDAwMCBgl1AQIDBBEFEgYhBxMiAAgxFEEyIxUJUUIWYSQzF1JxgRhikSVDobHwJjRyChnB0TUn4VM2gvGSokRUc0VGN0djKFVWVxqywtLi8mSDdJOEZaOzw9PjKThm83UqOTpISUpYWVpnaGlqdnd4eXqFhoeIiYqUlZaXmJmapKWmp6ipqrS1tre4ubrExcbHyMnK1NXW19jZ2uTl5ufo6er09fb3+Pn6EQACAQMCBAQDBQQEBAYGBW0BAgMRBCESBTEGACITQVEHMmEUcQhCgSORFVKhYhYzCbEkwdFDcvAX4YI0JZJTGGNE8aKyJjUZVDZFZCcKc4OTRnTC0uLyVWV1VjeEhaOzw9Pj8ykalKS0xNTk9JWltcXV5fUoR1dmOHaGlqa2xtbm9md3h5ent8fX5/dIWGh4iJiouMjY6Pg5SVlpeYmZqbnJ2en5KjpKWmp6ipqqusra6vr/2gAMAwEAAhEDEQA/ANUf37r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//QoBOMxp+uPoT/AK9JT/8AXv37r3WNsNiG/Visa1/rqoaU/wC9xe/de64/wPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XuvfwPC/86fF/wDnvpP+vPv3XutjL/hO30j0v2j2j8lqLszqHq/sWjxWwdh1WLpd9bA2pu6mxlTUbizcVRPj4dwYnIRUU9RGiq7RhWdVAJIHuj+XSi3AJaox1tW/7JX8OPp/spnxmt/4gfqz8/8Akq+6VPr0r0r/AAjr3+yV/Dc/9yl/GbgAD/jA/Vn0H0H/AB6v0Hv1T69e0r/COvH4V/Dc/X4l/GY/6/Q/Vn/2K/4+/VPr17Sv8I69/slfw4N/+cTPjNz9f+MD9Wc/U8/79Xnkn36p9evaV/hHXv8AZK/hve/+yl/Ga4+h/wBA/Vl/oB9f7q/0Fv8AW9+qfXr2lf4R12PhZ8OBYj4mfGcECwt0P1bwP6D/AH6v09+qfXr2lf4R10PhX8OAbj4l/GYEcgjofqwG/wDW/wDdW/v1T69a0r/COvD4V/DgWt8S/jMLXtbofqwWv9bf79X8+/VPr1vSv8I67/2Sz4ccf84mfGfgkj/jA/VvBP1I/wB+rwT79U+vXtK/wjrofCv4bi4HxL+Mwubm3Q/Vgufrc/79X63Hv1T69e0r/COvH4V/Dc3v8S/jMbm5v0P1Ybn6XP8Av1frb36p9evaV/hHXv8AZK/hx/3iZ8Zv/RD9Wf0t/wA8r/T36p9evaV/hHXv9kr+HBvf4mfGbn6/8YH6s5txz/v1eePfqn169pX+Ede/2Sv4b/T/AGUv4zW/p/oH6s/x/wCzV/xPv1T69e0r/COvf7JX8OP+8TPjN/6Ifqz+hH/PK/0Pv1T69e0r/COu/wDZLPhxe/8Aspnxnvxz/oH6svxyOf7q34I9+qfXr2lf4R10fhX8NyAD8S/jMQv0B6H6ssOSeB/dXjk+/VPr1rSv8I69/slfw3P1+Jfxm+t/+ZD9WfW97/8AHq/W49+qfXrelf4R17/ZK/hx/wB4mfGb6af+ZD9Wfp/1P/Hq/T/D36p9evaV/hHXZ+Fnw4P1+JnxnP556H6tP+97V9+qfXr2lf4R17/ZLPhx/wB4mfGf/wBEP1Z/W/8Azyv9ffqn169pX+EddH4V/Dc8n4l/GY/6/Q/Vh+nA/wCYV/A9+qfXr2lf4R17/ZK/hv8AT/ZS/jNbgW/0D9Wfg3H/ADCv4J9+qfXr2lf4R11/slXw2/7xK+Mv/ohurP8A7Ff8Pfqn169pX+Edd/7JX8N73/2Uv4zXP1P+gfqy/wBdX/PK/wCq5/1/fqn169pX+Ede/wBkr+G5vf4l/GbkEH/jA/VnINrg/wC/V+ht79U+vXtK/wAI68PhX8Nx9PiX8Zhf626H6sF+Qf8Anlf6ge/VPr1rSv8ACOvf7JX8NzcH4l/GazfqH+gfqzn/AF/9+rz9ffqn163pX+Ede/2Sv4bk3PxL+MxP9f8AQP1Zf/b/AN1ffqn169pX+Ede/wBkr+G//eJfxm4vb/jA/Vn5Gk/8wr+VFv8AW9+qfXr2lf4R14fCv4cAWHxL+MwH9B0P1Zb6afp/dW36eP8AW9+qfXr2lf4R17/ZK/hueD8S/jNY/Uf6B+rOf/XV9+qfXr2lf4R13/slnw4/7xM+M/1v/wAyH6s+p5J/49X639+qfXr2lf4R10fhX8Nze/xL+MxuCpv0P1Ybg3uD/v1fob+/VPr17Sv8I68PhX8NxyPiX8ZgeRcdD9WDgggj/j1fyD79U+vWtK/wjrv/AGSz4cf94mfGf8n/AJkP1b+RY/8AMK/ke/VPr1vSv8I69/slnw4/7xM+M/BuP+MD9W8H+v8Ax6v19+qfXr2lf4R11/slfw3P/cpfxm/r/wAyH6s/+xX36p9evaV/hHXv9kr+G/8A3iX8ZuDcf8YH6s+t73/49X639+qfXr2lf4R11/slXw2+v+ylfGW/9f8AQN1Z+Pp/zCvv1T69e0r/AAjrw+FXw2H0+JXxlF+TbobqwXt9P+YV/Fvfqn169pX+Edd/7JV8N/8AvEv4zc2/5oN1Z+LW/wCYV/Fvfqn169pX+Eddj4WfDgfT4mfGcW+luiOreOCP+eV/obe/VPr1rSv8I66/2Sv4cEAH4mfGawGkD/QP1ZYDVqsB/dXgajf/AF/fqn169pX+Ede/2Sv4ccf84mfGbj6f8YH6s4v9bf79Xi/v1T69b0r/AAjrv/ZLPhx/3iZ8Z/6f8yH6s+n9P+PV9+qfXr2lf4R11/slfw4uT/spnxmuTcn/AED9WXJta5P91frbj/W9+qfXr2lf4R14fCv4bg3HxL+MwIFgR0P1YCBe9v8Aj1fpcX9+qfXr2lf4R17/AGSv4ccf84l/Gbj6f8YH6s4/PH+/V459+qfXr2lf4R17/ZK/hv8AT/ZS/jNbn/mg/Vn5+v8AzCv59+qfXr2lf4R13/slnw45/wCcTPjPz9f+MD9W8/jn/fq88e/VPr17Sv8ACOuv9kq+G4+nxL+M3PB/4wP1Zz9P+zV/w9+qfXr2lf4R17/ZK/hx/wB4l/Gb8/8ANB+rP7X6v+YV/P59+qfXr2lf4R13/slnw44/5xM+M/BLD/jBHVvDEEEj/fq/Ug8+/VPr1rSv8I66/wBkr+G5+vxL+M3FwP8AjA/Vn0JBI/49X8kX9+qfXrelf4R17/ZKvhvz/wA4l/Gbk3P/ABgbqzkj8n/fq8n36p9evaV/hHXv9kr+G/H/ADiX8ZuL2/4wP1Zxf62/36vF/fqn169pX+EdePwq+G5Nz8S/jMT9bnofqy/H/kq+/VPr17Sv8I67/wBks+HBBB+JnxnIJJI/0D9WWJP1JH91eSffqn169pX+Eddf7JX8N/8AvEv4zf8Aoh+rP8B/zyv9B79U+vXtK/wjrx+Ffw3PJ+JfxmJIA56H6s+gAAH/AB6v0AHv1T69a0r/AAjr3+yV/Dfj/nEv4zcfT/jA/VnH54/36vHPv1T69b0r/COvf7JX8Nybn4l/GYk3BP8AoH6svY/UX/urfm3v1T69e0r/AAjrs/Cz4cE3PxM+M5P9T0P1YT/t/wC6vv1T69e0r/COuv8AZKvhvx/ziX8ZuOR/xgfqzj6fT/fq8fT36p9evaV/hHXv9kr+G/8A3iX8Zv6f8yH6s/xH/PK/0J/2/v1T69e0r/COvf7JX8N7W/2Uv4zW4Nv9A/Vlri9jb+6v4uf9v79U+vWtK/wjrx+Ffw4P1+JfxmNvpfofqw/0/wCzV/w9+qfXrelf4R12fhZ8OCbn4mfGckcgnofq0n8f9mr/AIe/VPr17Sv8I66/2Sv4cf8AeJfxm+lv+ZD9WfT+n/Hq/T36p9evaV/hHXf+yWfDgfT4mfGf/wBEP1Z+AQP+YV/oT79U+vXtK/wjr3+yWfDj/vEz4z/UH/mQ/Vv1BuD/AMer9Qef9f36p9evaV/hHXv9ks+HA/7lM+M/4/5oP1b+OR/zCv4Pv1T69e0r/COvf7Jb8Oef+cTfjP6jdv8AjBHVvqJOok/79Xkk8/6/v1T69a0r/COvf7Jb8ORyPib8ZwTa5/0EdW82+n/MK/j36p9et6V/hHXQ+Ffw3BuPiX8ZgbabjofqwG3HF/7q3tx79U+vXtK/wjr3+yV/Djn/AJxM+M3Nwf8AjA/VnN+Tf/fq/k+/VPr17Sv8I67/ANkt+HPH/OJvxn4+n/GCOreL3vb/AH6v5v79U+vWtK/wjrr/AGSv4bn6/Ev4zfj/AJoP1Z+Pp/zCv4v79U+vW9K/wjrv/ZLPhx/3iZ8Z+f8Avw/Vn/2K+/VPr17Sv8I66/2Sv4b/APeJfxm5tf8A4wP1Z+BYf8wr+Abe/VPr17Sv8I69/slfw45H+ymfGaxABH+gfqzkD6A/79XkD36p9evaV/hHXQ+FXw2AAHxK+MoA4AHQ3VgABFiB/v1eLge/VPr17Sv8I67/ANkr+G//AHiX8Zv/AEQ/Vn+P/Zq/4+/VPr17Sv8ACOvf7JX8OPr/ALKX8Zv6f8yH6s+nIt/x6v8AQ+/VPr17Sv8ACOu/9ks+HHP/ADiZ8Z+RY/8AGB+reR/Q/wC/V+nv1T69e0r/AAjrofCv4cDkfEv4zDm/HQ/Vg5PBP/Hq/Uj36p9evaV/hHXf+yWfDj/vEz4z8m5/4wP1bybWv/x6v1tx79U+vXtK/wAI66/2Sv4b/wDeJfxm/wDRD9Wf/Yr79U+vXtK/wjr3+yV/DcfT4l/GYcaf+ZD9WfT+n/Hq/Tj36p9evaV/hHXR+FXw2Nr/ABK+Mpsbi/Q3VnBFrEf79Xgi3v1T69e0r/COuX+yWfDj/vEz4z8/X/jBHVvP1P8Azyv9Sffqn169pX+Eddf7JX8OOP8AnEv4zcWt/wAYH6s4sbi3+/V/B9+qfXr2lf4R17/ZK/huTf8A2Uv4zXta/wDoH6svb+n/AB6v09+qfXr2lf4R13/slnw4tp/2Uz4z6f6f6B+rbf1+n91bfX36p9evaV/hHXR+Ffw4P1+JnxmP556H6s/w/wCzV/wH+29+qfXr2lf4R13/ALJZ8OPr/spnxn/9EP1Z/rf88r/Qe/VPr17Sv8I66/2Sr4b/AF/2Uv4zX+v/ADIfqz/7Fffqn169pX+Ede/2Sv4cf94l/Gb8/wDNB+rPz9f+YV/N/fqn169pX+Ede/2Sv4cD/uUv4zc8H/jA/VnIta3/AB6v9Pfqn169pX+Ede/2Sv4bj/uUv4zf+iH6s/pb/nlf6e/VPr17Sv8ACOu/9ks+HH/eJnxn/p/zIfqz6DgD/j1fwPfqn169pX+Ede/2Sz4ccf8AOJnxn45H/GB+reLWtb/fq/4e/VPr17Sv8I66/wBkr+G//eJfxm+t/wDmQ/Vn1/r/AMer9effqn169pX+EdeHwr+G4+nxL+Mw/wBbofqz/wCxX36p9evaV/hHXf8Aslnw45/5xM+M/N7/APGB+rOb/W/+/V5v79U+vXtK/wAI66/2Sv4b2t/spfxmt9bf6B+rLf7b+6vv1T69e0r/AAjr3+yV/DcfT4l/Gb/0Q/Vn4BH/ADyv9D79U+vXtK/wjr3+yV/Df/vEv4zfW/8AzIfqz6/1/wCPV+vv1T69e0r/AAjrv/ZLPhxz/wA4mfGfkEH/AIwP1byD9Qf9+r9D79U+vXtK/wAI66Hwr+G4+nxL+Mw4I46H6s+hJJH/AB6v0JPv1T69e0r/AAjr3+yVfDe1v9lL+M1h+P8AQP1Zb/3lffqn169pX+Edd/7JZ8OLW/2Uz4z2ve3+gfq21/6/8er9ffqn169pX+Ede/2Sz4cfX/ZTPjPf+v8AoH6s/rq/55X/AFXP+v79U+vXtK/wjrr/AGSv4cfX/ZS/jNe9/wDmQ/Vn1+t/+PV+vPv1T69e0r/COvf7JX8N/wDvEv4zfS3/ADIfqz6DkD/j1foD79U+vXtK/wAI68fhX8Nz9fiX8Zjc356H6s+v9f8Aj1frz79U+vXtK/wjr3+yV/Df/vEv4zf+iH6s/wDsV9+qfXr2lf4R17/ZK/hva3+yl/Gaw/H+gfqy3F7f8wr+Ln/b+/VPr17Sv8I6/9Gg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5mz8pv6f6Ouvf9b/AI+XPf7D3R/LpTbcX629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3/wDCaAf8ZY+U5tx/o868F/6X3Ln7D/Y290fy6U23F+tvX230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/06D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaG/+lr5S88f6OuveP8f7y52x/wBh7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/1KD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoP+Zs/KYf9+768/1+Ny5//ivuj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/VoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzb/AMJof+ZtfKX+v+jrr3/3pc7f3R/LpTbcW629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/1qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR315/rf8fLnv8AD3R/LpTbcW629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9eg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pv8AxHXXv/vS53/Y+6P5dKbbi/W3v7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmhv/pa+Uv0t/o669/17/wB5c7b8fT3R/LpTbcX629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv//RoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8pv6f6O+vP/elz3490fy6U23F+tvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r//0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoR/xln5TH+nXfXo/2+5c9/j/h7o/l0ptuLdbe3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/06D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbXym/wDEdde/+9Lnf8PdH8ulNtxbrb39t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//UoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzb/AMJof+ZtfKX/AMR117/T/npc7/sfz7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/9Wg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8AwmhP/GWvlMP69dden/bblzv5/wBj7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/1qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR317/70ue/x90fy6U23Fvs629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/XoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhv/pZ+Uwtx/o768uf8RuXPWH+xv7o/l0ptuLdbe3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pvr/AMy669/1v+Plzv8AvPHuj+XSm24t1t7+2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/RoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8pv8AxHfXv+v/AMfLnvz7o/l0ptuL9be3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbXyl5/5p117x/X/fy53n6/j/AIn3R/LpTbcW629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/ToP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzb/AMJof+Zt/KXj/mnXX3PFh/v5s5x/Xn/iPdH8ulNtxfrb39t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//9Sg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaD/mbPymH/fu+vP6f89Lnv8AC/590fy6U23F+tvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r//1aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR317zz/wA9Lnv8bc+6P5dKbbi3W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//WoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8phx/zLvr3/X43Lnv94590fgOlNt8TfZ1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf//XoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzf/AMJoP+ZsfKf/AMR513+f+zkz/wCPof8AX/Huj+XSq24v+XW3r7b6Vde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaD/mbHym/wDEd9efn/s5c/8Ai/uj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf//RoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmg/5mz8phf/AJp315x/W25c9z/sL/7z7o/l0ptuLDrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//9Kg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaG3+ln5Tf1/0d9e/7b+8uev7o/l0ptuL9be3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//06D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/KYf9+769/3jcue/4r7o/l0ptuLdbe3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//UoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmht/pZ+U39f9HfXv9fp/eXPX/FvdH8ulNtxb7Otvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/9Wg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym/r/o769/p/z0ue/2Puj+XSm24v1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/9ag/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5mz8pv8AxHXXv/vS573R/LpTbcX629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv//XoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmgI/wBLPymH5PXfXh/2A3Lnr/4/n3R/LpTbcX629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaG/+ln5Tf0/0d9efj8/3lz1ufdH8ulNtxbrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//9Gg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pf6/6Ouvfxz/AMfLnfz7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR317/T/npc9/sfdH8ulNtxb7Otvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/9Og/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaA/8ZY+Uw/r1314fpxxuXP8A5/2Puj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/1KD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoP+Zs/KYW4/0d9ec8cf7+XPcf15v/ALx7o/l0ptuL9be3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//1aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbPymP/fuuvf6353Lnv8AYfj3R/LpTbcX629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/WoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8pv8AxHfXv/vS578e6P5dKbbi3W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//16D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbXym5/5p117x/X/fy57m/wDh7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/0KD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR315/T/npc9/sfdH8ulNtxbrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//0aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbXym/wDEdde/0/56XPf7H3R/LpTbcW629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9Kg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaE/8ZZ+Uwtz/AKO+vTf+lty564/2N/dH8ulNtxfrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6/9Og/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym/wDEd9ef63/Hy573R/LpTbcX629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/UoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhH/GWflMeOOu+vB+b87lz30/FuPdH8ulNtxbrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//9Wg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym+n/Mu+vf9f8A4+XPfTi9v9j/AMao/l0ptuLfZ1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/1qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoQf8ASz8pjzYdd9eD/C53LniL/wCPHuj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/9eg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pj+f9HXXv/vS53/H/D3R/LpTbcW629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv//QoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8pv8AxHfXv9P+elz34tf/AHm3+8e6P5dKbbi32dbe3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/0aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoR/xln5TH/v3fXo+n9dy57824+nuj+XSm24t9nW3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoQf8ASz8pjbgdd9eC/wDidy56wv8AXm3uj+XSm24v1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/ToP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhv/pZ+U39P9HfXv+3/ALy563uj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf//UoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhv/pZ+Uxubf6O+vLji1/7y56x/rcW90fy6U23Fvs629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9Wg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pv6f6Ouvf6/89Lnbc/T3R/LpTbcW+zrb39t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//9ag/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym/wDEd9e8f+TLnvzb8e6P5dKbbi3W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691/9eg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym4/5p315z/T/fy57j/Y/8R7o/l0ptuLdbe3tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPym/wDEd9ef+9Lnv8fdH8ulNtxbrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//RoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhv/pZ+U39P9HfXv+3/ALy563+Puj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Job/6WflN/T/R315f/AF/7y5635/1/dH8ulNtxfrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//06D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoQf8ASz8pjxYdd9eD683O5c8Rx/Tj3R/LpTbcW+zrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6/9Sg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18puef9HXXvH/ky57m/wDh7o/l0ptuLdbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/1aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaG/+lr5Tf0/0dde3/wBf+8udt/j7o/l0ptuL9be/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X//1qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaED/AEtfKY/n/R116B9Pody52/8Aj+PdH4DpTbfE32dbe/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/9eg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaC3+ln5TcG/+jvryx5tb+8ueuL/S590fy6U23FvXrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//QoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzb/AMJof+ZtfKb6/wDMuuvf9b/j5c7/AIXv7o/l0ptuL9be/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/0aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaH/mbXym/wDEdde/1/56XO/7D3R/LpTbcX629/bfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9Kg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pvr/AMy6694/H/Hy53/eR7o/l0ptuL9be/tvpX1737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3X/9Og/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaD/mbPymPNv9HfXn/Bedy561/wDHjj/Y+6P5dKbbix629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv/9Sg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaEf8ZZ+Ux/P+jvr0XvzzuXPX4/2Huj+XSm24setvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/1aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb6/wDMu+vOPx/x8ue/3ke6P5dKbbi3W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//WoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmhP/GWflMP69d9eH/bblzw/rf8AP9PdH8ulNtxb7Otvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/16D/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaE/8Za+Uy/16669NuPxuXOj/AF+NXuj+XSm24t1t7+2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/9Cg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaE/8ZZ+Uw/79314f9tuXPf8AFfdH8ulNtxfrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6/9Gg/wB+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbNv8Awmh/5m18pv8AxHXXv/vS533R/LpTbcW+zrb39t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//0qD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s2/wDCaEf8Za+Ux/p1116Lf6+5c7/h+Le6P5dKbbi3W3v7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//ToP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8ph/37vrz/wB6XPfn/Y+6P5dKbbi/W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//UoP8Afuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzb/AMJobf6WvlKfz/o66944+h3Lnb/4/j3R/LpTbcW+zrb39t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6//1aD/AH7r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8Jof+Zs/Kb/AMR317zz/wA9Lnv8Le6P5dKbbiw629vbfSvr3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuv//WoP8Afuvde9+69097e21uLduUpsHtbBZfceZrHCUuKweOq8pkJ2JA/bpKKKadgL8nTYD6+1NpZXl/OlrY2sk1yxwqKWY/YFBPRHzFzPy5yhtN1v3Ne/2e27JAtZJ7qaOCFB/SklZUH5nPl1YD0Z/Lv+QmY35sHMdg7Ch2vsKn3Ztyv3V/eHNYSHIttynylPU5aIYOKuqck8s9FC8fieNGu4BAFyJa5Y9nucrrdNruNz2gQ7YJ42kEkkYYxhgWGjUWqVBFCK54dc7Pff8AvKfu0bDyHz9tHIfuPJunPMm03kVk1laXckK3jwtHbt9T4SQ6VlZX1rIVohIbhXj3r/Lx+QmE33v3M9f7Bj3TsKp3buKt2r/dzMYWoyK7cqcnUVOJjbAvW0+Sjkgopkj8SROwKEAEWJ9zR7Pc5Wm57pcbZs4m2wzyNGInjZhGWJUaNQbCkCgFccOt+w395R92jmDkXkHZ+fPch9r56j2mzivWvrS7iha8SFY7hvqvCeCjSqz62kC0YVataEA3BtvcO0srU4PdGDy+3czRuUqsVm8dV4vIU7Alf3aStihnUEg2Omx/HuJbuyvLCd7W+tZIbleKOpVh9oYA9dEuXeZ+XOb9ptd+5U3+z3PZJxWO4tZo54XH9GSJmQ/MA1Hn0ye03R51737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3XvfuvdbN3/CaH/mbPymP/fu+vP8AW53Lnv8Ainuj+XSm24t1t7e2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf//XoP8AfuvdHA+JXxE3b8ntyzutRNtrrjb9TAm694NT+VvI2iX+BYCKQeGtz1RTNq9V4aWMiSW+qOOWQuQfb+/52vXYuYNlhI8Walfn4cYOGkI/JR3N5BsNvvfffD5T+63yzawpapu3ujuaN+7tsD6agVU3d2Vq0VnG4oSKPO4MURBEkkWxz1V051d0NtyLbHW+16LBU/jQV9akaVG4M3OioPu9wZuWM1dbO+m4Vjoj+kaRrZRllt9ly5yVZ/QbPZLCKDVShlc+sshFanjp8vJVHXzz838x+9P3pOZ25u9yeZptybWfBD6l2+1Uk/pbfZIwQKtSvi1BelXmmap6W2RydRaJIoogXYk6kMpsvC3MhYcsfbQ5qrI2lEA4epqfmT0dP93nw7KFZ7i5djVmoVjUBR5LGq0H2k8OPXsbk6gLIkscXobUNKeLgmzf5sqLhh72eadMi6kQjgfI4+YI6bH3evEs5hBcXKOKMtSJFIYeayK1R9hHHj0h+1umOrO+9vSbb7J2rQ52MRSJj8iY46fceDldWH3e385FH93RzKzaigJikItJHItwXtwsOXOdrP6LebFJjpOkmglT5xSDNRx08D5qw6I+T+afe37rPMg5q9tOaZ9tXxFMyJqbbroAj9PcLFmMelqafFGUr2TQtQ9a4nyz+I28PjBuanE877j673DUTrtHeMcHi8rIHlOEzsKAx0OfpqcaiqnxVEYMkR4kSPEvn/2+3Dki8jYyePs07HwZgKVpnRIPwyAZpwYZXzA+h37oH3xuT/vUctXsSWY2r3O2qNP3ltrNq0hu0XVo5zNZyP2hiBJC5EcwzG8hQ/ce9Zlde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691s3f8JoR/wAZZ+Uxv9Ou+vRa5/O5c9zb6G1vr9R/sfdH8ulNtxfrb29t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6/9Cj7YGyc32RvbauwtuQ+fN7tzmPwWPUqzRxy19QkLVU+m5Wlo4i00zf2YkY/j2YbVtt1vO5WO1WSaru4lVFHzY0qfkOJPkAT0D/AHA542D205H5r9wOabrweXtnsJrudvPRChcqo83cgIi8WdlUZPW2r1Z1ptfo/rjbXXOz6VKfHbfoUp/P4o0qstlGVXy24MiVFpq/JVhaRibhSwVbIigZy29pYcm7DZbPtqgRwrQGmXf/AESZvVmbhXgccFA6+U7cN85u+837vcze4/PMrveblOJZU1Epb2tSLPbID+CGKIDXppqUFj3zuxWLGnp6aqyGRqoKKhpIJqytrq2eOmpKSmp43lqKmqqp3SGCGCJS7u7BVUEk29xbvO9lS7tJ9pJ/meuhntl7Uo8dpBBZAsdKqqr9gCqoH5AD8uivbD+YHUvbXfm0eheo6HcHZGe3Bkqimq90Yikhptn4agxdHUZLK5Zq+tmhrMlRUNNSsS8MHikNhG7lluD9k5ph3HmHa9ugDyI06l2GAEU6mOeOAesifdf2MvuSvZnn3nLdJLe0u4tqmjto3qZGuJ1MMC0UEKTLIppWopkAAkR97fMjqHqzv3efQva9HuDrnPbZy8eOh3NmKOGo2jl6XIUdNksXlRXUU01bi6HI0tXG6PNB4kB/cdLNb2881wbfzDum3TB40WdijH4SrHUpxwGkjrftd7D7jzj7Ncgc57bJb3d3PtEC3EaEiRbiFBDOtGADMJY2qK1rwBFCTUgwTQU9dj6mGso6qGKqpK2jmjqKWpp5lWWCopqiBnimhljYMjoSrCxHHsabNveoxsstRggg/wA69Y0+5/tVGkd5BPZAMNSsrL9oKsCPyII+R6RXbXWG2O9us9zdb7tpkkos9RNHDVBFaoxGahVpcNuChaxMdXj60K5tYONSNdJGBlC6sbDnTl+92bcQCkq0r5pJ/ocy+hDfFTjkHDHrnvsvMvN33W/ePln3K5JkdZtvmLrHU6LmzJAvNtmz3RyRZi1VKEK698CMNSrfezc113vPdGxdxQiDObSzuSwOSRdXjapxtVJTNNAWAL01SEEkTW9UbA/n3gtue3XO0bjfbXeJpureVo3H9JSQafI0qPl19XvIfOmx+43JXKnPvLNx4vL+87fBeW7eZiuI1kUMPJlDaXH4WBB4dJT2h6FnXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Wzd/wmh/5mz8pvr/zLvr3/W/4+XPf7z7o/l0ptuLelOtvb230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/0Sgfyr9k0u4vkJmN1VkUcqdf7Ey2UoRImrx5bNVVHt+nmQkFUePH11XY/Xnj6e5p9jNtS65sutxkUEWdo7r/AKdyIwftCs1OuX/96/zvc7D93nl/kqznZH5m5itraWhpW2tlkvJFNOKmWGAMOBBocHrYGqqmI1hi8gJjtFpX1G6j1XANgS5N/wDX9zfzZNL4siAdqjT+wZ/aanrmB93HarD93WVzIyma4cykAVPe1EBp6RhF/Lqq3+ZV2pvrKVGwfi31Zj8zl892NDDmtx4bbdJV5HP5+lfIvSbb2zSUVDFLVTU9fX0M1ROiC8ggjB9GsHGvm+9uZ7mLa7YM0j5IWpJzgUH2VP5ddwfu98t7Js+xXvPO9yRQWltVUklKokdFrJIWYgCgIUHyJamadW7fyhf5W2X+Me1sn233FQ0UnefYOMjx8GGpnjro+udmzfaVpwD1aK0M258rWwq+RkgZ4YliigjdtMryLeX7QctI91dEfvGQU/0i8dNf4ifip8gPOoG94+YZvfG4tdg5fWQ8m2kmsEgj6qcVUS6ePhopIiDAElmcgEqFhfzef5XeT+S+3Mf2509j6CHvTYmNkoKvETvHQR9j7PpvuqtMF90yrDFujE1czNjpJ2SKVJZKeR1BiaN7mDb15miju7On7yjWn+nXJC1/iH4TwOR6ULPZ7nB/Yi/vOXOZfEHJN3LrqASbWc6VMunj4TqAJlUFgVV1BOoNUx/LV7W3nRLv74xdn0Wbxm4utkly+3MRuGmmocxgsaletDuPbFTSVscNbCmNylZDPBG4OhaiQCyKg9oOTr25guZdquQyyx5AbBGcrQ5xx/b1IH3i+W9k3fY7DnvY3hnsboBXkiZXSSq1jlDKSpqAVJByQvnXq0+hq4PvBEJLNJeMq5KNdv0jng2cA/7D3kvynNN4kaGulhTHz4fzp1w1+8bte3nb724RlE9u4lAIoe00cCvrGXGPI9a9f80jZFNtn5Iw7ko4Uhi7C2RgtwVWiw15XHS1m26uQgE8vTYenYn8sT7gX3y21LPnNL6NaLe2sch/061jY/adAJ+Z66wf3UfPFxzJ92a55TvJ2km5Y5gvLGOprS3l8O+hUeiqbmRFHkqgDHQW/A7pfYXfPe67D7Hx9bktuHZu4sz9tQZOrxM/3+Okxq0shqqJ45/Gq1L3W4BuL/T2Re1XLe1c1832+z7zEz2LQSsQrFDVVqMjPUsf3gvvf7g/d7+7ju/uR7ZX8FtzTDutjArzQR3CeHPKVkHhyArUjgaVHl0Kf8w34jbY+OOd2PuHrPHZSk683bQVOLqUr8jU5c43d2Mkeokhatq2eoSPLYmdJIUcm7U05U2FgIveP27seSrza7vZIpBs1yhU6mL6ZkNSNRzR0IKg+av5DqF/7tD76XNf3pOWufuXPdTcLSX3N2S5SdTDCluJ9uuFCo4iTsLW9wjxysoACzW+oamJMr+Xh8QNqfIzI773R2hjcrV7A2vTUuExsWPyVTiP4lu2veKtlH3tGyVLx4bERhpY1KgtWREnixc9nPbmw51uN1vt9hkbZ7dQihWKapmofiGaIgqw9XXpH/eY/fY5u+69s/t9yp7T7jZxe5W8TyXMrTQpciDbYQ0dfCk7Q1zcsFidge22uABWjAM/mD8eMD178oKbpTpXb+UljzON2bDgcHLkKvL11XnNxRlDGtXWyPKElnZf1NojUEkhQbB73E5Qtdn59k5X5ZtJCriARpqLsXkRTSpzlj54H2dTH9y77yO/e5H3RLP3898+YrOO4t33SW9uxFHbQx21jcSqGMcYCjTFH+EanNAAWOTrUfwf+Ivxp2hh838v+ymym6s2geLAYnJZrHYvyQapaikweK23TPvDNQw3WKaucwQaiBoiLLeVI/az275I2+1uvcffC+4S8I0Z1XHEIsYMzgcDIdK1oKLUVwAvfv8Av30fvVc577sP3JPahbXk/b8Ne3UNrLcUcgJJcz3si7bbO2XjtFE8xQM2qUK2nHlfgv8AFX5IbGzW6/hz2JNS7kwisW23lsnlK7EPVSRvPS4rMUW5KWLd23pKzQ0cFY7z050n0SWZhW49quQeddrutw9td6IvYuMTszJUglUYSgSxlqUVyWXBwaEhzZ/7wb73/wB1zn7l/k/78fthHJytuBAW/toYIrgIrKktxbyWMj7deiIMHltUSGcal749SIxRvg/8btqdsfIfePUndmAy6LtTZ+6anIYaDJ1WGr8dufAbl27hpYaiooXDuKYVtQjIGKMxDAkAXjr2t5KsOYuctw5d5ntZQILeUsgYoyyxyRoQSPSrAj16zU+/396PnD2Y+7Pyb70exHMFhI+771t6Q3MkCXUM9jeWV5cq8aSUA8TwoXViNQWq0BJoA3yo6+2z1V8guz+vdm01TR7Z2vnKegxFNWVk2QqYqeTEY2rdZayoJmnJqKlyCxJAIH49hTn7ZrHl/nDfdm21GWxglCoGJYgFFOScnJPWQf3Qvcvmv3i+7b7T+5nPFxDLzXu+3vNcPFGsMbOLmeMaY17VGhFBAwTU+fR4NjfE3pbO/wAvrL/ITI4bMydmUe2N9ZOHJR7hyMWOFbgt5ZnDY2Q4hHFEY4qCjjVlKkOwLHk8SjtHt7y1ee0F3zjPBKd8SGdgwkIWscrIvZw+ECo8+PWBHuL98v3z5b/vI+X/ALtG17tt6+1Vzue0wvE1nG1x4d3t1vczj6gnxATLI5VhTSCFAoMrj4afBnpX5C/Gio3luinztDv7JZbeODx24qHNVqUuLnotEOHrWwgdaGsWimlDyRvbzAFSwuCDH2x9rOWudOSrncr8TJuxmmjSRXOlSFXQxTg2ktUj8QxUdAn79/8AeAe+P3YfvRbFyTyi+3XPt3Ftm3XtzZzWyGW4WSecXMS3VDJCZY4tCOoPhMdYVqEGq7srq/dPUXY+d613xQNQ53b2WWgqgpJp66kldJKLKY+YcTY/KUUiTQuOdDgEBgQII3vYr/l3errZN1h0XcMmk+hByrKfNWUhlPofXHXW72t92eUPej2w5f8AdT2/3IXPLu52ZmiJFHjkUFZYJlOVmglVopVPB0NCVIY2cfL/AOLnxg+Ou/fjc74zcOH683juPci9nSS57O5qqk29hxtdlNCIfLkKeaJcpNqNONb6hYcD3OPuNyFyLyXvHJWqKePZrmaT6omR3Php4Xw0GoHvb4cn8uuUf3Kfvdfex+897b/ekEW4bVee5uxbbY/uJUtLa1jF5c/X18YuwhdWNvFQSkKtDUnV0N3Sfxs/lt/IXJ53EdV0m9M7XbcoaXI5aOrzG/8ADLBSVtRJTU7pJllpI52aaMgqhJH1PsV8r8j+ynOM93b7At1LLAgZwXnSgY0GXArkeXWP3v396r+9J+7RtXL29e78+wbft+6XEkFu0dttN0XkiQSOCtu0hQBWBBYAHgM9AV3Fs3+WLsii7P2hjJ960vaO1aHeWBxdDLP2LVUkG+8PS5Ghx9NJUvTPi56ZM9Ail2cwMvJOm59hPmTbPYva4t926CW6XfrdJo1Um4IE6BlUE6dJHiAZrpI86dZDeyPPf97Fz9fe03Om7WXL8vtJvFxtl3cSquzRyNtNzJDLM6xiUTo5tHZgoUShsBdVB1XZ8fOpa7vLuTYXWNEKlIdyZyBMzV00ZeTG7cog1fuDI3sUjalxNNKULWUylV+rAGGeUOXpuaeZNo2KGoE8oDkfhjHdI3+1QMRXzoPPrpx95D3k232A9kPcX3Z3ExtJtO3u1vG5oJ7yUiGzg9SJbmSJWpkIWbgp6tQ+YXwI6T2R0VuzsPoinr23D1xnKWTd0D7qqtxxjCU4FNuDHT0s0s60OVw5yVNWyAlHjpopLg6l9z37j+0nK+1cqbhvXKWtruymHjAy+IAgxIpH4XTUjkYIUGoyOuQ/3Jf7xb379wPvBcne2P3ifpouXOadtc7Y67eLJmuX/Us50cU8a2uRBcWyMAyPO0elxpatRPU238Zu3tTrPaubilnw25uwNm7fy8MEz0002MzO48bjq+KGoj/cglkpKlwrryhNxyPeO2wWcG4b7sthcgm2nu4Y3ANCVeRVah8jQmh8uu0Xu/zHuvJ3tL7o83bFIib3tXLu5XluzqHRZ7azmmiLIcOokRSynDCoPHo8v8xT419U/HPcHVeO6uxWUxlNurFbnq8yuSzddmTNLjK3CwUZhatd2pwkdZJqCmzXF/p7lX3l5H2Dkq82CHYoZES4jlL63L1KsgFK8PiP29c+/wC7J+9X7wfei5b93tz92tysri52e8sI7b6e1jtgqzxXTyaxHh6mJKV+Gh9ejU/KP+WnsTE9Ozb3+P8Ais7Du7bFJ/Hsvt2sy9fnn3PgVo1myUGMjqmeSLMY5EM8EcQvUqHiCl2jsPee/ZDa7blo7tyhDN+8IE8R42cyeLHpqwSuQ6/EAPiFVpXT1iL90v8AvVef9798U9vfvH7ntp5N3W4Nna3sVvHaCxu/FKQNcFe1rackRSyOQIG0SswjEp6r4+C3TWxe9u+6LYHY1DXV+3J9rbkyklNQZKqxNT97jIad6V/uqRkmCIZGutwD+fcO+1nLe1818322zbzG7WTQysQrFDVVqMjPHrpb9/8A97+fvu9/dw3v3L9tL22g5og3OxhR5oUuI/DuJtEg8N+0krwPl5dBx8p+vts9V/IHs/r7Z1NU0e2dr52CgxFNWVk1fUxU74nHVbrLWVBM05M9Q5BYkgED8eyTnjabPYubd+2jb1ZbK3uCiAnUQABxJyepU+6n7h8ze7P3dfaP3I5zuIpeaN42iO4uXijWKNpGZwSsa9qCgGBjov8A7CnWQXWzd/wmh/5mz8pv/Ed9e/6//Hy578e6P5dKbbi3W3t7b6V9e9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691//Sbv5afcGyN5Sb169wPSmw+udxYfbFBm63d+1ajMzZHddFT5OLHtQ5c7grMvkrU9RXJMumraHU7WjQW95R+ynMO2Xs+77dbcuWtncx2isZIi5aQK4B1eIzniwOGpXgo64G/wB6J7Nc88rbZ7c86b371b9zJsd3zDLDHZ36WqQ2Ly2skiGD6OG2ioViaMaoPE05eVzUm0QwywZSZ5ItaGd2F7spBYkX0EEXB9jTml0W5nJRXox7WqAf2EH9h6x4+73aXM+w7OkVzPbF4EAliEbMppQkeIkqVBH4kPQeVvzv6O+JvyJ2rt7uPqsY+n3htOgqKHvnD4rH5XL4aimymSoarBZOmWhG4n27jpoVllNJUTun3BP2z/X3EW4c47PY3yWlxsENqWSnjRCvmfi1Venr3H7OuhvK/wB2X3K5p5Wl5i2j3i3Lf4orot+7twZYwNKIR4Pg6LXXSoTVBECOMgNethXAZrA90dL15637IGKx/Y+zy20Oz9m1MORNFT5SGOWkzmGqaaohEwaMFG0SxSqrMAyOAQksr+32vmPZOYLjaody22C4SVreQ/pToPwsaMKHiKqwqBqUio6kzY7eW55M5s5Hk3m52DmO4tZLYXUakXFlKSAXVdSMGFCp0ujaWYo6tRgybny+C6X6boG7I7ITKUHXm0h/ezs3edVDjTWw42J3qs1mKmqqJhF6SEXXLLKwVQWeQ3K663K13PmDeuYINqh2zbJ52lWCM/pQIfwKaKKDiaBRUnSoFB0E+Z7C6h5W5Y5Hh3e65g5kgto7Y3LqWub2WpAdl1O2o1CjU7tpUF3Y6m6198V84+i/lb8h914XqHq3+JLszaldNU97ZfHY7EZTM48ZTG0UWCx1KaI7j/u/kKh/LCaqohc/b3NOn19qLHnXY9yv3tLfl2K6kRKePLVaZHw6aSU9Ksv2dR/zD91v3Y5M5Uh5k3f3q3Hl61nutR2vbyswYlHP6pnL2gemJNNvLStBIT0J8UE0+UheODxoJ0c6dQUAOCbF2LEAD/H3LPK0iNcQFUVAWHatSB+0k/tPXPr7wVldQ7HuySXM10VgcGWURqzGlASI0iSpJ/CgHy6q5/mY9ubK2dV7M2DnOlNgdh7izG0KrL0O9Nz1GWiy216ObL1dCmPxgwVTicj42no2nGqsEJf6xML3A/vdzDtljPte23HLVrd3UtoWWWUuHiBkYAL4ZRuIJy+mvFT1kr/dY+zXPXNW1+4HO2y++HMHLnL9jzIkMu32CWz299IlpC7PP9ZFcxDtkWM6bcSaQNEqGhBTP5VzB/lZrVFiV+vN5MsSaikYafDkImtnfSgNhck2+pPuP/Yg19w7UgUH00//AB3rML+9vUp9zHfkZyzDfNqFTSppMcmgAqeJoAPQdWI78ytR8xMb8wvi3m6rGv2D1ju9c91N/ktFikOMxtHQjAQTTxgLNJDm46ihrqyRdYp8qhYsfczbpJJ7iw+5XIN1KrbvY3Rms60XtAGha4wr6o2Y8EmFeHXMfkOztPuV7n9xz732w2MsHtvzZsC7dzIIzJKDM7MLmbRUnVLb+FeQwJh7jbpCAC9OnDYWXm+IzfDn4n4CooP74dg5jJ7g7XmNHR5FJKKrxGYqcxDBUzKwhep3NIlPR1UQ8opcVbUoazO7XK3t83tp7eWkirud3MZrwihqCrFlr6NJ2Kwzohp59F/P1hD98eL78/3zOY7KWbkblzbY9s5bSQyRhXWeBIZtAKkPDZE3MsL1QXO5aiCUwH2d/gX/AA7dtj+N+Hzf6Nk/gHn8fj/jv9xs14LeT0+b7H7jxW9Xl025t7J5vpP+CLh+qpq+lHh14eJ9Kafnp1U+dKZ6kzbP6x/8mVdx/q/r8H9+v9Zorq+j/f66+GdPi+D4nl4evV216Ij/ADSm3K3yqyq5pakYdNnbSXaBl8v2z4c48tXtSamMX/HxtWrJoA9S8/1MWe/Jvf8AXAuvqg30/wBPD4Na00aO7T5U8TXWnnXrPz+6NXlYfc65fOxNAd5O9bl+8dGnxBdfUERCag1avovpSmqv6ZWhpgPn8qF9xj5M1yYpqwYN+utxndKxE/ZNSLV4n+FmsBPjMq5hovDb9zl7enX7U+wBvv69kWpb6Q2cvjemnt018q+Jop58aefRJ/fCpym33Slk31Yv6wLzJYfu0mniC4Im8fR+LSbIXGv8OErnR0eDqdsGv81nv1cM1Osb9VgVogNkfNii6rbLgA/qqBVhzKF/3YHP4PuU+XjaD7wPNotSuk2JrTh4mm11/nqrq+desAPeZOYX/udfu6Nv6ymdea1MeuuoWnjb8LUmv4PBMYi8vDMdMU6qi+eKMny57vV1Kk7nonAPB0ybcwsiN/rMjAj/AAPuAfdn/p4nNP8AzXH/AFbTrsL/AHeNP+Au9gKEH/dRJwIP/Ey59PMcCOIODkdWd9cEUP8AKLzzVv8Akqy7K7HERn/bEhr+yc7FRaC1tX3ck6LH/qiwt9fc57H+j93a+MvaDbXNK4rquHC/tJAHqeuT3uv/ALsv76TlWOwHjPHveyagncV8LZbZ5a04eGisz/wgEnh1x+Ju69wbF/lodj7x2pkXxG5NtN2bmMLko4aed6PIUT0ssE4hq4aiml0svKujKR+Paf273C82n2R5j3Lb5zFewS3DowAJVgsdDQgg/mOjf76PJvLPuH/eo+yXJHOW1JfcrbpY7Nb3UDM6rLFJLfBlLRsjr6gqwIIGem7f2z8D/Mo6G2x2311R4jG/InrlqDB7l2/JW0+Oirb1Anr8DU1dVM60+IrC82QwtRUNpUmWB3BMrJ7eNutPezlKy5j2aONOc7HSksdQurNWQkk0Ru6SBmP8SE11Ee9tuc+Yf7rP7xXNHsn7nX15c/di5qM11YX3hvMYDp0RXapGi67iICGz3WGJakCC5jQqIkkYv5xilaH4+Kfqs3Y6n88iHY4PPtL95LCcnV4/4z/1g6Pv7j8hrv7yxU1U/uT/AI9uvSI/k8f8zB7n/wDDN23/AO7us9lv3bv+SxzN/wA80f8Ax89Dn++8/wCna+xX/S8vv+0WPquf5Qf9lH97/wDiXOwP/eoyfuFuev8AldObP+ljcf8AV1uun/3Tf/EX/u8/+KXs3/aBB1aP/Kw6jq9tbM7S+Slbt+tzGQbEZXbHXuMo4/LkcxTYWF8ruVcVTkqs9RmMtSUmPp2DAmanmTgE+539hOXpLHbOYOeJLJ5ZlieK3RRVn0DXJoHmXYJGpHmHHXJX+9495LPmnnn2c+6pZ8zW+37XLfW1/vNzK2mC1NzJ9LYm4fgkdtA9zeTqwP6b28nkOhZ+BOI72q8p8htlfIXrDe+G2127Plt9z5DcuJyGOw82a3HLJht34SKauYv5czjMjTmFFuUhoXueF9iH2jtOa2uOcdp5x2G6jsdy1zlpUZUMkhKTJU+ciupA9Iz1Dn94zzB93u32f7tHuF92j3b5fvObORxbbTHBY3UU90lpZqtztlzpjPwWc9vKHYjMl2nqa1M0PV+S6Y+aG0escoJDNtLvvYdBRzykF67DTbxwlbgMkSqRqf4lhKqnn4UWMlrC1vePcWxXHLPuVt+xXNddtu0Cg/xJ4yGN/wDboVb8+uyd97s7R76fce5w92NkKi03r283WZ0Br4NwNtuo7u3Pzt7pJoSfMpUYPR8f5xH/AB9/RH/ag3t/7s9ue5c+8j/yUeU/+aM//Ho+ucn9yD/ypX3hv+lltf8A1YvejyfLD5R5j4w5f44ZZoo63YO7MhuPG9iYxaKGfIz4qjxW21ocjiath9xTVmEmyUlQIkISqUGN/qrLJ/uDz7dciXXItxTXs9x4q3KUBYqqw6WQ8QyFy1AaN8J8iMDfucfdD2D72vL33sNo8Q2/uRsy7fNstyZHWGOeafczNBcRg6Hhult44i7KXgNJY+DI6S65+L+J2l8vtt/I/p8UGR6a7R2HujL1rYmWkTH7fz2co6CrparGwaoHfA7rjn88CRI5p6gTI4jjMQBfs3Idvt3uRtnOvLSq/LF9ayu2igWJ5EBBUY/TlqCoAOltQIVdPQy9y/vc7xzn9yHnn7rvvjJPa++/Ke/bfaxC5WRp761tLkpIkz94+s28xtHO0jIs0PgSI0svjHql352f9lb94f8Ah1U//uhw/vGP3Q/6eDzX/wA9bf4B13e+4T/4hz93r/xXYf8Aj8nRSvYC6y662bf+E0P/ADNr5Tf+I669/P8A2cud/Huj+XSm24sOtvf230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r//Trd+AvaVJ1X8mNlVmWqko8FvBK3YWZqZXMcMEe4/CMVNKf0iKPcNLR62ayohZibD3I3tTvkexc67XJcPptLnVbuTwAlFFJ+Qk0E+gr1hR/eD+1d77rfdd58tNmtTNzFspi3e1RRV2awYyTog463tDcIoGWYqvn1s51sbRyrL/AFGhrcWdAAb/AE/Utj/jz7yg5r29m/WpxFD/AKZcfzFD+3064Ufd152gt/8AdaZqaH8RM8YpWLCn+kkLofSi+o6KF8zfi/TfJ7rSnp8I9HRdnbLepyGya+rdKemyCVQi/im2cjUkN4aPLJAjRSH0wVUaMSIzJfHrmvYHvYy0YpcJXT8/UH7f8PXZ32C92rfl+dYryQttNwAJQMlSPhkA9VqQR5qTTIHQNfyg/k/3B8WO7qn4d92Y3P7b2d2HU5Cq2bjtzQ1FMu2ew4aU1QiwdVIrUlThN6UVI0ZEEjwPXJDJEbyylyPkqaaPdE5f3CMrHMTo1eUlOA9Q4FMfipTiepA+9JtVhd8h3Pu/yfcpLebYim6MRB8W01aWdwMh7YsHJIBEXiasItG7+bp8m+3Plj3xB8QelcVuPcOzutchSHdWN28lRLBunsWWmWpmlzk8Wiip8JsulqlgT7qRIYq01ErkFIinudbiafdG5f2+NjHAaOF/FJ5g+VE4Z4NWvAde+65sm37ZyJb+8POF1HHf7tEWtjKQPCs60VkByXuaa6qCTEYwvxNUc/hr8Xqb4w9b1VNmpKLIdmbykp6/emRonE9JRJSeUYrbWLqWSNpqPFJO7SyWtNUyuw9AjsIOUuX3sU1SCtw9C3y9APs/meo0+8D7u2/Mc7RWbldotgREDgkn4pGHkWoAB5KB5k9HBoImklaX+noS/wDq3BFh/wAEQk/7b3kRypt7L+rp4Cg/0zY/kKn9nXFr7xfOsM+rb/FB1t4j54RREMa/6dwiD17qcD1rF/PrtSl7W+TG9q3F1KVeC2etHsDC1EUnkhni235kylRCw9Bim3BVVhQrw0ek/n3iz7r77Hv3O26y276rS3pboa1BEQ0sR8i+sj1BHXej+729p732k+6zyBYbzamHmPehLu92rLpdX3BhLDG4464rQW8bg5DKw8uufwM7m2D0R3wu++yMlV4rbg2buLD/AHVFjK3LTff5GTGtSxfaY+Kao0OtM9306VsL/X3v2p5j2nlXm+33fepmjsVglUlVZzVloMKCemf7wn2S9wvvBfdv3f249sNshu+apt1sZ1jlnit18OCUtIfEmZUBA4Amp8q9PdV8s166+aG/fkV1kkm5ttZ3PZ+2Jrnq8BHufbmZoIoDR13mo6iroRFXwQ1SaoSwlpkJHtbN7g/uf3M3bnPYP1rOWaSitqQSxuoFGFNQyAwqMMoNOgztv3OB7l/cW9u/ux+7x/dfMlhttoGng8K6exvbWZnEsLBvCkLRtJA+l6NFNIobNesmA+WMG+Pmxs35H9pCXbe3MXmaPyY2h+7z6bZ2/jtv1uNpMfQLT0kFXWo1bUNM5WFWaWokfTyfdrL3BG6e5+2c68wHwbRJRVV1OI41jKqqgDUcmpoMliadM81fc5k5D+4bzz91z2cUbnzFdWEgSacxWrXt5PeR3Es0zM3hRnw1EaanIWOKOPVgHqF8wvkNgOwPlDTd1dKbiySxYXG7MnwWcfH1uHraTObdjLs60lfFDOUinVQdS6JFJBupN0/uLzfa7tz9JzRyzeOEQQNHJpZGDxoorRgDhh5ih+zo6+5Z92/f/br7oNn7Ce+vLVq1zcvukV7aCWK5iktr24mYKZImZDqikB7W1I1DUMMHYT5p/Dj5QbPwWH+XWwqrb+78HDLEmfxmOz9fjYZZhGKmq2/mdqSPuvG09e0KySUM0MsKOANcpUP7lMe53ttz5t1rbe4uztDuUQI8RFkZRXiUeI+KoalTGysoPm3Hrn+/3EfvwfdI5z5g3v7lfuTDufJN/IrmyuJrSGdgpOiO7ttwUbdO8QYol3FLDM617IgxQ4Mh83Pid8ZNh5vanw52RUZjc+eRi+58tjs3RYuOtSEw0eTzldutk3Vnv4f5naChWKCmDavXHrbXWb3R9vuRtpu9u9tdraS/m/0V1cKGpQO5l/Vk01OmPSqVrkVNX9s+4R98f72HuJy7zl9+Pn2Gz5R20im3201rJO0RfVJb20e3D6C08bSomuzLNcFdA0uUTw6wep/kJvvqvu6g71gq23Duo5vJ5XcqZOd0TdUO4Gn/ALxUWQmhQ+I5NKp2V1QiCYI6odAX3BXL/OG7bBzTDzWknjbh4rPJqP8AaiSviKxHDWCc07WoQMDrrP7x/dt9vPd/2D3P7vdzZLtvJ52+C2sjAoJ29rMJ9DLCrEavp2jQFCw8WPXGzgSMerVN2fID+Wh8kaug7B7p2/n9rb9jpqZMpTzYrfEddkfs4khhoa7J7AWrx2bp4IoFjhmn8Mwi0r6ANKz7uHN/sfztLFvPM9nNbbuFGsFZ6vpFArNb6lcAAAM2lqUGBgcgOT/u4f3q/wB1qxv/AG09iOZts3r25aWQ27rPtRjt/FYu0sUG8GKa1d2ZnkihM0Pia3GtjrYu3zJ+cOzuzOv8V0B0Dt6p2v1JiRi4cjWVFCmGbMUWCMMuGweHw8M0j0GApKmFJpGqLVFRNEl0QK3kBnuT7p7bvuz2/KHKNi1vy5HpDEqELhKFERATpjUgMS3czAYWh1ZO/ce+4Dzt7Te5O9feP+8dzVFvPvXem4MKRytcrbSXYZbm6uLl0XxryWNmhRYh4MELyAPKXUQy+oPlF05tH4Gdj9EZzO5Kn7I3HRdgQ4rFRYDLVNFNLnhD/C1fLQ07Y+ETFDqLONH5965a555c232n5g5Uu7p13q4M+hBG5U6wgWrgaRXSeJx1v3x+6j71c7/3h/s394Xl3YLWX2v2ZNrF1cNd28cqfTS3LTabdnEz6RKlNKHVWg4GhSvil8jM/wDGntfFb0oGqavbGQMOH33t+KQiPNbamqI3qHihZhCcviiDPRyMLrIClwksgIC5A50vOR+YLfdICzWT0SeMHEkZOflrX4kPkcV0swOXv3wPuw8s/er9nt45C3VIoOabcNc7TestWtL5UISpHcbecfo3MdSGQiQKZYomU138x35P9PfIym6jTqvOZLMSbTm3q+cWvwOVwop0zUe1xjzE2Sp4BUmU4ybUEvo0i/6h7kD3q565b50Xl39wXbyG3M+vVG6U1+FppqArXS3Dh+fWHP8AdefdM97fuwTe9B94OX7WxTeV2v6Uw3dvdazam/8AG1CB30afHipqpqqafCekr/Lo+RfVfx23d2Vl+08xX4ei3JtzCY3EyUGFyWZaeqosnU1NQkkeNgqHgCQyKQWAB+gN/aH2Y5z2Dkzcd7ud/uXjingRU0oz1IYk/CDTHr0LP7zz7sXu/wDeb5L9q9l9odjtr6/2vdLqe4E11BahI5YEjQgzugclgQQpJHEinRO+8Nz4TfXdfae79vVjzbd3Z2JuzPYeuqaappXkxOYz1bW0VTPRyR/dwMaWdWaMoXXkWv7jXmm+td05m3/cbOQtZ3F7NIjEEEo8jMpIORg8OPWcPsDynv3t/wCxPtDyTzLZrHzLs/LO3WdzEjpIq3FtaRRSosinQ4EiEK4bSwoa06s77L+d2xOpfj31d1J8Qd21UW4ttDHY/cG46zZUlGgoaKgnnzNfDR7ox81HLkN1bjqTUSHxytGpkF1LKfc5737r7Vy7ydy/y57c7iwvINIklMNMBSXIEqkFpZWLGgNACMVHXKD2s/u9Of8A3o+8t7w+9f31+SoH5a3Uyy2NjHuYkPiSyqltG0m3zpIsW3WMSwKC6iRijUbQ3RdNk/zKflDjN47WyG8Owjn9p0e4MTUbmwi7T2bTPlcBHWwnMUMNTRYKkqoKiox/kWN0kQrJpN+Lewbtnvbz7BuNjNuG8+NYLMhkTwoRrjDDWtVjBBK1AIIIOesm+ev7rH7ou78l82bXyd7ZDbebbjbrhLK7/eG5SfTXbRMLeYpLePG6xy6GdHRlZQVIz0JXyw79+N3ZnyJ6C7x663Nlambbm59px9lw1G1M1jWjwe0904zNY7N04qqSFslXJQy1cEkaapGSCBV4Hs69wubeSt85z5T5q2W+dmini+pBideyKVXWQVA1Np1KQM0VAOou+5p93b70XtT92P7w33f/AHR5VtIodw2vcDsTJuFrODc7jY3FvPaP4cjeBF44gmR2omqe4ZqHil/5jHyR6n+RO4eqsh1ZmshmKba2K3PSZp6/B5PCmnmyddhZ6MRLkqeneoEkdHISUBC2F/r7R+9HO3L3Od5sE2wXTyJbxyh9SMlCzIRTUBX4Tw4dCT+7B+617yfdj5a94Ns93titrG63e8sJLUQ3UF1rWCK6WUsYHcJQyoAGILVNMDp+/mG/KDp75B7a6cxnV2dyWYqtnVG5nziV+AyuGFOmTx+3KejMT5GngSoLyY6UEIW025+ou/7x888uc4WnK8Ow3TyPaibxNUbpTWsIWmoCvwNw4U+fRV/dn/dQ96/u18xe/e4+7ewWtlbb+23GzMN3b3Wv6eXcml1CB2KUFzFTWF1VNOBo9fAP54YHpDBZfq3ufJ5Rdg0yy5bZGWpMbU5ifAV89Trye35KWhilrWxeSeZqmFgGWCdZBwso0rPaL3YteVba42HmaeT9zirwOFLmNiatHRatoepYcdLV8mqAz/ePf3eG/wD3gN82X3d9jNrtD7luVt90t5Jo7ZL2BU0wXYklKxC5twqwyBipmgKGuqALISD5U7/2x2l8g+0OwNmVc9ftjc+egr8PV1NHUY+eenTE46ldpKOrSOpgInp3ADqCQL259xVzzutlvnN2/wC77dIXsbi4LoSCpKkDiDQj8+ug/wB0/wBvuaPan7uXtB7dc62cdvzVtGzx29zGkiTKkqs5IWSMlHFCMqSPn0X32E+shetmv/hNCyjtz5SIT62646/ZR/VU3Nmw5/2Bcf7f3R/LpTbcW+zrb59t9K+ve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de6/9ShFHeN0kjdkkRldHRiro6kMroykMrKwuCOQfewSCCDQjqrokiNHIoaNgQQRUEHBBB4g+Y62YPg38rcX8huv6bau56+GLtrZmNp6PcNFLIqVG58VSrHT0m8McjBfPLJwtci38VSbnSk0fvMn2551tedtkXbdwlH9YLeMLKpPdKi0CzpXiw4SDyNSe1xT5m/vofdi5h+6v7py878nWTj2e3i9aWwmVSYtvuZyzzbTc6alIWILWjkAGLQqEy27ajsuk1IwIN4yfRIo9Lf7cXV/wCqnke7b1y46Etp1RngRwP+Y+oOelvtV732t5FHEZzFexga4mPenzxh0P4ZFqrDzBqBEqqulaqx9RWUlLVSUNTHV0M1TTwzS0VTE3pqaOSRC9NUKp4dCrC/19ha12UwX1ndiIa4ZkcY/hYH/J1kHzF7px7zylzJy3JfEWu6bXc2zgMQD4sLx5pxoWU/l1xpaqlWqyFRR0dLSy11RLV101NTwwS1tVMxBqKySONXqp2Qcu5Zjb6+93OyeNe3d2Yx4k0zuccdTE/5eq7H7qxbTyry9y7Hek2m17ZbWqAsSAIYkjxU4rpJ+09TY45qtwSbRj9UjfpA/p9PU3+0jk/7z7Fmy8uu7A6aRjiSMD/OfQDJ6xy91ve61tIZYxOZLxwdESnvf5/0EH4pGoqj1NASSfOn5X4v49dfVO0NrV8Mnbe9MbNR4KkhlRqnamHrFkgrN4ZFUJNPUBQyUCGxkqbOAY4XBp7kc623JGxnbdvkH9YLmMrGte6JGBDTtTgx4Rj+LIqENUn3Jvuw7796v3WTnjnO1c+0GyXiTX0pU+FuF1CweDabcsO+Bah7xs0i1I2mS5UrrQu7yO8kjtJJIzO7uxZ3diWZ3ZiWZmY3JPJPvDMkkkk56+m5ESNFjjUKigAACgAHAAeQHkOuPvXVuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3XutmD/AITSKT3X8m350r1bs1T/AEu+7Ksj/C9kPuj+XSm24v1uCe2+lfXvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvde9+691737r3Xvfuvdf/9Wg/wB+690odqbs3JsbcOK3XtDNZDb248JVx1uLy+Mnanq6SoiNwVYXWSJxdZI3DRyoSjqykgq7C/vdru4L/brp4byJqq6mjA/b/Ig4IwQR0HObuUOWOfeXN25Q5z2K23Llm/hMU9tcIJIpEPkVPAg0ZHUh0YB0ZWUEXZdB/wA1DbGUo6PAfIPD1G38yiR07762xj3r8BkQiooqM1t6DyZDF1EjXLmkSphJN1jiXj3kpyv74bddxpac32phuaUM8S6o3+ckWSp+aBh6Ko64f+/P91Vzry5fXPMX3b9/Tc9kDF02q/nMF7bVJOiz3AkRzIMBVuWhkoKNLM2erGdpdzdHdpGjp9jdmbF3JXVk8UNFi8duWjp81UzzAhKePCVcseVNQ9uIxAH4+nHuWNr3Pk3epF/dm8WUrvgKsoViTwHhsdVflTrnvz5yd95j2uspf68+2/NG32tuC7Sz7e88CIMMxu4FMOkVFWMpHDPXt19z9HdWmrp99dmbF21X0U8sVbishuWjqczTVEAtJTyYSkklyoqEH1jMBa5+nve6bnybssjfvLebKJ0wVaUM4PmPDU6q/KnVORuTfvM+6NlD/UX235o3G2uQHSWDb3gt3Q4VhdzqIQpoe4SgcaHquXvz+ajtrF0tXt/494SXP5V43gXfm6KCSgwWOZ1ZTPhNuTrHXZOeM2KtWLTQhh6o5l4MT80e+W32kb2fJ9qZrihAnlXTGnzji4sfm4UV4qw66Eewv91LzlzFe2/Mf3k+YE27Zi4d9psJzPeXIFDovdwBKQqchltmmkI+CaFqEUnbq3XuTfG4cruvd+ayG4dx5urkrcpl8nO1RWVdRIeWd2sscaKAscaBY4kARFVQAMa7+/vd0vJ7/cbp5ryVtTOxqST/AKsAYAwAB13D5R5Q5Y5C5b2jlDkzYrbbeWbCFYoLaBAkcaLwAA4k8WdiXdiXdmYklPe0fQj697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917rZr/AOE0IH+l35SNbkdcbAF/8G3NmiR9Pzp90fy6U23F+tvn230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r//1qD/AH7r3Xvfuvde9+690L3x+/vJ/p06c/uh5P7z/wCk7Y/8D8fkv/Ef7yY7wa/F6/Bq/wA5b/dd/Z7yx9V/WTYPoqfV/WQ6K8NXiLSvy9flXqJ/fj9w/wCsj7vf1p1f1d/q1uX1Gmmvwvo5tXh1x4n++v8AhmmmevfID+8n+nTuL++Hk/vP/pN3x/HPJ5L/AMR/vJkfPo8v7ng1f5u/+69Pv3M/1X9ZN/8ArafV/WTa6cNXiNWny9PlTr3sP+4P9ZL2i/qtq/q7/Vrbfp9VNfhfRw6fEpjxaf2v/DNVc9BD7IupY697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de62bf+E0P/M2vlL/AE/0ddff1/56XO2/w90fy6U23F+tvf230r697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r3v3Xuve/de697917r/9k===                

            """)
        # Decode the base64 string e crie um objeto de imagem
        image_data = base64.b64decode(canto_preto_base64)
        image = Image.open(io.BytesIO(image_data))

        # Redimensiona a imagem para 1200x700
        resized_image = image.resize((1200, 700), Image.Resampling.LANCZOS)

        # Converta a imagem redimensionada para um formato que o Tkinter possa usar
        self.imagem = ImageTk.PhotoImage(resized_image)

        # Crie um rótulo em tela cheia para exibir a imagem
        self.tela_cheia = tk.Label(self, image=self.imagem)
        self.tela_cheia.place(x=0, y=0, relwidth=1, relheight=1)

        # Variáveis
        self.directory_path = tk.StringVar()
        self.key_path = tk.StringVar()
        self.sensitive_files = []
        self.blacklist_directories = []
        self.scan_reports = []

        # Carregar configurações salvas, se houver
        self.load_settings()

        # Criar e exibir widgets
        self.create_widgets()

        # Definindo o modo de aparência inicial
        ctk.set_appearance_mode("light")

        # Agendar a execução do escaneamento a cada 5 minutos
        schedule.every(2).minutes.do(self.scan_blacklist_directories)

        # Iniciar o loop de agendamento
        self.after(100, self.start_schedule_loop)

    def create_widgets(self):

        interrogacao_base64 = """
iVBORw0KGgoAAAANSUhEUgAAABgAAAAYCAYAAADgdz34AAAABHNCSVQICAgIfAhkiAAAAAlwSFlzAAAAwgAAAMIBT4kc1wAAABl0RVh0U29mdHdhcmUAd3d3Lmlua3NjYXBlLm9yZ5vuPBoAAAGySURBVEiJndYxa1RBFAXgb1Zr46JVioAWVmKlhaDYBLa0Ezst0psU+QXaWkuqjY1VGhXxB5h6QbAQhCjYhcRVCKaIYSx2Ft5OZva9zcCBx5t7zpl35859E2KMaiOEsISnuI/ruJamvmMPn7AdY/xTFYkxngGWsYUjxBYcpdjlolZBfBX7HYRz7GN1rgHWcXoO8SlOsVE0SCsviZ/gAzZxG30M8ByHFZPBjEHKeSkt//ColNvEu4FvBd4BVpoGWxXxxzXxhskV/Cjwh9MKXapUy3Ym1MMdPMCFbO5Fgf83pdOzyoYdYy0JXMS7xtz7zOBmRWMddloq4w3eZu8+F9JU4u7AqMUgx0fZocLDSuwIxh2FT0xKNWTil/GlwhkvYvCqUEF38XUOZ7xIim5l4mvaT/2oZ9IVu4yf04fUZV+alO68sdczabldxq8QQgwhRPzGpQ6cXeoHLUc/29i2+MlBm9MqZlDY4DaDYZdmd16D2WbX0q6nuNqI7c+JO9uuG8SNFpM21H84DZNB+sRFxQ+aK68aJJMVDFMldKmW4TTnOULLtaWPJ7infG3ZxesY47im8R/7mK77QRMEdgAAAABJRU5ErkJggg==
"""
        # -----------------------------------------------------------------------------------------------------

        sair_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDowQzZGMjBENUUxN0MxMUVFQjczNUVDRjcyNUI2NDZFQyIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDowQzZGMjBENkUxN0MxMUVFQjczNUVDRjcyNUI2NDZFQyI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOjBDNkYyMEQzRTE3QzExRUVCNzM1RUNGNzI1QjY0NkVDIiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOjBDNkYyMEQ0RTE3QzExRUVCNzM1RUNGNzI1QjY0NkVDIi8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+h3ijewAAARdQTFRFAwMDycnJFhYWAAAAcXFxkJCQ/Pz8ZmZm5OTkAwMD29vbHh4eHx8flpaWjo6OkZGRi4uLOzs70tLSRkZGenp6ubm5bGxsMTExeXl5CgoKhYWF+vr6tLS0+/v78vLy2NjYIyMjpaWlDAwMwMDAtbW1np6e9fX1vr6+iIiIfHx8BAQEs7Oz/f39paWlkpKSAAAAPz8/zMzMioqKrKysHBwce3t7fn5+LCwsysrKGxsb6+vrODg4Pj4+dnZ2l5eX7e3tra2t8PDwCwsLExMTYmJiKioqMDAwEhISFBQUfX19x8fHYGBgjY2NT09PVlZW8/Pz2dnZCQkJt7e3AgIC9/f3ERER7+/vVVVV+fn5oqKiSUlJAAAA////StUKfAAAAF10Uk5T//////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////8A4a+dkAAAAOlJREFUeNpiiCENMIyqH7LqVXydOGFAigj1AibRcCBGhHrBaE4mGOAjQj1ntB8jDHAQoT46WoQk/0ZHs6DI6fhbhZCinkPXSJ0E9WzGrtzejCSod+CJlmQ3lCXaPY6q4sGB7BZsxKqPUYu0jlZi9yKgXlqGFQL07CTEo6NDXZTxq5cTjkYBkqakqXcWINI9/LwaBtHR3BJCxPo3RjQqmkkzgNjwkQ2y0XI34yI6vrTNoz08uYiPX11LBVsuEtKbPD+vECnpM8I+jI2a+YUzOpyk/EhqfvdxI608EWXnIam8Gi3/h596gAADAA890kW1rojaAAAAAElFTkSuQmCC
        """

        # --------------------------------------------------------------------------------------------------------

        save_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDpFRkNENjcyQ0UxODIxMUVFOTRBQkI2NTg5MUQyRjYwNSIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDpFRkNENjcyREUxODIxMUVFOTRBQkI2NTg5MUQyRjYwNSI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOkVGQ0Q2NzJBRTE4MjExRUU5NEFCQjY1ODkxRDJGNjA1IiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOkVGQ0Q2NzJCRTE4MjExRUU5NEFCQjY1ODkxRDJGNjA1Ii8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+P6ZW+gAAAL1QTFRF5OTkv7+/j4+PUVFRYGBgcXFxlZWVT09PBQUFmZmZt7e3Hh4e4uLi2NjY39/f4eHhvr6+vb29xsbGTk5OHBwcbGxsubm5k5OT1tbWcHBwq6urwsLCXFxc7+/vjY2NSUlJqKiolJSUkJCQpaWljo6Ou7u7sbGxgICAc3NzdHR0oaGh+/v7sbGxr6+vAAAAExMTJycn0NDQKCgokZGRAAAAFxcX9vb2dnZ2PDw8AgICd3d3BgYGMDAwAAAA////9ryNxQAAAD90Uk5T//////////////////////////////////////////////////////////////////////////////////8AjiZ8FwAAAOxJREFUeNrs1ckOgjAQgGFwRQsIggso7kQFjFvc7bz/Y7lGU9MpctPof2v5DgOZBAmSJf39p3jLJoTYHfZ5xyZdxK/opcHjPF/MyjA4Xy35/nj1hcdZpjSCwuWu+5YvU1q9eZrQ797yJqWTm1cEPnpe7DdriOJ8xS/eM+p12a/EeV6/5g+YD0xDfskwA3wfQt7Kh7hXYdTOM7VHoOJeB7eWY6q5oONe482j4b4ELSfL5LSgJPLDZoapORT5pPPo0O+lmHp90fuqMLbSTNZY9D093jwe7hsSpwbiCbpvfL9VFGSdp///xbf6kwADAACA3CQneTTlAAAAAElFTkSuQmCC        """

        # --------------------------------------------------------------------------------------------------------
        chave_base64 = """
iVBORw0KGgoAAAANSUhEUgAAAIcAAAAtCAYAAACEYHiCAAAAA3NCSVQICAjb4U/gAAAAAXNSR0IArs4c6QAAAARnQU1BAACxjwv8YQUAAAAJcEhZcwAAF50AABedATRi8NsAAAAZdEVYdFNvZnR3YXJlAHd3dy5pbmtzY2FwZS5vcmeb7jwaAAAE2UlEQVR4Xu2cS4iVZRjHx9QyyyxNs6zMLt7DdNEFhSiksEULF20EEVxEIApCS5cKuhMXgjshNBehiEFmRUqoZeWtvJUpapZ3M0szu/x+Z86Hx68zzpnhgJ3nff/w4/B5Zg7MfP95bu/z2ZaVlZWVldU89ai+RtQdcD88AmPhZXgC/oEdsB5OwI9wHq5BVnDdBg/ATFgJv4CGKPM3aI4l8BIMBL83K6h6wUhYDGehninK/AVfwNvwOESOpsnKv/pnYDlcgnpG6AgNchQWgmkoC/WsvkbQCJgDb8A98BOsgq/gOAyAu6CejBb9YBicgZ1g2skKoP4wF46BN/U3mAVDwPpjHCyAzlKNEeRreBGygmgifAR/gjd5PhgJCply/JoPwK6kbIpa/Ay/v6Mok9VC6gsz4HsobrBGKBeWphWjx2WoNUM9bG1fgaQVoXXzpk+GhypX7boXyuZw7nE3NNKN3A5Gnkg1WZcVwRzeSIddfSpX7TKSPAi2tv6MRpcJoIl6Q2fSFLa1SaeWCH8ZRo5XYXTlql0WoEYIC1Pffw6mwwtgBOlMphZb2y1gW5zVonoS1kK5btAYn8FG2F29Ln9NR1iUvgtJzzwipJXihpZlKpkEU+Dp6nWj8vMsSjVJsopiDmcb9QzSXfl5dj9Gm2QVwRxX4TTYojZLDsOOwO+Vq0QVwRxOPT8Fj96bJY12ETRJsopgDm/kQbC7aMZOhunJFnYMOJZvZC4SUlGGPLabzjvGg4du3b2hRgojka8WsxfA2uMKJKco5rDuMA0MBre9agdiXdE5WA0HwLmJY3g/9xA0s6ZpCUUaD9t6HgZTggs/RpKuyu5kEzg38XfjVNUdEc1nBEnKIJHMYa3gLsa3YC3lAMuaoStyejoUXB/8EBy1O13VJM48voE/IAlFMocqhle7wFbUNOEZS+3xfSHrCvc/ToH7o0pT3QfD4QfwiN/DuufBNGPkcN8jiUWgaOYoZHqwbvgcPoYN1WvbXY3zDiytvu4HO5OyQdwK2wvvgSe+RhANYnH6JWQFkTfcGuROsFj1tLboaLyeDa4VGnkKjA6eybwGrgC8D0abk/AmhFfUyFGWN9sb6xxEatOC1/vA940gxc6HuGL4FFjHrABb5VFgkWrL+x00Y7byv1Qq5uhMdiOmFyOMnU7tUpAF6sNgMWoXY6tserEuMYq4vBzygC6b47osNo0ExZNyzjc8W/FaQxhFfFLOItUi1xmIHZEdkt2NBgulbI4b9Ss4irfW+AQ0i2nEAvUx0CB2K5vBgduzoEF+hnARJJvjvzJiaArTjLXKVBgEFrGPgmlmO2wDH31wu8y045DMCBKmzc3muLmMCq+D5rAG0SD+m62tK4S2tBrDcxhrkD1gHRJC2Rw3lzWHgzJnIz4XY/pwq8wZiKZwVcCn4yxQNYjvmXIcxGUlIP+AnJHYwbwFjs9NN9YXDtemwRqwIHVAtgiyEpPDsnlQmKMezjzWQQjZ12c1Jm++6ULc85Byd+Lv0zY3hHLN0bjsQjyks5vZChafFqbFmYyyaPWwb1nlKitZOfewSC2nFk0TQjmtNFdGF9NNMXpvaYX4IW6RfMzS/3fMRy99NMK2V3M4YXXVsOWVzdF9Wa95lO8agAvOFqemFU93k9kWy8rKysrKakhtbf8CIVlED9WLREIAAAAASUVORK5CYII=
            
            """
        # --------------------------------------------------------------------------------------------------------
        scan_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAABIAAAASCAYAAABWzo5XAAAABHNCSVQICAgIfAhkiAAAAAFzUkdCAK7OHOkAAAAEZ0FNQQAAsY8L/GEFAAAACXBIWXMAAACvAAAArwErsBW5AAAAGXRFWHRTb2Z0d2FyZQB3d3cuaW5rc2NhcGUub3Jnm+48GgAAAWdJREFUOE+t078rRlEcx/FDhEhJycKEGFgM2Ez+ASmLwmYik8VE2Swmg8WPMhhlEIvJRMmPTFhQlPIrP/Lj/Tnnnuu43ed5eur51Ot2znOfc84993uuKXTGcYOfwBzCNCK8/4QV1MMU6UKusYBN23O5w61r2pSi1TVtWrCGYazrB0UrtLtmnBI0owe6V45kNG5MjWJdUqInHcAW9rCDeVQhNZkmqsA0TtGARWgLfcgavVj70qJU4w2dtmdMLb4wYXt/0bhu10xPJc6wgTpM4RX9yCva8igu8YF7LEFbTo0vf1pU7jZoy484xDuyRodNA320QBdUrRccYxBlCKNx/yqZPEd6yRfYxgxW8YxJhJWOz5FPcqLdiKqlp9PTzuIcKoRPzgN5hWU8QH/+hF72Pr6RMfoA8y1tE7TAiO1F0VesA6jVvVxfvyY5QQfi8qvEvaixPReVW1vxUXWGXNNGCx/gCJq4EDHmF655ULR9leqMAAAAAElFTkSuQmCC
        
        """
        # Decodifique as strings base64 em dados binários
        image_data_chave = base64.b64decode(chave_base64)
        image_data_scan = base64.b64decode(scan_base64)
        image_data_sair = base64.b64decode(sair_base64)
        image_data_interrogacao = base64.b64decode(interrogacao_base64)
        image_data_save = base64.b64decode(save_base64)

        # Crie as imagens a partir dos dados decodificados
        imagem_botao_chave = PhotoImage(data=image_data_chave)
        imagem_botao_sair = PhotoImage(data=image_data_sair)
        imagem_botao_interrogacao = PhotoImage(data=image_data_interrogacao)
        imagem_botao_save = PhotoImage(data=image_data_save)

        #Crie os botões com as imagens como conteúdo
        botao_chave = ctk.CTkButton(master=self, image=imagem_botao_chave, text= "", text_color="",
                                            fg_color="transparent", hover_color="#f7f3f2", font=("Times New Roman", 17),
        command=self.choose_key_file).grid(row=3, column=1, padx=10, pady=10, sticky="s")

        botao_sair = ctk.CTkButton(master=self, image=imagem_botao_sair, text="", text_color="", fg_color="transparent",
                                   hover_color="#f7f3f2", font=("Times New Roman", 17), command=self.close_program)
        botao_sair.grid(row=3, column=1, padx=10, pady=10, sticky="e")

        botao_interrogacao = ctk.CTkButton(master=self, image=imagem_botao_interrogacao, text="", text_color="",
                                           fg_color="transparent", hover_color="#f7f3f2", font=("Times New Roman", 17),
        command=self.tutorial)
        botao_interrogacao.grid(row=0, column=1, padx=10, pady=10, sticky="e")

        botao_save = ctk.CTkButton(master=self, image=imagem_botao_save, text="", text_color="", fg_color="transparent",
                                   hover_color="#f7f3f2", font=("Times New Roman", 17), command=self.save_settings)
        botao_save.grid(row=3, column=1, padx=10, pady=10, sticky="w")

        # --------------------------------------------------------------------------------------------------------
        # String base64 da imagem do logo do grupo
        base64_image = """iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDpEN0YxMzFBQkUxODUxMUVFQTBERjhGRUJDMjYxOERBNSIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDpEN0YxMzFBQ0UxODUxMUVFQTBERjhGRUJDMjYxOERBNSI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOkQ3RjEzMUE5RTE4NTExRUVBMERGOEZFQkMyNjE4REE1IiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOkQ3RjEzMUFBRTE4NTExRUVBMERGOEZFQkMyNjE4REE1Ii8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+S3PHGwAAAwBQTFRFAPb1AMPFaf7+U/79APf1AN3dh/38xf79ALu+sdvdg/Lylf/+AMrKANbXMf/+APTxAPn2ANDRAL/DU9na8P//AOPh1f7+AMjJAMLEAOvpANraJP/+4P/+9f//APX0VcTG7P//3v7+xuLjAMHDAMjLAMzOAPv4AMDCAP//ANXWAP//ANbWAPTxAPr5AMXGzf/+Xf/+qv//APb2AOXjAOLhAMfIuP7+AO/tn/n5AODfAM/Q7//+APHtANLSAPbz7/f3AO7sAMvNAM/QAPTy6P//ANbWAMzNAPz5of79pNnb+P//AO3rAO/sQ8LEAOrny/r6AMvMAOfmAOjmjf39AOvqAPn35f7+APr5vv//AOblAMXHmv//ef//Qc3OAOjnAPXz8f//ANHRAMjJAMXHANvbAM7Ph///AN/eAMTGAPz7APn5ANfXAOfjAMHDANbW/v///P//////+///x///0f7+AP//AO7sAPj2AM7PAP79ANXUANDRsfn5pO/6svP7AP38/v/+Rd73AN7eyvHxAPz7APLwf//+AP//q/P7sf39s/79ALu+ANHSmv38APfzAPf0AOTjAObkAPXyAPz6APDuAPr6AMPGcPz83+/xAMTH2O3uSMDEpdXXAPPxq/j8wN/g+P39tvz7+f39lerqANLSANXWAPn4fuvsANPULP372f/+yujoy+fpAMnKVvz7z/PzAM3NAODgdv37f/37ANTUYsXIasrLAPr4APr6APz6AOTiAMfJAPz7AOHg/Pz8Af387f39AOHe/v7+/v39RMPGAPz6AODgANrbANzdAN3bAPLwAPLyAN3anv39TPz79vv7gc7PAP39AMLEAMbIAN3dAPv3AN7cAPn3AOfnU8LEAPn5k+j6+vv8gf38NMHC3e3t6e/wD8DCAO/sdcjKAMnKAMnKAPPxANrb8vLzAPn1APf2ANjZ/f7/7Pj4sff8rv79r///APj2sv39sv7+AL7CAPr3AL7AAPj2AOXki//+AMzNAO3rAPj1AO7slPv8////rBwyFgAAAQB0Uk5T////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////AFP3ByUAAAJgSURBVHjaYvhPGmAYVU+WesVV+SSoL0ow85CRJVq9iOKno0eZ2cOIU684Xe0rg7NOPMvXPWZeBNVb9+5geBQyUVR0QtBfnVvMJQTUi7B9N4+zzF7i5NSz9Erkb7s9LgV41Z/9/uiHVLXjpYykjKdOS4+F9K90waO+QMHd3DjPcVr3Pd4rsd1fVBKbfr/qlMWtXoRNZ2JeQO69ObO45JqVHn9JTjQx36GPW73Yn98pAdMect7J524XKbzveykpO0jHALd6dp3IpTO7lRb8L+ZrTxMpeLE1d8le83ovXOplNz06VpF70/5/BFPbT74Y60JP9adSkuGK+djVF0Wo/T3BkyzsX8z8811HsDR3yWJVvby4WzHWDdjUF4btf/n3eZLjjfB2vuLbjfMM2i+efp20JfB75cJFWM33+nz978mkZOEaQ77outp/7QYflgPND/ouHXYIh/unPzpYkSv/Nt9/809uZosDb8DuV9PPxxU+Zs+Mq5K7Vx/+z97eHnO+0H6rUOIx86MJONWXGEmmBOQemb8PGP4KhQ+2XUrKC3zGl48z/BWmn9q1NEDIl3MGl9z6NRe+JVdb/r3Fjif9uGx/ZLLUMbd7NiT9VF/53fdeFl/65Gb5bZK3IYnx8mVG019LrvnZ7QnDm54V7rI8yjmTV13xq2LFc1FJ5z1iRfjzl0Lw1Ul//UTT00Ujf5+6vp+dcP5ddW6ZQPwUG50/Au7vZYnI70Uiivoe4dfXdZUoyOYTV/4UJHiwKcqSUr55yQ5UeZufT0gNSEF+fn4BiCTW/IL8AmDE5ReM1l8Dqx4gwACHbFzwX2V6UgAAAABJRU5ErkJggg==
        """
        # Decodificar a string de base64 em uma imagem
        image_data = base64.b64decode(base64_image)
        image = Image.open(io.BytesIO(image_data))
        image = ImageTk.PhotoImage(image)

        # Exibir a Logo
        #image_label = ctk.CTkLabel(master=self, image=image, text="", height=0)
        #image_label.grid(row=3, column=1, padx=0, pady=0, sticky="s")
        # --------------------------------------------------------------------------------------------------------

        #Criação de botões na lateral
        frame = ctk.CTkScrollableFrame(master=self, fg_color="#484444", border_color="#962CCA", border_width=2,
                                       height=530)
        frame.grid(row=0, column=0, rowspan=3, padx=10, pady=10)
        frame.grid_columnconfigure(0, weight=1)  # Configura a coluna 0 para expandir horizontalmente

        botao_escanear = ctk.CTkButton(master=self, text="Escanear", text_color="black",
                                       fg_color="#9370DB", width=820,
                                       hover_color="#53DEC9", font=("Times New Roman", 17), command=self.start_scan)
        botao_escanear.grid(row=0, column=1, padx=10, pady=10, sticky="w")

        ctk.CTkButton(master=frame, text="Escolher Diretório", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.choose_directory).grid(row=3, column=0, padx=0, pady=10,
                                                                                 sticky="ew")

        ctk.CTkButton(master=frame, text="Excluir Arquivos", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.open_delete_window).grid(row=5, column=0, padx=0, pady=10,
                                                                             sticky="ew")
        ctk.CTkButton(master=frame, text="Mover Arquivos", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.move_files).grid(row=6, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Adicionar Blacklist", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.choose_blacklist_directory).grid(row=7, column=0, padx=0,
                                                                                           pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Lista Blacklist", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.show_blacklist).grid(row=8, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Relatório Excel", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.open_report).grid(row=9, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Relatório em Gráfico", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.graphic).grid(row=10, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Outlook", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.informacoes_outlook).grid(row=11, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Sobre", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.sobre).grid(row=12, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Filtrar Info", text_color="black", fg_color="#53DEC9", hover_color="#9370DB",
                      font=("Times New Roman", 17),
                      command=self.filtrado).grid(row=1, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Procurar Info", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.Escanear_info_especifica).grid(row=2, column=0, padx=0, pady=10, sticky="ew")

        # Quadrado Vazio
        quadrado_vazio = ctk.CTkFrame(master=self, width=900, border_color="#962CCA", border_width=2)
        quadrado_vazio.grid(row=2, column=1, padx=10, pady=(0, 0))
        quadrado_vazio.grid_rowconfigure(0, weight=1)
        quadrado_vazio.grid_columnconfigure(0, weight=1)

        self.output_text = ctk.CTkTextbox(master=quadrado_vazio, wrap=tk.WORD, border_color="#962CCA", border_width=1,
                                          height=540, width=900)
        self.output_text.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")
        messagebox.showinfo("DICA",
                            "Clique no botão de interrogação no topo do aplicativo para se ter mais informações!")

    def sobre(self):
        messagebox.showinfo("Créditos",
                            "Aplicação desenvolvida por Grupo Mountain. 2024 - v.6.0> https://www.linkedin.com/in/lucas-portugal-/ \nhttps://www.linkedin.com/in/danilo-rocha-mendes/")

    def Escanear_info_especifica(self):
        # Cria uma janela modal para entrada de texto
        root = tk.Tk()
        root.withdraw()  # Esconde a janela principal

        # Pergunta ao usuário que tipo de informação ele deseja procurar
        info_to_search = simpledialog.askstring("Busca de Informação",
                                                "Que tipo de informação deseja procurar?")

        messagebox.showinfo("Scan em progresso",
                            "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar. Caso tenha recusado a escolha, ignore esta mensagem)")
        # Se o usuário cancelar a entrada, info_to_search será None
        if info_to_search is not None:
            directory_path = self.directory_path.get()
            if directory_path:
                # Limpa o texto antigo
                self.output_text.delete(1.0, tk.END)

                results = {}
                process_directory(directory_path, results)

                sensitive_files = []

                # Variável para armazenar os diretórios onde a informação foi encontrada
                directories_with_info = []

                for path, data in results.items():
                    results[path] = list(set(data))
                    sensitive_files.append(path)

                    # Chamada para a função de detecção de informações sensíveis
                    results = extract_sensitive_info_from_image(path, results)

                    # Verifica se rostos foram detectados e exibe uma mensagem
                    if 'Rosto' in data:
                        self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

                    # Verifica se algum texto contém a informação especificada
                    for info in data:
                        if info_to_search in info[1]:  # info[1] contém o texto detectado
                            directories_with_info.append(path)
                            break

                # Mostra os diretórios onde a informação foi encontrada
                if directories_with_info:
                    self.output_text.insert(tk.END, f"Diretórios com {info_to_search} encontrada:\n")
                    for directory in directories_with_info:
                        self.output_text.insert(tk.END, f"{directory}\n")
                else:
                    self.output_text.insert(tk.END, f"Nenhum diretório com {info_to_search} encontrada.\n")

                self.sensitive_files = sensitive_files

                # Aumenta o tamanho da fonte
                self.output_text.configure(
                    font=(
                        "Times New Roman",
                        20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

                # Adiciona os dados do escaneamento a self.scan_reports
                current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
                for path, data in results.items():
                    self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

                # Gera o relatório
                self.generate_report()
        else:
            print("Busca cancelada pelo usuário.")

    def show_blacklist(self):
        blacklist_window = tk.Toplevel(self)
        blacklist_window.title("Lista Blacklist")
        blacklist_window.geometry("800x600")  # Definindo a geometria da janela

        listbox_frame = tk.Frame(blacklist_window)
        listbox_frame.pack(fill=tk.BOTH, expand=True)  # Faz o frame expandir para preencher toda a janela

        listbox = tk.Listbox(listbox_frame, selectmode=tk.MULTIPLE)
        listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)  # Faz a Listbox preencher todo o frame

        for directory in self.blacklist_directories:
            listbox.insert(tk.END, directory)

        scrollbar = tk.Scrollbar(listbox_frame, orient="vertical", command=listbox.yview)
        scrollbar.pack(side="right", fill="y")
        listbox.config(yscrollcommand=scrollbar.set)

        def remove_selected():
            selected_indices = listbox.curselection()
            selected_directories = [listbox.get(index) for index in selected_indices]
            for directory in selected_directories:
                self.blacklist_directories.remove(directory)
            messagebox.showinfo("Removido", "Os diretórios selecionados foram removidos da blacklist.")
            self.save_settings()

        def clear_blacklist():
            self.blacklist_directories = []
            messagebox.showinfo("Blacklist Esvaziada", "A lista de blacklist foi esvaziada com sucesso.")
            blacklist_window.destroy()  # Fechar a janela após esvaziar a blacklist

        remove_button = tk.Button(blacklist_window, text="Remover Selecionados", command=remove_selected)
        remove_button.pack(padx=10, pady=10)
        esvaziar_button = tk.Button(blacklist_window, text="Esvaziar Blacklist", command=clear_blacklist)
        esvaziar_button.pack(padx=10, pady=10)

    # save
    def load_settings(self):
        if os.path.exists("settings.json") and os.path.getsize("settings.json") > 0:
            with open("settings.json", "r") as f:
                settings = json.load(f)
                self.key_path.set(settings.get("key_path", ""))
                if self.key_path.get():
                    os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = self.key_path.get()
                self.blacklist_directories = settings.get("blacklist_directories", [])
        else:
            messagebox.showwarning("Aviso", "O arquivo de configuração está vazio ou não existe.")

    def save_settings(self):
        settings = {
            "key_path": self.key_path.get(),
            "blacklist_directories": self.blacklist_directories
        }

        with open("settings.json", "w") as f:
            json.dump(settings, f)

        messagebox.showinfo("Salvo", "As configurações foram salvas com sucesso.")

    def tutorial(self):
        popup = tk.Toplevel(self)
        popup.title("Tutorial")
        popup.geometry("900x600")

        scrollbar = tk.Scrollbar(popup)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        tutorial_text = tk.Text(popup, wrap=tk.WORD, yscrollcommand=scrollbar.set, font=("Arial", 12))
        tutorial_text.pack(fill=tk.BOTH, expand=True)

        scrollbar.config(command=tutorial_text.yview)

        # Adicione o conteúdo do tutorial aqui
        tutorial_content = """
COMO PEGAR A CHAVE JSON (I.A):\n

01- Crie uma conta no Google Cloud (cloud.google.com).\n
02- Adicione a forma que quer que seja feito os pagamentos.\n
03- Crie um novo Projeto.\n
04- Após criado, clique nas 3 barrinhas no canto superior esquerdo. Escolha "APIs e serviços" e em seguida "Biblioteca".\n
05- Pesquise por "Cloud Vision API".\n
06- Selecione a que se parece com um olho azul.\n
07- Clique em ativar e recarregue a pagina e verifique se a opção "ativar" mudou para "gerenciar".\n
08- Clique nas 3 barrinhas no canto superior esquerdo. Escolha "APIs e serviços" e em seguida "Credenciais".\n
09- Clique em "Criar Credencial" e escolha "Contas de Serviço".\n
10- Escolha um nome para a conta de serviço e crie.\n
11- Após isso clique na conta de serviço e vá em "chaves".\n
12- Clique em "adicionar chave" e crie uma nova chave JSON.\n
13- Coloque a chave em algum diretório que você irá se lembrar para quando for usar o aplicativo.\n

        FUNCIONAMENTO DOS BOTÕES:\n


01- ❓ > Abre o tutorial (onde você está agora).\n
02- Escanear > Escaneia o diretório escolhido.\n
03- Filtrar Info > Permite o usuário escolher que tipo de informação ele deseja ver dentre as opções: CPF, Gênero, Religião, RG, Telefone e Rosto (necessário escolher o diretório antes).\n
04- Procurar Info > Permite o usuário procurar por uma informação especifica dentro do diretório escolhido (necessário escolher o diretório antes).\n
05- Escolher Diretório > Escolhe o diretório a receber o scan.\n
06- Escolher Chave > Seleciona a chave JSON da I.A dentro do diretório que foi salvo. (necessário para Scan de imagens)\n
07- Excluir Arquivos > Exclui todos arquivos sensiveis encontrados no diretório escolhido com exceção dos não sensíveis (necessário ter feito o scan antes).\n
08- Mover Arquivos > Move todos arquivos sensiveis encontrados para um diretório de sua escolha (necessário ter feito o scan antes).\n
09- Adicionar Blacklist > Escolhe 1 ou mais diretórios para ser constantemente monitorado em busca de arquivos sensiveis, caso algum arquivo sensível apareça será aberto um pop-up falando que foi encontrado um arquivo sensível e te mostrará o diretório que ele se encontra a cada 5 minutos.\n
10- Lista Blacklist > Mostra a lista de todos diretórios salvos na blacklist, caso queira remover um ou mais basta clicar nos desejados e então clicar em "Remover Selecionados", caso queira esvaziar toda blacklist, clique em "Esvaziar blacklist".\n
11- Relatório Excel> Mostra um relatório em uma planilha excel contendo informações de escaneamentos anteriores, dentro do relatório contem a data que foi feito o escaneamento, o horário, que tipo de arquivo sensivel foi encontrado e o diretório onde estava.\n
12- Relatório em Gráfico > Mostra um relatório no formato de pizza contendo uma porcentagem em cada fatia para mostrar quantos % de cada informação foi encontrada.\n
13- Outlook > Pedirá seu login e senha para acessar o seu Outlook, após dar as informações o app vai ler e baixar todos e-mails para a pasta "Anexos_email" que vai estar dentro do local onde o executável original se encontra. Dentro da pasta "Anexos_email" terá várias sub-pastas com os nomes referentes aos titulos dos e-mails respectivos que foram baixados. Com essa pasta o usuário pode fazer o que quiser, desde scans até mover ou excluir arquivos. Há também uma opção para excluir a pasta "Anexos_email" que vai deletar por completo todos e-mails baixados.\n
                        LEMBRE-SE DE SEMPRE SALVAR AS CONFIGURAÇÕES!!! 
        """

        tutorial_text.tag_configure("bold", font=("Arial", 12, "bold"))
        tutorial_text.insert(tk.END, tutorial_content)

    def choose_key_file(self):
        key_file = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if key_file:
            self.key_path.set(key_file)
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = key_file
            self.save_settings()  # Salvar as configurações sempre que o caminho da chave for definido

    def choose_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.directory_path.set(directory)
            messagebox.showinfo("Concluído", "Diretório escolhido com sucesso.")

    def generate_report(self):
        # Abre ou cria um arquivo para o relatório
        report_file_path = "scan_report.csv"
        with open(report_file_path, mode="w", newline="") as report_file:
            report_writer = csv.writer(report_file)

            # Escreve o cabeçalho do relatório
            report_writer.writerow(["Data", "Horário", "Informação encontrada", "Diretório"])

            # Escreve os dados de cada escaneamento
            for scan_data in self.scan_reports:
                data_scan, _, diretorio, data, _ = scan_data
                for info_type, info in data:
                    data = data_scan.split()[0]  # Apenas a data
                    hora = data_scan.split()[1]  # Apenas a hora
                    report_writer.writerow([data, hora, info_type, diretorio])

        # Obtém o diretório atual do arquivo Python em execução
        current_directory = os.path.dirname(os.path.realpath(__file__))
        excel_file_path = os.path.join(current_directory, "resultado_scan.xlsx")

        # Após gerar o relatório CSV, cria o relatório Excel
        self.create_excel_report(report_file_path, excel_file_path)
        # Exclui o arquivo CSV após a criação do arquivo Excel
        os.remove(report_file_path)

    def open_report(self):
        if os.path.exists("resultado_scan.xlsx"):
            os.system("start resultado_scan.xlsx")
        else:
            messagebox.showwarning("Aviso", "O relatório ainda não foi gerado.")

    def create_excel_report(self, csv_file, excel_file):
        # Abre o arquivo CSV e lê os dados
        with open(csv_file, 'r') as f:
            csv_reader = csv.reader(f)
            data = list(csv_reader)

        # Cria um novo arquivo Excel
        wb = Workbook()
        ws = wb.active

        # Copia os dados do CSV para o Excel
        for row_index, row in enumerate(data):
            for col_index, value in enumerate(row):
                ws.cell(row=row_index + 1, column=col_index + 1, value=value)

        # Salva o arquivo Excel
        wb.save(excel_file)

    def graphic(self):
        messagebox.showinfo("DICA",
                            "O gráfico que lista o caminho dos diretórios é separado em cores, onde: Azul - TOP 1, Amarelo - TOP 2, Verde - TOP 3!")
        if not os.path.exists("resultado_scan.xlsx"):
            messagebox.showwarning("Aviso", "O arquivo 'resultado_scan.xlsx' não existe.")
            return

        # Carregando os dados do arquivo Excel
        data = pd.read_excel("resultado_scan.xlsx")

        if data.empty:
            messagebox.showwarning("Aviso", "O arquivo 'resultado_scan.xlsx' está vazio.")
            return

        # Converte as colunas 'Data' e 'Horário' para datetime e combina-as
        data['Data'] = pd.to_datetime(data['Data'].astype(str) + ' ' + data['Horário'].astype(str),
                                      format='%d-%m-%Y %H:%M:%S', dayfirst=True)

        # Ordena os dados pela coluna 'Data'
        data.sort_values(by='Data', ascending=False, inplace=True)

        # Seleciona as linhas com a data mais recente
        latest_time = data['Data'].max()
        latest_data = data[data['Data'] == latest_time]

        # Atualiza o campo 'Diretório' para refletir o caminho da pasta, não do arquivo
        latest_data['Diretório'] = latest_data['Diretório'].apply(lambda x: str(Path(x).parent))

        # Contagem de informações por diretório
        directory_info_counts = latest_data.groupby('Diretório')['Informação encontrada'].count().sort_values(
            ascending=False)
        top_directories = directory_info_counts.head(3)
        total_info = directory_info_counts.sum()
        top_directory_percents = (top_directories / total_info) * 100
        other_percent = 100 - top_directory_percents.sum()

        # Tamanho da figura
        desired_width_px = 1200
        desired_height_px = 700
        dpi = plt.rcParams.get('figure.dpi')
        figsize_inches = (desired_width_px / dpi, desired_height_px / dpi)
        fig = plt.figure(figsize=figsize_inches, dpi=dpi)

        # Função para quebrar texto nas legendas
        def wrap_labels(labels, width=30):
            return ['\n'.join(textwrap.wrap(label, width=width)) for label in labels]

        # Gráfico geral de informações sensíveis
        total_info_latest = latest_data['Informação encontrada'].value_counts()
        wrapped_labels_total_info = wrap_labels(total_info_latest.index)
        ax1 = fig.add_subplot(2, 2, 2)
        ax1.pie(total_info_latest, labels=wrapped_labels_total_info, autopct='%1.1f%%', startangle=90)
        ax1.set_title(
            f'Total de Informações Sensíveis Encontradas - Todos diretórios (Volume: {total_info_latest.sum()} )',
            fontsize=10, fontweight='bold')

        # Preparação dos dados e rótulos para o gráfico de pizza dos TOP diretórios
        pie_labels = [f'TOP {i + 1} - {dir}' for i, dir in enumerate(top_directories.index)] + ['Outros']
        pie_data = list(top_directory_percents) + [other_percent]
        wrapped_labels_pie_data = wrap_labels(pie_labels)

        ax_total_comparison = fig.add_subplot(2, 2, 1)
        ax_total_comparison.pie(pie_data, labels=wrapped_labels_pie_data, autopct='%1.1f%%', startangle=90,
                                labeldistance=1.3)
        ax_total_comparison.set_title('Comparação dos TOP Diretórios com os demais', fontsize=10, fontweight='bold')

        # Gráficos detalhados para cada TOP diretório
        for i, directory in enumerate(top_directories.index):
            specific_data = latest_data[latest_data['Diretório'] == directory]
            info_counts = specific_data['Informação encontrada'].value_counts()
            wrapped_labels_info = wrap_labels(info_counts.index)
            ax = fig.add_subplot(2, 3, i + 4)
            ax.pie(info_counts, labels=wrapped_labels_info, autopct='%1.1f%%', startangle=90)
            ax.set_title(f'TOP {i + 1} (Volume: {info_counts.sum()})', fontsize=10, fontweight='bold')

        plt.tight_layout()
        plt.show()

    def start_scan(self):
        directory_path = self.directory_path.get()
        if directory_path:
            # Limpa o texto antigo
            self.output_text.delete(1.0, tk.END)

            messagebox.showinfo("Scan em progresso",
                                "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar")

            results = {}
            process_directory(directory_path, results)

            sensitive_files = []

            for path, data in results.items():
                results[path] = list(set(data))
                sensitive_files.append(path)

                # Chamada para a função de detecção de informações sensíveis
                results = extract_sensitive_info_from_image(path, results)

                # Verifica se rostos foram detectados e exibe uma mensagem
                if 'Rosto' in data:
                    self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

            for path, data in results.items():
                path = os.path.normpath(path)
                self.output_text.insert(tk.END, f"Informações sensíveis encontradas em: {path}\n")

                types_found = set(info[0] for info in data)
                for info_type in types_found:
                    if info_type == 'Rosto':
                        self.output_text.insert(tk.END, f"{info_type} encontrado\n")
                    else:
                        self.output_text.insert(tk.END, f"{info_type} encontrado\n")

                self.output_text.insert(tk.END, "\n")

            self.sensitive_files = sensitive_files

            # Aumenta o tamanho da fonte
            self.output_text.configure(
                font=(
                "Times New Roman", 20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

            # Adiciona os dados do escaneamento a self.scan_reports
            current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
            for path, data in results.items():
                self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

            # Gera o relatório
            self.generate_report()

    def filtrado(self):
        # Cria uma janela modal para entrada de texto
        root = tk.Tk()
        root.withdraw()  # Esconde a janela principal

        tipos_arquivo = simpledialog.askstring("Tipos de arquivo",
                                               "Digite os tipos de arquivo que deseja filtrar (separados por vírgula) - Opções: CPF, Gênero, Religião, RG, Telefone, Rosto: ")
        messagebox.showinfo("Scan em progresso",
                            "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar. Caso tenha recusado a escolha, ignore esta mensagem)")  # Se o usuário cancelar a entrada, tipos_arquivo será None
        if tipos_arquivo is not None:
            # Remove espaços em branco extras e divide os tipos de arquivo
            tipos_arquivo = [tipo.strip().lower() for tipo in tipos_arquivo.split(',')]

            # Inicia o scan apenas se pelo menos um tipo de arquivo for especificado
            if tipos_arquivo:
                directory_path = self.directory_path.get()
                if directory_path:
                    # Limpa o texto antigo
                    self.output_text.delete(1.0, tk.END)

                    results = {}
                    process_directory(directory_path, results)

                    sensitive_files = []

                    for path, data in results.items():
                        results[path] = list(set(data))
                        sensitive_files.append(path)

                        # Chamada para a função de detecção de informações sensíveis
                        results = extract_sensitive_info_from_image(path, results)

                        # Verifica se rostos foram detectados e exibe uma mensagem
                        if 'rosto' in [info[0].lower() for info in data]:
                            self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

                    # Filtra os resultados para mostrar apenas os diretórios com informações sensíveis especificadas
                    filtered_directories = []
                    for path, data in results.items():
                        tipos_encontrados = [info[0].lower() for info in data if info[0].lower() in tipos_arquivo]
                        if tipos_encontrados:
                            filtered_directories.append(path)

                    # Mostra os diretórios filtrados na saída
                    if filtered_directories:
                        self.output_text.insert(tk.END, "Diretórios com informações sensíveis filtradas:\n")
                        for directory in filtered_directories:
                            self.output_text.insert(tk.END, f"{directory}\n")
                    else:
                        self.output_text.insert(tk.END,
                                                "Nenhum diretório com informações sensíveis filtradas encontrado.\n")

                    self.sensitive_files = sensitive_files

                    # Aumenta o tamanho da fonte
                    self.output_text.configure(
                        font=(
                            "Times New Roman",
                            20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

                    # Adiciona os dados do escaneamento a self.scan_reports
                    current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
                    for path, data in results.items():
                        self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

                    # Gera o relatório
                    self.generate_report()
            else:
                print("Nenhum tipo de arquivo especificado para filtrar.")
        else:
            print("Entrada cancelada pelo usuário.")

    def open_delete_window(self):
        self.delete_window = Toplevel(self.master)
        self.delete_window.title("Excluir Arquivos")

        tk.Button(self.delete_window, text="Excluir todos arquivos", command=self.delete_files, width=20,
                  height=2).pack()

        tk.Button(self.delete_window, text="Excluir individualmente", command=self.open_individual_delete_window,
                  width=20, height=2).pack()

    def delete_files(self):
        directory_path = self.directory_path.get()
        if directory_path:
            confirmation = messagebox.askyesno("Confirmação",
                                               "Tem certeza de que deseja excluir todos os arquivos sensíveis no diretório?")
            if confirmation:
                for sensitive_file in self.sensitive_files:
                    if os.path.isfile(sensitive_file):
                        os.remove(sensitive_file)

                messagebox.showinfo("Concluído", "Todos os arquivos sensíveis foram excluídos com sucesso.")
                self.delete_window.destroy()

    def open_individual_delete_window(self):
        self.individual_delete_window = Toplevel(self.delete_window)
        self.individual_delete_window.title("Excluir Arquivos Individualmente")

        scrollbar = Scrollbar(self.individual_delete_window)
        scrollbar.pack(side="right", fill="y")

        self.listbox = Listbox(self.individual_delete_window, selectmode="multiple", yscrollcommand=scrollbar.set)
        for file in self.sensitive_files:  # self.sensitive_files deve ser a lista dos arquivos
            self.listbox.insert("end", file)
        self.listbox.pack(side="left", fill="both", expand=True)

        scrollbar.config(command=self.listbox.yview)

        Button(self.individual_delete_window, text="Remover Selecionados", command=self.remove_selected_files).pack()

    def remove_selected_files(self):
        selected_indices = self.listbox.curselection()
        selected_files = [self.listbox.get(i) for i in selected_indices]
        for file in selected_files:
            if os.path.isfile(file):
                os.remove(file)
        messagebox.showinfo("Concluído", "Os arquivos selecionados foram excluídos com sucesso.")
        self.individual_delete_window.destroy()

    def move_files(self):
        if not self.sensitive_files:
            messagebox.showwarning("Aviso", "Nenhum arquivo sensível foi encontrado.")
            return

        destination_directory = filedialog.askdirectory()
        if destination_directory:
            for sensitive_file in self.sensitive_files:
                try:
                    shutil.move(sensitive_file, destination_directory)
                except Exception as e:
                    self.output_text.insert(tk.END, f"Erro ao mover arquivo '{sensitive_file}': {str(e)}\n")

        messagebox.showinfo("Transferência concluída!", "Todos os arquivos foram transferidos com sucesso!")

    def close_program(self):
        self.destroy()  # Fecha a janela principal do aplicativo

    def choose_blacklist_directory(self):
        blacklist_directory = filedialog.askdirectory()
        if blacklist_directory:
            # Salvar o diretório escolhido na lista de diretórios de lista negra
            self.blacklist_directories.append(blacklist_directory)
            messagebox.showinfo("Concluído", "Diretório escolhido para a blacklist com sucesso.")

    def start_schedule_loop(self):
        # Loop para verificar e executar os agendamentos
        self.scan_blacklist_directories()  # Executar imediatamente antes de entrar no loop
        # Agendar a próxima execução após 5 minutos (300.000 milissegundos)
        self.after(300000, self.start_schedule_loop)

    def scan_blacklist_directories(self):
        sensitive_files = []
        for directory in self.blacklist_directories:
            results = {}
            process_directory(directory, results)
            for path, data in results.items():
                if data:  # Verifica se há informações sensíveis encontradas
                    path = os.path.normpath(path)
                    sensitive_files.append(path)

        if sensitive_files:
            sensitive_message = "\n".join(f"Informações sensíveis encontradas em: {path}" for path in sensitive_files)
            messagebox.showinfo("Aviso", sensitive_message)

    def informacoes_outlook(self):
        def limpar_nome(nome):
            nome_limpo = re.sub(r'[\\/*?:"<>|]', '_', nome)
            return nome_limpo

        def get_body(message):
            if message.is_multipart():
                for part in message.walk():
                    ctype = part.get_content_type()
                    cdispo = str(part.get('Content-Disposition'))

                    if ctype == 'text/plain' and 'attachment' not in cdispo:
                        return part.get_payload(decode=True).decode('utf-8')
                    elif ctype == 'text/html':
                        return part.get_payload(decode=True).decode('utf-8')
            else:
                return message.get_payload(decode=True).decode('utf-8')

        root = Tk()
        root.title("Login no Outlook")

        # Função para lidar com o clique no botão de login
        def handle_login():
            username = username_entry.get()
            password = password_entry.get()
            if username and password:
                root.destroy()
                try:
                    Objeto_conexao = imaplib.IMAP4_SSL("imap.outlook.com")
                    Objeto_conexao.login(username, password)
                    Objeto_conexao.select(mailbox='inbox', readonly=True)
                    resposta, idDosEmails = Objeto_conexao.search(None, 'All')

                    for num in idDosEmails[0].split():
                        resultados, dados = Objeto_conexao.fetch(num, '(RFC822)')
                        texto_do_email = dados[0][1].decode('utf-8')
                        mensagem_email = email.message_from_string(texto_do_email)

                        titulo_email = mensagem_email['Subject']
                        nome_pasta = limpar_nome(titulo_email)
                        pasta_anexos = os.path.join("Anexos_email", nome_pasta)
                        if not os.path.exists(pasta_anexos):
                            os.makedirs(pasta_anexos)

                        corpo_email = get_body(mensagem_email)
                        with open(os.path.join(pasta_anexos, f"{nome_pasta}_conteudo.txt"), 'w',
                                  encoding='utf-8') as txt_file:
                            txt_file.write(corpo_email)

                        for part in mensagem_email.walk():
                            if part.get_content_maintype() == 'multipart' or part.get('Content-Disposition') is None:
                                continue
                            fileName = part.get_filename()
                            if fileName:
                                caminho_arquivo = os.path.join(pasta_anexos, fileName)
                                with open(caminho_arquivo, 'wb') as arquivo:
                                    arquivo.write(part.get_payload(decode=True))
                                print(f"Anexo '{fileName}' salvo em '{pasta_anexos}'")

                    # Mostrar a mensagem de sucesso após a conclusão do download dos e-mails e anexos
                    show_success_message()

                except Exception as e:
                    print("Ocorreu um erro:", e)

                finally:
                    if 'Objeto_conexao' in locals():
                        Objeto_conexao.logout()
            else:
                messagebox.showerror("Erro", "Por favor, insira o login e a senha.")

        def handle_delete_folder():
            try:
                shutil.rmtree("Anexos_email")
                messagebox.showinfo("Sucesso", "Pasta 'Anexos_email' excluída com sucesso!")
                root.destroy()
            except Exception as e:
                messagebox.showerror("Erro", f"Ocorreu um erro ao excluir a pasta: {e}")

        def show_success_message():
            messagebox.showinfo("Sucesso",
                                f"Pasta 'Anexos_email' criada com sucesso em:\n{os.path.abspath('Anexos_email')}")

        # Labels e campos de entrada para login e senha
        Label(root, text="Login:").grid(row=0, column=0, padx=5, pady=5)
        username_entry = Entry(root)
        username_entry.grid(row=0, column=1, padx=5, pady=5)

        Label(root, text="Senha:").grid(row=1, column=0, padx=5, pady=5)
        password_entry = Entry(root, show="*")
        password_entry.grid(row=1, column=1, padx=5, pady=5)

        # Botão de login
        login_button = Button(root, text="Login", command=handle_login)
        login_button.grid(row=2, column=0, columnspan=2, padx=5, pady=5)

        # Botão para excluir a pasta "Anexos_email"
        delete_folder_button = Button(root, text="Excluir pasta Anexos_email", command=handle_delete_folder)
        delete_folder_button.grid(row=3, column=0, columnspan=2, padx=5, pady=5)

        # Verifica se a pasta existe e mostra a mensagem de sucesso
        def check_and_show_message():
            pasta_anexos = "Anexos_email"
            if not os.path.exists(pasta_anexos):
                os.makedirs(pasta_anexos)

        show_success_message("E-mails baixados com sucesso!")
def process_directory(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Verifica se o arquivo é uma imagem, PDF, TXT ou DOCX e chama a função correspondente
            if filename.endswith(('.jpg', '.png', '.bmp')):
                image_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_image(image_path, results)
            elif filename.endswith('.pdf'):
                pdf_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pdf(pdf_path, results)
            elif filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)
            elif filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)
            elif filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)
            elif filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)

if __name__ == "__main__":
    app = MeuApp()
    app.mainloop()
