# Antes de tudo instale no terminal > pip install PyMuPDF
# Antes de tudo instale no terminal > pip install google-cloud-vision
# Antes de tudo instale no terminal > pip install python-docx
# Antes de tudo instale no terminal > pip install python-pptx
# Antes de tudo instale no terminal > pip install openpyxl
# Antes de tudo instale no terminal > pip install Pillow
# Antes de tudo instale no terminal > pip install schedule
# Antes de tudo instale no terminal > pip install customtkinter
# Antes de tudo instale no terminal > pip install pandas
# Antes de tudo instale no terminal > pip install matplotlib

import io
import imaplib #biblioteca para se conectar na caixa de e-mail
import email #decodificar partes do e-mail

import threading
import os  # funções para manipular caminhos de arquivos
import re  # ajuda a usar padrões de busca do scan
import shutil  # Copia e/ou move os arquivos
import fitz  # PymuPDF
from google.cloud import vision
from docx import Document  # Para lidar com arquivos DOCX
from pptx import Presentation  # powerpoint
import openpyxl  # Para lidar com arquivos XLSX
import tkinter as tk
from tkinter import Tk, Label, Entry, Button, messagebox
# imports de design
from tkinter import filedialog
from tkinter import messagebox
from tkinter import ttk
from tkinter import Canvas, Entry, Text, Button, PhotoImage
from pathlib import Path
import sys
from customtkinter import *
from PIL import Image, ImageTk
import time  # contador
import schedule  # contador
import customtkinter as ctk
from tkinter import simpledialog
import json  # transforamar o save em um arquivo JSON
import base64  # para a imagem do design virar base64
import io  # load do save
import csv  # relatório
from datetime import datetime
from openpyxl import Workbook
import pandas as pd
import matplotlib.pyplot as plt

#--------------------------------------------------------------------------#

religioes = [
    'Cristianismo',
    'Cristão',
    'Cristã',
    'Islamismo',
    'Islâmico',
    'Islâmica',
    'Hinduísmo',
    'Hindu',
    'Hinduísta',
    'Budismo',
    'Budista',
    'Sikhismo',
    'Sikh',
    'Judaísmo',
    'Judeu',
    'Judaica',
    'Bahaí',
    'Bahaísta',
    'Jainismo',
    'Jainista',
    'Espiritismo',
    'Espírita',
    'Ateísmo',
    'Ateu',
    'Ateia'
]

cores_etnias = [
    'Branco',
    'Negro',
    'Pardo',
    'Indígena',
    'Amarelo',
    'Asiático',
    'Outro / Não Declarado',
]


def extract_info_by_pattern(pattern, text, info_type, results):
    matches = re.findall(pattern, text, flags=re.IGNORECASE)
    if matches:
        results.extend([(info_type, match) for match in matches])
    return results


def format_rg(rg):
    # Remove pontos e hífens do RG e adiciona o hífen no formato desejado
    rg_limpo = re.sub(r'[^\d]', '', rg)
    return f'{rg_limpo[:-1]}.{rg_limpo[-1]}'


# Modifica a função extract_sensitive_info_from_xlsx para incluir formatação de RG
def extract_sensitive_info_from_xlsx(xlsx_path, results):
    wb = openpyxl.load_workbook(xlsx_path)

    sensitive_info = []

    # Itera sobre todas as folhas no arquivo Excel
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]

        # Itera sobre todas as células na folha
        for row in sheet.iter_rows(min_row=1, max_col=sheet.max_column, max_row=sheet.max_row, values_only=True):
            for cell_value in row:
                if cell_value:
                    # Aplica expressões regulares para encontrar informações sensíveis
                    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})|\d{8,9}-\d{1,2}', cell_value)
                    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', str(cell_value))
                    matches_email = re.findall(r'\S+@\S+', str(cell_value))
                    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', str(cell_value))
                    matches_genero = re.findall(
                        r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', cell_value)
                    valid_telefones = []

                    for telefone in matches_telefone:
                        numero_limpo = re.sub(r'[^\d]', '', telefone)
                        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                            valid_telefones.append(telefone)

                    for rg in matches_rg:
                        rg_formatado = format_rg(rg)
                        rg_in_cpf = any(rg_formatado in cpf for cpf in matches_cpf)
                        if not rg_in_cpf:
                            sensitive_info.append(('RG', rg_formatado))

                    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
                    sensitive_info.extend([('Email', email) for email in matches_email])
                    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
                    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
                    # Extrai informações sobre religiões
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', cell_value,
                                                             'Religião', sensitive_info)
                    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', cell_value,
                                                             'Cor/Etnia', sensitive_info)
    if sensitive_info:
        results[xlsx_path] = results.get(xlsx_path, [])
        results[xlsx_path].extend(sensitive_info)

    return results


def process_directory_with_xlsx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos XLSX
            if filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


def extract_sensitive_info_from_pptx(pptx_path_or_text, results):
    if os.path.isfile(pptx_path_or_text):  # Verifica se é um caminho de arquivo
        presentation = Presentation(pptx_path_or_text)
        text = "\n".join(
            [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

    else:
        text = pptx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)
    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pptx_path_or_text] = results.get(pptx_path_or_text, [])
        results[pptx_path_or_text].extend(sensitive_info)

    return results


# Modifica a função process_directory para incluir arquivos PPTX
def process_directory_with_pptx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos PPTX
            if filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)


# Função para extrair informações sensíveis de um arquivo PDF
def extract_sensitive_info_from_pdf(pdf_path, results):
    # Inicializa o documento PDF
    doc = fitz.open(pdf_path)

    sensitive_info = []

    # Itera sobre as páginas do documento PDF
    for page_number in range(doc.page_count):
        page = doc[page_number]

        # Extrai texto da página
        text = page.get_text()

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
        matches_email = re.findall(r'\S+@\S+', text)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

        valid_telefones = []

        # Filtra e formata números de telefone válidos
        for telefone in matches_telefone:
            numero_limpo = re.sub(r'[^\d]', '', telefone)
            if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                valid_telefones.append(telefone)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
            if not rg_in_cpf:
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
        # Extrai informações sobre religiões
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião',
                                                 sensitive_info)
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                                 sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[pdf_path] = results.get(pdf_path, [])
        results[pdf_path].extend(sensitive_info)

    return results


# Função para extrair informações sensíveis de uma imagem
def extract_sensitive_info_from_image(image_path, results):
    # Inicializa o cliente Google Cloud Vision
    client = vision.ImageAnnotatorClient()

    # Lê o conteúdo da imagem
    with open(image_path, 'rb') as image_file:
        content = image_file.read()

    image = vision.Image(content=content)

    # Envia a imagem para análise de texto
    response = client.text_detection(image=image)
    # Extrai texto identificado na imagem
    texts = response.text_annotations

    response = client.face_detection(image=image)
    faces = response.face_annotations

    sensitive_info = []

    # Verifica se há rostos na imagem
    if faces:
        sensitive_info.append(('Rosto', 'Rosto encontrado'))

    # Itera sobre os textos identificados
    for text in texts:
        text = text.description

        # Aplica expressões regulares para encontrar informações sensíveis
        matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
        matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2})', text)
        matches_email = re.findall(r'\S+@\S+', text)
        matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
        matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)
        valid_telefones = []

        # Filtra e formata números de telefone válidos
        for telefone in matches_telefone:
            numero_limpo = re.sub(r'[^\d]', '', telefone)
            if len(numero_limpo) == 11 or len(numero_limpo) == 12:
                valid_telefones.append(telefone)

        # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
        for rg in matches_rg:
            rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
            if not rg_in_cpf:
                sensitive_info.append(('RG', rg))

        # Adiciona informações sensíveis encontradas
        sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
        sensitive_info.extend([('Email', email) for email in matches_email])
        sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
        sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
        # Extrai informações sobre religiões
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião',
                                                 sensitive_info)
        sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                                 sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[image_path] = results.get(image_path, [])
        results[image_path].extend(sensitive_info)

    return results


# Função para extrair informações sensíveis de um arquivo TXT

def extract_sensitive_info_from_txt(txt_path, results):
    with open(txt_path, 'rb') as txt_file:
        text = txt_file.read().decode('utf-8', errors='ignore')  # Decodificar explicitamente como UTF-8

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-\d{1,2}|\d{8}-\d{1,2}|\d{7,9}-\d{1,2}', text)
    matches_cpf = re.findall(r'\b(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})\b', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in matches_telefone])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[txt_path] = results.get(txt_path, [])
        results[txt_path].extend(sensitive_info)

    return results


# Função para processar um diretório e seus subdiretórios, incluindo arquivos TXT
def process_directory_with_txt(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos TXT
            if filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)


# Função para extrair informações sensíveis de um arquivo DOCX
def extract_sensitive_info_from_docx(docx_path_or_text, results):
    if os.path.isfile(docx_path_or_text):  # Verifica se é um caminho de arquivo
        with open(docx_path_or_text, 'rb') as docx_file:
            doc = Document(docx_file)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    else:
        text = docx_path_or_text

    sensitive_info = []

    # Aplica expressões regulares para encontrar informações sensíveis
    matches_rg = re.findall(r'\d{2}\.\d{3}\.\d{3}-(?:\d{1,2})', text)
    matches_cpf = re.findall(r'(\d{3}\.\d{3}\.\d{3}-\d{2}|\d{9}/\d{2}|\d{11})', text)
    matches_email = re.findall(r'\S+@\S+', text)
    matches_telefone = re.findall(r'\(\d{2}\)\d{5}-\d{4}|\(\d{2}\)\d{4,5}-\d{4}', text)
    matches_genero = re.findall(r'\b(Masculino|masculino|M|Homem|homem|Feminino|feminino|Mulher|mulher|F)\b', text)

    valid_telefones = []

    # Filtra e formata números de telefone válidos
    for telefone in matches_telefone:
        numero_limpo = re.sub(r'[^\d]', '', telefone)
        if len(numero_limpo) == 11 or len(numero_limpo) == 12:
            valid_telefones.append(telefone)

    # Verifica se um RG não está contido em um CPF e adiciona à lista de informações sensíveis
    for rg in matches_rg:
        rg_in_cpf = any(rg in cpf for cpf in matches_cpf)
        if not rg_in_cpf:
            sensitive_info.append(('RG', rg))

    # Adiciona informações sensíveis encontradas
    sensitive_info.extend([('CPF', cpf) for cpf in matches_cpf])
    sensitive_info.extend([('Email', email) for email in matches_email])
    sensitive_info.extend([('Telefone', telefone) for telefone in valid_telefones])
    sensitive_info.extend([('Gênero', genero) for genero in matches_genero])
    # Extrai informações sobre religiões
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(religioes) + r'\b', text, 'Religião', sensitive_info)
    sensitive_info = extract_info_by_pattern(r'\b' + r'\b|\b'.join(cores_etnias) + r'\b', text, 'Cor/Etnia',
                                             sensitive_info)

    # Adiciona as informações sensíveis extraídas ao dicionário de resultados
    if sensitive_info:
        results[docx_path_or_text] = results.get(docx_path_or_text, [])
        results[docx_path_or_text].extend(sensitive_info)

    return results


# Função para processar um diretório e seus subdiretórios, incluindo arquivos DOCX
def process_directory_with_docx(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Adiciona suporte para arquivos DOCX
            if filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)


def process_directory(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Verifica se o arquivo é uma imagem, PDF, TXT ou DOCX e chama a função correspondente
            if filename.endswith(('.jpg', '.png', '.bmp')):
                image_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_image(image_path, results)
            elif filename.endswith('.pdf'):
                pdf_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pdf(pdf_path, results)
            elif filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)
            elif filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)
            elif filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)
            elif filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


# Define o caminho do diretório a ser processado
caminho_diretorio = ''
results = {}

# Chama a função para processar o diretório e seus subdiretórios, incluindo arquivos PDF, TXT e DOCX
process_directory(caminho_diretorio, results)

# Remove duplicatas nas informações sensíveis
for path, data in results.items():
    results[path] = list(set(data))

# Exibe as informações sensíveis encontradas
for path, data in results.items():
    path = os.path.normpath(path)
    print(f"Informações sensíveis encontradas em: {path}")


class MeuApp(ctk.CTk):
    def __init__(self):
        super().__init__()
        self.title("Eye Guardian")
        self.geometry("1200x700")
        self.resizable(False, False)

        imagem_base64 = (
            """iVBORw0KGgoAAAANSUhEUgAABLAAAAK8CAMAAADI7ba8AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDozQzJGMDFBRkU3MDMxMUVFODNFMkQ5QjU0M0ExOUQwQiIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDozQzJGMDFCMEU3MDMxMUVFODNFMkQ5QjU0M0ExOUQwQiI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOjNDMkYwMUFERTcwMzExRUU4M0UyRDlCNTQzQTE5RDBCIiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOjNDMkYwMUFFRTcwMzExRUU4M0UyRDlCNTQzQTE5RDBCIi8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+E5WbxQAAAIdQTFRFKSkpVlRULi4uWFZWV1VVrampWVdXq6mprqqqrKioraurqqioVVNTr6urqaenVFJSWllZsKysNDQ0qKWlMDAwMzMznpubNTU1rKqqW1paU1FRr62tqKOjXFtbnpmZoZycnJmZqqWlMTExpqOjoZ6em5iYo6CgMjIypKGhSEZGq6enQkJC8u3teVpJIgAADlJJREFUeNrs3Yty2za3gNGgQNDWpiXbSZqm6f1+4fs/X7E30z5AzJmQxlpnTidNbUkW/H8DgiD14gUAAAAAAAAAAABP8vj27ePj45s3b/74+fXr159N4a9aWiu11lZb6bWUVsafahn/VnqrS88/LHWJ/zi+qtTexzeML6v345/ju/v4D+Pvxn8oy9LGP2uPvy0lH6jFo/Txx/Gn8TUlv3n8y/iyFt8Vf1XKu3iW0vPxlni4Eg/XxyPVfJLxXeOxxl+V+Lp49PEf2of/G98XL2+pNb97/Dm+7D6fcjxGu374scbTL+M19nbNFzyeLdTxN9srrfEErfVv8lmW+xY/WL7+7anjx6r58sv2TOP/49WMR7/Gi6zxMsZblm9nL/fj779r8dO1eKT4oWvNd7rXa/xM8Xflmi+kx8vb3sO2PXrPV1/i6eOB4+nj57zv154DdY3vbtccp3jCeOZ8s5aSDzi+asnXNF5wu49X2XJA4hHjmfMV5ODGF8arHH/Ib+jb25dvc7yesr1Z+R6NHy5/B+KrYmxiDJcYnPia+i6GP973eNeG8WDLtW9fF29Bbff50+f7UPJ3Iv/c8qfLHyhf3BiA8fPEAG+/ifk12+/AeIJ8SfmVPd+Y7XdjjGRZYnTydyefucXLzvex5duaf4w3peSry5fU8x2KYa91eyvLNkr5lsSA5y9dPHrfXkKNn6XFr3GPJ81fyfiu7Str/PLnO1f//92McYyhqT3/5zCebOlL/FvZfovz3YkHy6/t8RXxg2y/CCXfnpavOv6RLyV/F0p78c+EvliBMxIsQLCO7HPjDoIlWIBgCRYwa7BujTsI1lm8NO4gWM4SAoK1sy+NOwiWQ0JAsMywADMsQLDswwIEyyEhCJZtDYBgWcMCBMsMCwTLDAsQrEO5Me4gWGZYgGAJFrBadAcE69jeG3cQLDMsQLBcSwis7ukOCJZgAYLlbg0gWPZhAYIlWIBgOUsIgmUfFiBY7tYACJZggWBZdAcEy6I7IFgW3UGwHBICgiVYgGC5+BkEywwLECx3awAEywwLBEuwAMGycRQQLMECwbKtARAsl+YAguX2MiBYDgkBwXJICEweLPuwQLAECxAs+7AAa1iAYDkkBARLsECwBAsQrINyAz8QLGcJAcFyaQ4wbbBsawDBsoYFCJazhIBDQkCwHBICgiVYIFizuRh3ECyL7oBgCRawujQHEKxj8zFfIFgOCQHBEixAsADBsugOCJYZFgiWi58BwTLDAgTradwiGQTL3RoAwXJICJhhAYLlWkJAsBwSgmCZYQGCZYYFCJZrCUGwnCUEBMshISBYFt1BsFxLCAiWQ0JAsCy6g2CZYQGCZR8WIFhmWCBYZliAYJlhAYIlWCBYz8DFuINgmWEBguXSHGD1yc+AYAkWIFj2YYFgWcMCBMshISBYZlggWO6HBQiWYAGCZac7CJZtDYBgHYCP+QLBMsMCBMs+LMAMCxAswQIEyz4sECxrWIBg2dYACJZggWBZdAcEy6I7IFiCBYLlbg2AYNnWAAiWGRYI1rPknu4gWBbdAcGyDwuwhgUIlhkWIFiCBYJlHxYgWGZYgGCZYYFgOUsICJaNo4BgmWGBYFnDAgTLDAsQLMECwRIsQLDcDwsQLDvdQbCek4txB8GyrQEQLIvugGABguUsISBYLn4GwRIsQLDswwIEy6I7CJZDQkCwDuWVcQfBsoYFCJZDQmC16A4I1rHdGXcQLGtYgGDt7MG4g2BZwwIEy1lCYHXHUUCwLLoDgmUNCwTLDAsQLBtHAcF6Cvd0B8GyhgUIljUsQLAAwbJxFBAsMywQrNncGHcQLMECBMvdGgBrWIBgCRYgWIIFgiVYgGC5WwMgWM4SgmCZYQGC5VpCQLDcDwsEyyEhIFhH4KPqQbAcEgKC5W4NwLTBsg8LBMsMCxAs+7AAMyxAsJwlBATLDAsEywwLECyX5gCC5SwhCJZDQkCwjsU93UGwHBICgmVbA+CQEBAswQIEy+1lQLDMsADBsg8LECzbGkCwBAsQLIeEgGCZYYFg2dYACJZLcwDBsoYFgmXjKCBYFt0BwXJICILlkBAQLDMsQLDswwLBckgICJZgAYJlDQsEy6U5gGA5JAQE6z93xh0E6yxujTsIlkV3QLCsYQFmWIBgCRYgWO6HBYIlWIBg2dYACJYZFgiW+2EBguUsISBYggWCJViAYLk0BxAsMywQLGcJAcESLECwbBwFwRIsQLAECxAswQLBOj+fmgOCZYYFCJZtDcC0wbLTHQTrNNxxFATLDAsQLIvuwLTBcnsZECxrWIBgWcMCrGEBgmWGBQjWTm6MOwjWWVyMOwiWNSxAsOzDAgQLECxnCQHBEiwQLIeEgGAdkzuOgmAJFiBY9mEB1rAAwXItISBYZlggWIIFCJaNo4BgmWGBYFl0BwTrWO6MOwiWjaOAYAkWYA0LECzbGgDB2oePqgfBsoYFCJZgAYIFCJZgAYIlWCBYtjUAgnVQF+MOguWQEBAswQJWN/ADBOvYfJAqCJazhIBgmWEBggUI1tG5pzsIljUsQLAEC5g2WPZhgWAJFiBYDgkBwQIEyyEhIFg2joJgzcYnP4NgWcMCBMsaFmCGBQjW0d0adxAsZwkBwdqZD6EAwRIsQLAsugPTBsu2BhAsi+6AYLk0B7CGBQiWGRYgWGZYIFjOEgKCZeMoIFhmWCBY1rAAwXKWEBAs98MCwXJICAiWYAGC5SwhCJYZFiBYbi8DCJZggWC5NAcQLBtHAcEywwLBcpYQECyX5gCC5SwhCJZgAYJl0R0QLMECwXLxMyBYZljA9MGyDwsEyyEhIFhmWIA1LECwbBwFBMshIQiWQ0JAsNytARAs2xpAsCZxZ9xBsKxhAYK1Mx9CAYJlHxYgWGZYwLTBsq0BBMuiOyBYe7sYdxAs+7AAwbLoDqwuzQEEy+1lAMGyrQEEyyEhIFiCBQiWjaMgWC5+BgTLDAsQrI/w3riDYNnWAAiWYAHWsADBcpYQECwzLBAsO90BwXK3BkCwBAsEyyEhIFgW3QHB+hg+NQcE6zQejDsIljUsQLAEC1gtugOCZYYFCJbby4BgOSQEBMvtZQDBEiwQLIvugGAdyY1xB8EywwIEywwLWC26A4JlHxYgWGZYIFhzcU93ECzBAgTLtgZg2mBZdAfBOg33dAfBsoYFCJZgAYIFCJazhIBg7cQ93UGwHBICguWQEJg2WO6HBYIlWIBgWcMCpg3WxbiDYFl0BwTL3RoAMyxAsJwlBATLISEIlm0NgGBZwwIES7BAsNxeBhCsQ/HJzyBYzhICguWQEJg2WDaOgmA5SwgIlhkWYA0LECyX5gCCJVggWPZhAYLlLCEgWC7NAcEywwIEy6I7IFg2joJgOUsICJad7oBgWcMCwRIsQLCsYQGCZYYFgmWnOyBYZliAYAkWCJad7oBgufgZECwzLBAsa1iAYNnWAAiWQ0IQLIvugGA5JATMsIw7CJZgAYLl9jKAYAGCJViAYNnpDoJlhgUI1jH5mC8QLIeEgGDZ6Q5MG6w74w6C5W4NgGC5NAewhgUIlrOEgGA5JATBsnEUECxnCQHBepIH4w6CZYYFCJazhIBgAYJlWwMgWHa6g2C5WwMgWNawAMGyhgWC5dIcQLAcEgKC9TF8zBcIlm0NgGBZwwLMsADBsq0BECyL7iBYtjUAgnVQF+MOgmUNCxAswQKsYQGCdXA+hAIEyyEhIFhmWMC0wXJpDgiWQ0JAsAQLmDZYPjUHBMvFz4BgmWEB0wbLTncQrNNwi2QQLBtHAcFySAhMGyxnCUGwrGEBgmWGBVjDAgTL3RoAwbKGBYLlkBAQLIvugGBZwwLBMsMCBMsMCxAswQLBcpYQEKxD8CEUIFiCBQiWs4TAtMGy6A6CJViAYDlLCEwbLB9CAYLlkBAQLNsaADMsQLCsYQGC5ZAQBGs27ukOguXSHECwrGEB0wbLWUIQLJfmAIK1twfjDoLlLCEgWIIFrM4SAoJlHxYgWLY1gGCZYQGCJViAYNk4CoLlLCEgWGZYgGA5SwiC5ZAQECxnCQHBsoYFgiVYgGB9Oj41BwTLojsgWIIFCBYgWBbdAcHayZ1xB8E6i4txB8GyrQEQLBc/A6uzhIBgHZtPfgbBckgICJZgAdMGy1lCECwzLECwBAuYNlh2uoNgufgZEKy9+RAKECzBAgRLsIBpg2UfFgiWRXdAsMywgGmDZeMoCJYZFiBYggUIFiBY1rAAwbKtAQTLISEgWC7NAQTLDAsEaxK3xh0EywwLEKyduUUyCJZtDYBgCRYwbbBujDsIlhkWIFiCBax2ugOCJViAYDkkBMESLECwbGsABMsMCwRrCm6RDIJlhgUIlvthAdMGy/2wQLDcDwsQLDvdAcECBMsaFiBYZlggWLY1AIJl4yggWC7NAcGyhgUIlmABgmXRHQTreXMDPxAsMyxAsMywgGmDZR8WCJZrCQHBMsMCpg2WfVggWIIFCJZrCYFpg2VbAwiWYAGCZQ0LmDZYr4w7CJZ9WIBgCRaw2tYACNax3Rl3ECxnCQHBsoYFCBYgWBbdAcHahxv4gWAJFiBYLn4GzLAAwRIsQLBsHAXBsg8LECyHhIBgOSQEwTLDAgTLGhYgWIIFgmWnOyBY7tYACJZDQhAsMyxAsMywAMESLBAsG0cBwfoEbo07CJZ9WIBgOSQEVtsaAMESLECwHBKCYM3la+MOgiVYgGDZhwXMGqxfjDsI1ln8btxBsM7iK+MOgnUWvxl3EKyz+NG4g2CdxffGHQTLGhYgWDv7wbiDYJ3Fr8YdBOss/jTuIFhn8a1xB8E6i1fGHQTrLC7GHQTrLNxxFATrNH4y7iBYZ+EWySBYp+EGfiBY1rAAwdrbS+MOguWQEBAsMyxg1mD9bdzhlP4VYAAKAvY/mmCmwgAAAABJRU5ErkJggg==                

            """)
        # Decode the base64 string and create an image object
        image_data = base64.b64decode(imagem_base64)
        image = Image.open(io.BytesIO(image_data))

        # Resize the image using Image.Resampling.LANCZOS
        resized_image = image.resize((1200, 700), Image.Resampling.LANCZOS)

        # Convert the resized image to a format Tkinter can use
        self.imagem = ImageTk.PhotoImage(resized_image)

        # Create a full-screen label to display the image
        self.tela_cheia = tk.Label(self, image=self.imagem)
        self.tela_cheia.place(x=0, y=0, relwidth=1, relheight=1)

        # Variáveis
        self.directory_path = tk.StringVar()
        self.key_path = tk.StringVar()
        self.sensitive_files = []
        self.blacklist_directories = []
        self.scan_reports = []

        # Carregar configurações salvas, se houver
        self.load_settings()

        # Criar e exibir widgets
        self.create_widgets()

        # Definindo o modo de aparência inicial
        ctk.set_appearance_mode("light")

        # Agendar a execução do escaneamento a cada 5 minutos
        schedule.every(2).minutes.do(self.scan_blacklist_directories)

        # Iniciar o loop de agendamento
        self.after(100, self.start_schedule_loop)

    def create_widgets(self):

        interrogacao_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDo4NDlGQjJGMkUxN0IxMUVFOEIzNDg5MzZGNzFENDk4NiIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDo4NDlGQjJGM0UxN0IxMUVFOEIzNDg5MzZGNzFENDk4NiI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOjg0OUZCMkYwRTE3QjExRUU4QjM0ODkzNkY3MUQ0OTg2IiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOjg0OUZCMkYxRTE3QjExRUU4QjM0ODkzNkY3MUQ0OTg2Ii8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+et2fAwAAAmRQTFRF+vr6+Pj4/f39AAAA+/v7AAAAAQEBAgICAAAA/Pz8GxsbKioqHBwc9vb2GRkZCQkJlpaWQEBA+fn5n5+fxsbGenp6iYmJjIyMcHBwCgoKeHh4YGBgq6urXFxc5eXlPz8/2dnZvr6+MDAwWlpaKSkpGBgYBwcH9PT0DAwM7+/vAQEBAgIC6+vrAwMDvb29ra2t4+Pj09PT1tbW3d3d6Ojoy8vLzs7OnZ2d7u7um5ubQ0NDqKio4eHhqampcXFx39/f9fX1AQEBLS0t1NTUHx8fIiIiBAQEtbW17OzsoaGhjY2NZWVli4uLWVlZf39/eXl5JiYmY2NjLi4uSUlJv7+/RkZGR0dHSkpKioqKZ2dnBQUFS0tLvLy8EBAQdXV1wMDAT09PoKCgUlJSExMT8fHx0NDQHh4eMjIyDg4ODQ0N9/f37+/vV1dXJycnAAAAODg45ubm0tLS6enp1dXV6urqmZmZh4eH0dHRycnJMzMzysrK5OTkAwMDaWlpu7u78vLy4uLi8/PznJycs7OzsbGxKCgoNDQ0x8fHFBQUo6OjTU1NOzs7Pj4+uLi4OTk5mJiYhoaGc3Nzr6+vBgYGfHx8j4+PiIiI2NjYPT09PDw8s7Oz19fXp6enYmJipaWl3t7eUFBQX19f5+fnzc3Nubm5jo6OZmZmRUVFoqKiNjY2TExM4ODgQkJCVVVV7e3tg4ODe3t7QUFBFhYWFxcXYWFhERERsLCwMTExurq6ISEh3NzcEhISTk5Od3d3bGxsNTU1U1NTkJCQzMzMxcXFXV1dUVFR8PDwrKysfn5+lJSUAAAA////yBmZWwAAAMx0Uk5T//////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////8AeUA2nwAAAnBJREFUeNrsldVTG1EUxpNMXhoXkhBCSIBAmFkGSYDgdHF3d7fA4B6g7dSou1B3d3e3b/+pbsJCX3bJ9K2d4b7c73znN3Nn7zn3LI/6u8Xb4P8tXuD8ota0Bu381qxNit7snbfVXrRkg15k74vg/AAvfMDCVwJosihPi7vfAMTDDOd6/B1NCCAzTsyOBnQ1nr0nR32iw8nN20uaoPcNtEtXUhmKBkB8MJ6T/xSGnuNRFCWcOj/h30Unp4szMdPIxZ85BdmPTRQVdfPtUfTVRObS3jWCrBtj52PqSL3mMkXxtSJklz6B5AF91NRH0qXhs/KOENww0wWI2wuJokU9hJkcHtUfmQ5TNCtvIJ8p6E0Yl9ru4FGqX4QkRUVRyfN4tY2Vv43WEfcuHL9Pu/GTrqFd9MdYE3aHXmflb6F95E90JAS+7qsX+HxAOSsvx/uYtSBWibs+HrWlGmGsPKBblcLoATJov8CjrcEAB39CwMjRZiIrgrlFaRkHT8JvkZHmcpiYXhb6HEApKz+IYDMjh4vPLTGSv5zlesrKG7Fj+6obvubzU2RyX1ZeFwo/Jvy+EDjMuOMFENWy8s8HIfbv98jPiUHaXI9aVMswMMbKvz4UhppYj8zpRfVKLa72INOfo5/Di9BRcpLuAepKkSTVnZHajCD84rneS6AJj+fdVeXbY+fceNIjFxqmOd9XZ4IFhLJwNtx9Rqct8p0Y2BPB5+RVaXlVoUDVMUXly7ZCgwjouNQiWGc+CFRzkyISJCnX15P0OOlucwq9zKu0vJ8Fyj56XrnSDfmVW73ON2tyXIT6cMW+ijLd8oWN/8X/zv8WYAAK9OXhTVh+wgAAAABJRU5ErkJggg==        """

        # -----------------------------------------------------------------------------------------------------

        sair_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDowQzZGMjBENUUxN0MxMUVFQjczNUVDRjcyNUI2NDZFQyIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDowQzZGMjBENkUxN0MxMUVFQjczNUVDRjcyNUI2NDZFQyI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOjBDNkYyMEQzRTE3QzExRUVCNzM1RUNGNzI1QjY0NkVDIiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOjBDNkYyMEQ0RTE3QzExRUVCNzM1RUNGNzI1QjY0NkVDIi8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+h3ijewAAARdQTFRFAwMDycnJFhYWAAAAcXFxkJCQ/Pz8ZmZm5OTkAwMD29vbHh4eHx8flpaWjo6OkZGRi4uLOzs70tLSRkZGenp6ubm5bGxsMTExeXl5CgoKhYWF+vr6tLS0+/v78vLy2NjYIyMjpaWlDAwMwMDAtbW1np6e9fX1vr6+iIiIfHx8BAQEs7Oz/f39paWlkpKSAAAAPz8/zMzMioqKrKysHBwce3t7fn5+LCwsysrKGxsb6+vrODg4Pj4+dnZ2l5eX7e3tra2t8PDwCwsLExMTYmJiKioqMDAwEhISFBQUfX19x8fHYGBgjY2NT09PVlZW8/Pz2dnZCQkJt7e3AgIC9/f3ERER7+/vVVVV+fn5oqKiSUlJAAAA////StUKfAAAAF10Uk5T//////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////8A4a+dkAAAAOlJREFUeNpiiCENMIyqH7LqVXydOGFAigj1AibRcCBGhHrBaE4mGOAjQj1ntB8jDHAQoT46WoQk/0ZHs6DI6fhbhZCinkPXSJ0E9WzGrtzejCSod+CJlmQ3lCXaPY6q4sGB7BZsxKqPUYu0jlZi9yKgXlqGFQL07CTEo6NDXZTxq5cTjkYBkqakqXcWINI9/LwaBtHR3BJCxPo3RjQqmkkzgNjwkQ2y0XI34yI6vrTNoz08uYiPX11LBVsuEtKbPD+vECnpM8I+jI2a+YUzOpyk/EhqfvdxI608EWXnIam8Gi3/h596gAADAA890kW1rojaAAAAAElFTkSuQmCC
        """

        # --------------------------------------------------------------------------------------------------------

        save_base64 = """
        iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDpFRkNENjcyQ0UxODIxMUVFOTRBQkI2NTg5MUQyRjYwNSIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDpFRkNENjcyREUxODIxMUVFOTRBQkI2NTg5MUQyRjYwNSI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOkVGQ0Q2NzJBRTE4MjExRUU5NEFCQjY1ODkxRDJGNjA1IiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOkVGQ0Q2NzJCRTE4MjExRUU5NEFCQjY1ODkxRDJGNjA1Ii8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+P6ZW+gAAAL1QTFRF5OTkv7+/j4+PUVFRYGBgcXFxlZWVT09PBQUFmZmZt7e3Hh4e4uLi2NjY39/f4eHhvr6+vb29xsbGTk5OHBwcbGxsubm5k5OT1tbWcHBwq6urwsLCXFxc7+/vjY2NSUlJqKiolJSUkJCQpaWljo6Ou7u7sbGxgICAc3NzdHR0oaGh+/v7sbGxr6+vAAAAExMTJycn0NDQKCgokZGRAAAAFxcX9vb2dnZ2PDw8AgICd3d3BgYGMDAwAAAA////9ryNxQAAAD90Uk5T//////////////////////////////////////////////////////////////////////////////////8AjiZ8FwAAAOxJREFUeNrs1ckOgjAQgGFwRQsIggso7kQFjFvc7bz/Y7lGU9MpctPof2v5DgOZBAmSJf39p3jLJoTYHfZ5xyZdxK/opcHjPF/MyjA4Xy35/nj1hcdZpjSCwuWu+5YvU1q9eZrQ797yJqWTm1cEPnpe7DdriOJ8xS/eM+p12a/EeV6/5g+YD0xDfskwA3wfQt7Kh7hXYdTOM7VHoOJeB7eWY6q5oONe482j4b4ELSfL5LSgJPLDZoapORT5pPPo0O+lmHp90fuqMLbSTNZY9D093jwe7hsSpwbiCbpvfL9VFGSdp///xbf6kwADAACA3CQneTTlAAAAAElFTkSuQmCC        """

        # --------------------------------------------------------------------------------------------------------
        # Decodifique as strings base64 em dados binários
        image_data_sair = base64.b64decode(sair_base64)
        image_data_interrogacao = base64.b64decode(interrogacao_base64)
        image_data_save = base64.b64decode(save_base64)

        # Crie as imagens a partir dos dados decodificados
        imagem_botao_sair = PhotoImage(data=image_data_sair)
        imagem_botao_interrogacao = PhotoImage(data=image_data_interrogacao)
        imagem_botao_save = PhotoImage(data=image_data_save)

        # Crie os botões com as imagens como conteúdo
        botao_sair = ctk.CTkButton(master=self, image=imagem_botao_sair, text="", text_color="", fg_color="transparent",
                                   hover_color="#f7f3f2", font=("Times New Roman", 17), command=self.close_program)
        botao_sair.grid(row=3, column=1, padx=10, pady=10, sticky="e")

        botao_interrogacao = ctk.CTkButton(master=self, image=imagem_botao_interrogacao, text="", text_color="",
                                           fg_color="transparent", hover_color="#f7f3f2", font=("Times New Roman", 17),
                                           command=self.tutorial)
        botao_interrogacao.grid(row=0, column=1, padx=10, pady=10, sticky="e")

        botao_save = ctk.CTkButton(master=self, image=imagem_botao_save, text="", text_color="", fg_color="transparent",
                                   hover_color="#f7f3f2", font=("Times New Roman", 17), command=self.save_settings)
        botao_save.grid(row=3, column=1, padx=10, pady=10, sticky="w")

        # --------------------------------------------------------------------------------------------------------
        # String base64 da imagem do logo do grupo
        base64_image = """iVBORw0KGgoAAAANSUhEUgAAAC8AAAAuCAMAAACPpbA7AAAAGXRFWHRTb2Z0d2FyZQBBZG9iZSBJbWFnZVJlYWR5ccllPAAAAydpVFh0WE1MOmNvbS5hZG9iZS54bXAAAAAAADw/eHBhY2tldCBiZWdpbj0i77u/IiBpZD0iVzVNME1wQ2VoaUh6cmVTek5UY3prYzlkIj8+IDx4OnhtcG1ldGEgeG1sbnM6eD0iYWRvYmU6bnM6bWV0YS8iIHg6eG1wdGs9IkFkb2JlIFhNUCBDb3JlIDkuMS1jMDAxIDc5LjE0NjI4OTk3NzcsIDIwMjMvMDYvMjUtMjM6NTc6MTQgICAgICAgICI+IDxyZGY6UkRGIHhtbG5zOnJkZj0iaHR0cDovL3d3dy53My5vcmcvMTk5OS8wMi8yMi1yZGYtc3ludGF4LW5zIyI+IDxyZGY6RGVzY3JpcHRpb24gcmRmOmFib3V0PSIiIHhtbG5zOnhtcD0iaHR0cDovL25zLmFkb2JlLmNvbS94YXAvMS4wLyIgeG1sbnM6eG1wTU09Imh0dHA6Ly9ucy5hZG9iZS5jb20veGFwLzEuMC9tbS8iIHhtbG5zOnN0UmVmPSJodHRwOi8vbnMuYWRvYmUuY29tL3hhcC8xLjAvc1R5cGUvUmVzb3VyY2VSZWYjIiB4bXA6Q3JlYXRvclRvb2w9IkFkb2JlIFBob3Rvc2hvcCAyNS4zIChXaW5kb3dzKSIgeG1wTU06SW5zdGFuY2VJRD0ieG1wLmlpZDpEN0YxMzFBQkUxODUxMUVFQTBERjhGRUJDMjYxOERBNSIgeG1wTU06RG9jdW1lbnRJRD0ieG1wLmRpZDpEN0YxMzFBQ0UxODUxMUVFQTBERjhGRUJDMjYxOERBNSI+IDx4bXBNTTpEZXJpdmVkRnJvbSBzdFJlZjppbnN0YW5jZUlEPSJ4bXAuaWlkOkQ3RjEzMUE5RTE4NTExRUVBMERGOEZFQkMyNjE4REE1IiBzdFJlZjpkb2N1bWVudElEPSJ4bXAuZGlkOkQ3RjEzMUFBRTE4NTExRUVBMERGOEZFQkMyNjE4REE1Ii8+IDwvcmRmOkRlc2NyaXB0aW9uPiA8L3JkZjpSREY+IDwveDp4bXBtZXRhPiA8P3hwYWNrZXQgZW5kPSJyIj8+S3PHGwAAAwBQTFRFAPb1AMPFaf7+U/79APf1AN3dh/38xf79ALu+sdvdg/Lylf/+AMrKANbXMf/+APTxAPn2ANDRAL/DU9na8P//AOPh1f7+AMjJAMLEAOvpANraJP/+4P/+9f//APX0VcTG7P//3v7+xuLjAMHDAMjLAMzOAPv4AMDCAP//ANXWAP//ANbWAPTxAPr5AMXGzf/+Xf/+qv//APb2AOXjAOLhAMfIuP7+AO/tn/n5AODfAM/Q7//+APHtANLSAPbz7/f3AO7sAMvNAM/QAPTy6P//ANbWAMzNAPz5of79pNnb+P//AO3rAO/sQ8LEAOrny/r6AMvMAOfmAOjmjf39AOvqAPn35f7+APr5vv//AOblAMXHmv//ef//Qc3OAOjnAPXz8f//ANHRAMjJAMXHANvbAM7Ph///AN/eAMTGAPz7APn5ANfXAOfjAMHDANbW/v///P//////+///x///0f7+AP//AO7sAPj2AM7PAP79ANXUANDRsfn5pO/6svP7AP38/v/+Rd73AN7eyvHxAPz7APLwf//+AP//q/P7sf39s/79ALu+ANHSmv38APfzAPf0AOTjAObkAPXyAPz6APDuAPr6AMPGcPz83+/xAMTH2O3uSMDEpdXXAPPxq/j8wN/g+P39tvz7+f39lerqANLSANXWAPn4fuvsANPULP372f/+yujoy+fpAMnKVvz7z/PzAM3NAODgdv37f/37ANTUYsXIasrLAPr4APr6APz6AOTiAMfJAPz7AOHg/Pz8Af387f39AOHe/v7+/v39RMPGAPz6AODgANrbANzdAN3bAPLwAPLyAN3anv39TPz79vv7gc7PAP39AMLEAMbIAN3dAPv3AN7cAPn3AOfnU8LEAPn5k+j6+vv8gf38NMHC3e3t6e/wD8DCAO/sdcjKAMnKAMnKAPPxANrb8vLzAPn1APf2ANjZ/f7/7Pj4sff8rv79r///APj2sv39sv7+AL7CAPr3AL7AAPj2AOXki//+AMzNAO3rAPj1AO7slPv8////rBwyFgAAAQB0Uk5T////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////////AFP3ByUAAAJgSURBVHjaYvhPGmAYVU+WesVV+SSoL0ow85CRJVq9iOKno0eZ2cOIU684Xe0rg7NOPMvXPWZeBNVb9+5geBQyUVR0QtBfnVvMJQTUi7B9N4+zzF7i5NSz9Erkb7s9LgV41Z/9/uiHVLXjpYykjKdOS4+F9K90waO+QMHd3DjPcVr3Pd4rsd1fVBKbfr/qlMWtXoRNZ2JeQO69ObO45JqVHn9JTjQx36GPW73Yn98pAdMect7J524XKbzveykpO0jHALd6dp3IpTO7lRb8L+ZrTxMpeLE1d8le83ovXOplNz06VpF70/5/BFPbT74Y60JP9adSkuGK+djVF0Wo/T3BkyzsX8z8811HsDR3yWJVvby4WzHWDdjUF4btf/n3eZLjjfB2vuLbjfMM2i+efp20JfB75cJFWM33+nz978mkZOEaQ77outp/7QYflgPND/ouHXYIh/unPzpYkSv/Nt9/809uZosDb8DuV9PPxxU+Zs+Mq5K7Vx/+z97eHnO+0H6rUOIx86MJONWXGEmmBOQemb8PGP4KhQ+2XUrKC3zGl48z/BWmn9q1NEDIl3MGl9z6NRe+JVdb/r3Fjif9uGx/ZLLUMbd7NiT9VF/53fdeFl/65Gb5bZK3IYnx8mVG019LrvnZ7QnDm54V7rI8yjmTV13xq2LFc1FJ5z1iRfjzl0Lw1Ul//UTT00Ujf5+6vp+dcP5ddW6ZQPwUG50/Au7vZYnI70Uiivoe4dfXdZUoyOYTV/4UJHiwKcqSUr55yQ5UeZufT0gNSEF+fn4BiCTW/IL8AmDE5ReM1l8Dqx4gwACHbFzwX2V6UgAAAABJRU5ErkJggg==
        """
        # Decodificar a string de base64 em uma imagem
        image_data = base64.b64decode(base64_image)
        image = Image.open(io.BytesIO(image_data))
        image = ImageTk.PhotoImage(image)

        # Exibir a imagem
        image_label = ctk.CTkLabel(master=self, image=image, text="", height=0)
        image_label.grid(row=3, column=1, padx=0, pady=0, sticky="s")
        # --------------------------------------------------------------------------------------------------------

        frame = ctk.CTkScrollableFrame(master=self, fg_color="#484444", border_color="#962CCA", border_width=2,
                                       height=485)
        frame.grid(row=0, column=0, rowspan=3, padx=10, pady=10)
        frame.grid_columnconfigure(0, weight=1)  # Configura a coluna 0 para expandir horizontalmente

        botao_escanear = ctk.CTkButton(master=self, text="Escanear", text_color="black",
                                       fg_color="#9370DB", width=820,
                                       hover_color="#53DEC9", font=("Times New Roman", 17), command=self.start_scan)
        botao_escanear.grid(row=0, column=1, padx=10, pady=10, sticky="w")

        ctk.CTkButton(master=frame, text="Escolher Diretório", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.choose_directory).grid(row=3, column=0, padx=0, pady=10,
                                                                                 sticky="ew")
        ctk.CTkButton(master=frame, text="Escolher Chave", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.choose_key_file).grid(row=4, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Excluir Arquivos", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.delete_files).grid(row=5, column=0, padx=0, pady=10,
                                                                             sticky="ew")
        ctk.CTkButton(master=frame, text="Mover Arquivos", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.move_files).grid(row=6, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Adicionar Blacklist", text_color="black", fg_color="#53DEC9",
                      font=("Times New Roman", 17),
                      hover_color="#9370DB", command=self.choose_blacklist_directory).grid(row=7, column=0, padx=0,
                                                                                           pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Lista Blacklist", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.show_blacklist).grid(row=8, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Relatório Excel", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.open_report).grid(row=9, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Relatório em Gráfico", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.graphic).grid(row=10, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Outlook", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.informacoes_outlook).grid(row=11, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Filtrar Info", text_color="black", fg_color="#53DEC9", hover_color="#9370DB",
                      font=("Times New Roman", 17),
                      command=self.filtrado).grid(row=1, column=0, padx=0, pady=10, sticky="ew")
        ctk.CTkButton(master=frame, text="Procurar Info", text_color="black", fg_color="#53DEC9",
                      hover_color="#9370DB", font=("Times New Roman", 17),
                      command=self.Escanear_info_especifica).grid(row=2, column=0, padx=0, pady=10, sticky="ew")

        # Quadrado Vazio
        quadrado_vazio = ctk.CTkFrame(master=self, width=900, border_color="#962CCA", border_width=2)
        quadrado_vazio.grid(row=2, column=1, padx=10, pady=(0, 0))
        quadrado_vazio.grid_rowconfigure(0, weight=1)
        quadrado_vazio.grid_columnconfigure(0, weight=1)

        self.output_text = ctk.CTkTextbox(master=quadrado_vazio, wrap=tk.WORD, border_color="#962CCA", border_width=1,
                                          height=540, width=900)
        self.output_text.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")

    def Escanear_info_especifica(self):
        # Cria uma janela modal para entrada de texto
        root = tk.Tk()
        root.withdraw()  # Esconde a janela principal

        # Pergunta ao usuário que tipo de informação ele deseja procurar
        info_to_search = simpledialog.askstring("Busca de Informação",
                                                "Que tipo de informação deseja procurar?")

        messagebox.showinfo("Scan em progresso",
                            "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar. Caso tenha recusado a escolha, ignore esta mensagem)")
        # Se o usuário cancelar a entrada, info_to_search será None
        if info_to_search is not None:
            directory_path = self.directory_path.get()
            if directory_path:
                # Limpa o texto antigo
                self.output_text.delete(1.0, tk.END)

                results = {}
                process_directory(directory_path, results)

                sensitive_files = []

                # Variável para armazenar os diretórios onde a informação foi encontrada
                directories_with_info = []

                for path, data in results.items():
                    results[path] = list(set(data))
                    sensitive_files.append(path)

                    # Chamada para a função de detecção de informações sensíveis
                    results = extract_sensitive_info_from_image(path, results)

                    # Verifica se rostos foram detectados e exibe uma mensagem
                    if 'Rosto' in data:
                        self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

                    # Verifica se algum texto contém a informação especificada
                    for info in data:
                        if info_to_search in info[1]:  # info[1] contém o texto detectado
                            directories_with_info.append(path)
                            break

                # Mostra os diretórios onde a informação foi encontrada
                if directories_with_info:
                    self.output_text.insert(tk.END, f"Diretórios com {info_to_search} encontrada:\n")
                    for directory in directories_with_info:
                        self.output_text.insert(tk.END, f"{directory}\n")
                else:
                    self.output_text.insert(tk.END, f"Nenhum diretório com {info_to_search} encontrada.\n")

                self.sensitive_files = sensitive_files

                # Aumenta o tamanho da fonte
                self.output_text.configure(
                    font=(
                        "Times New Roman",
                        20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

                # Adiciona os dados do escaneamento a self.scan_reports
                current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
                for path, data in results.items():
                    self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

                # Gera o relatório
                self.generate_report()
        else:
            print("Busca cancelada pelo usuário.")

    def show_blacklist(self):
        blacklist_window = tk.Toplevel(self)
        blacklist_window.title("Lista Blacklist")
        blacklist_window.geometry("800x600")  # Definindo a geometria da janela

        listbox_frame = tk.Frame(blacklist_window)
        listbox_frame.pack(fill=tk.BOTH, expand=True)  # Faz o frame expandir para preencher toda a janela

        listbox = tk.Listbox(listbox_frame, selectmode=tk.MULTIPLE)
        listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)  # Faz a Listbox preencher todo o frame

        for directory in self.blacklist_directories:
            listbox.insert(tk.END, directory)

        scrollbar = tk.Scrollbar(listbox_frame, orient="vertical", command=listbox.yview)
        scrollbar.pack(side="right", fill="y")
        listbox.config(yscrollcommand=scrollbar.set)

        def remove_selected():
            selected_indices = listbox.curselection()
            selected_directories = [listbox.get(index) for index in selected_indices]
            for directory in selected_directories:
                self.blacklist_directories.remove(directory)
            messagebox.showinfo("Removido", "Os diretórios selecionados foram removidos da blacklist.")
            self.save_settings()

        def clear_blacklist():
            self.blacklist_directories = []
            messagebox.showinfo("Blacklist Esvaziada", "A lista de blacklist foi esvaziada com sucesso.")
            blacklist_window.destroy()  # Fechar a janela após esvaziar a blacklist

        remove_button = tk.Button(blacklist_window, text="Remover Selecionados", command=remove_selected)
        remove_button.pack(padx=10, pady=10)
        esvaziar_button = tk.Button(blacklist_window, text="Esvaziar Blacklist", command=clear_blacklist)
        esvaziar_button.pack(padx=10, pady=10)

    # save
    def load_settings(self):
        if os.path.exists("settings.json") and os.path.getsize("settings.json") > 0:
            with open("settings.json", "r") as f:
                settings = json.load(f)
                self.key_path.set(settings.get("key_path", ""))
                if self.key_path.get():
                    os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = self.key_path.get()
                self.blacklist_directories = settings.get("blacklist_directories", [])
        else:
            messagebox.showwarning("Aviso", "O arquivo de configuração está vazio ou não existe.")

    def save_settings(self):
        settings = {
            "key_path": self.key_path.get(),
            "blacklist_directories": self.blacklist_directories
        }

        with open("settings.json", "w") as f:
            json.dump(settings, f)

        messagebox.showinfo("Salvo", "As configurações foram salvas com sucesso.")

    def tutorial(self):
        popup = tk.Toplevel(self)
        popup.title("Tutorial")
        popup.geometry("800x600")

        scrollbar = tk.Scrollbar(popup)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        tutorial_text = tk.Text(popup, wrap=tk.WORD, yscrollcommand=scrollbar.set, font=("Arial", 12))
        tutorial_text.pack(fill=tk.BOTH, expand=True)

        scrollbar.config(command=tutorial_text.yview)

        # Adicione o conteúdo do tutorial aqui
        tutorial_content = """
        COMO PEGAR A CHAVE JSON (I.A):


 01- Crie uma conta no Google Cloud (cloud.google.com).

02- Adicione a forma que quer que seja feito os pagamentos.

03- Crie um novo Projeto.

04- Após criado, clique nas 3 barrinhas no canto superior esquerdo. Escolha "APIs e serviços" e em seguida "Biblioteca".

05- Pesquise por "Cloud Vision API".

06- Selecione a que se parece com um olho azul.

07- Clique em ativar e recarregue a pagina e verifique se a opção "ativar" mudou para "gerenciar".

08- Clique nas 3 barrinhas no canto superior esquerdo. Escolha "APIs e serviços" e em seguida "Credenciais".

09- Clique em "Criar Credencial" e escolha "Contas de Serviço".

10- Escolha um nome para a conta de serviço e crie.

11- Após isso clique na conta de serviço e vá em "chaves".

12- Clique em "adicionar chave" e crie uma nova chave JSON.

13- Coloque a chave em algum diretório que você irá se lembrar para quando for usar o aplicativo.

        FUNCIONAMENTO DOS BOTÕES:



01- ❓ > Abre o tutorial (onde você está agora).

02- Escanear > Escaneia o diretório escolhido.

03- Filtrar Info > Permite o usuário escolher que tipo de informação ele deseja ver dentre as opções: CPF, Gênero, Religião, RG, Telefone e Rosto (necessário escolher o diretório antes)

04- Procurar Info > Permite o usuário procurar por uma informação especifica dentro do diretório escolhido (necessário escolher o diretório antes)

05- Escolher Diretório > Escolhe o diretório a receber o scan.

06- Escolher Chave > Seleciona a chave JSON da I.A dentro do diretório que foi salvo. (necessário para Scan de imagens)

07- Excluir Arquivos > Exclui todos arquivos sensiveis encontrados no diretório escolhido com exceção dos não sensíveis (necessário ter feito o scan antes).

08- Mover Arquivos > Move todos arquivos sensiveis encontrados para um diretório de sua escolha (necessário ter feito o scan antes).

09- Adicionar Blacklist > Escolhe 1 ou mais diretórios para ser constantemente monitorado em busca de arquivos sensiveis, caso algum arquivo sensível apareça será aberto um pop-up falando que foi encontrado um arquivo sensível e te mostrará o diretório que ele se encontra a cada 5 minutos.

10- Lista Blacklist > Mostra a lista de todos diretórios salvos na blacklist, caso queira remover um ou mais basta clicar nos desejados e então clicar em "Remover Selecionados", caso queira esvaziar toda blacklist, clique em "Esvaziar blacklist".

11- Relatório Excel> Mostra um relatório em uma planilha excel contendo informações de escaneamentos anteriores, dentro do relatório contem a data que foi feito o escaneamento, o horário, que tipo de arquivo sensivel foi encontrado e o diretório onde estava.

12- Relatório em Gráfico > Mostra um relatório no formato de pizza contendo uma porcentagem em cada fatia para mostrar quantos % de cada informação foi encontrada.

13- Outlook > Pedirá seu login e senha para acessar o seu Outlook, após dar as informações o app vai ler e baixar todos e-mails para a pasta "Anexos_email" que vai estar dentro do local onde o executável original se encontra. Dentro da pasta "Anexos_email" terá várias sub-pastas com os nomes referentes aos titulos dos e-mails respectivos que foram baixados. Com essa pasta o usuário pode fazer o que quiser, desde scans até mover ou excluir arquivos. Há também uma opção para excluir a pasta "Anexos_email" que vai deletar por completo todos e-mails baixados.

LEMBRE-SE DE SEMPRE SALVAR AS CONFIGURAÇÕES!!! 
        """

        tutorial_text.tag_configure("bold", font=("Arial", 12, "bold"))
        tutorial_text.insert(tk.END, tutorial_content)

    def choose_key_file(self):
        key_file = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
        if key_file:
            self.key_path.set(key_file)
            os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = key_file
            self.save_settings()  # Salvar as configurações sempre que o caminho da chave for definido

    def choose_directory(self):
        directory = filedialog.askdirectory()
        if directory:
            self.directory_path.set(directory)
            messagebox.showinfo("Concluído", "Diretório escolhido com sucesso.")

    def generate_report(self):
        # Abre ou cria um arquivo para o relatório
        report_file_path = "scan_report.csv"
        with open(report_file_path, mode="w", newline="") as report_file:
            report_writer = csv.writer(report_file)

            # Escreve o cabeçalho do relatório
            report_writer.writerow(["Data", "Horário", "Informação encontrada", "Diretório"])

            # Escreve os dados de cada escaneamento
            for scan_data in self.scan_reports:
                data_scan, _, diretorio, data, _ = scan_data
                for info_type, info in data:
                    data = data_scan.split()[0]  # Apenas a data
                    hora = data_scan.split()[1]  # Apenas a hora
                    report_writer.writerow([data, hora, info_type, diretorio])

        # Obtém o diretório atual do arquivo Python em execução
        current_directory = os.path.dirname(os.path.realpath(__file__))
        excel_file_path = os.path.join(current_directory, "resultado_scan.xlsx")

        # Após gerar o relatório CSV, cria o relatório Excel
        self.create_excel_report(report_file_path, excel_file_path)
        # Exclui o arquivo CSV após a criação do arquivo Excel
        os.remove(report_file_path)

    def open_report(self):
        if os.path.exists("resultado_scan.xlsx"):
            os.system("start resultado_scan.xlsx")
        else:
            messagebox.showwarning("Aviso", "O relatório ainda não foi gerado.")

    def create_excel_report(self, csv_file, excel_file):
        # Abre o arquivo CSV e lê os dados
        with open(csv_file, 'r') as f:
            csv_reader = csv.reader(f)
            data = list(csv_reader)

        # Cria um novo arquivo Excel
        wb = Workbook()
        ws = wb.active

        # Copia os dados do CSV para o Excel
        for row_index, row in enumerate(data):
            for col_index, value in enumerate(row):
                ws.cell(row=row_index + 1, column=col_index + 1, value=value)

        # Salva o arquivo Excel
        wb.save(excel_file)

    def graphic(self):
        if not os.path.exists("resultado_scan.xlsx"):
            messagebox.showwarning("Aviso", "O arquivo 'resultado_scan.xlsx' não existe.")
            return

        # Carregando os dados do arquivo Excel
        data = pd.read_excel("resultado_scan.xlsx")

        if data.empty:
            messagebox.showwarning("Aviso", "O arquivo 'resultado_scan.xlsx' está vazio.")
            return

        # Contamos as ocorrências de informações sensíveis por diretório
        directory_counts = data['Diretório'].value_counts()

        # Identificamos o diretório com a maior quantidade de informações sensíveis
        top_directory = directory_counts.idxmax()
        top_directory_count = directory_counts.max()

        # Preparamos os dados para o gráfico, selecionando apenas as informações do diretório mais crítico
        top_info = data[data['Diretório'] == top_directory]['Informação encontrada'].value_counts()

        # Preparamos o segundo conjunto de dados para o total de informações encontradas
        total_info = data['Informação encontrada'].value_counts()

        # Configuramos a figura e os eixos para os subplots
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 7))

        # Primeiro gráfico de pizza
        ax1.pie(top_info, labels=top_info.index, autopct='%1.1f%%', startangle=90)

        # Quebramos o título do diretório em linhas se necessário
        max_length = 50
        wrapped_directory = '\n'.join(
            [top_directory[i:i + max_length] for i in range(0, len(top_directory), max_length)])
        ax1.set_title(f'Diretório com mais informações sensiveis:\n{wrapped_directory}', fontsize=10)

        # Segundo gráfico de pizza
        ax2.pie(total_info, labels=total_info.index, autopct='%1.1f%%', startangle=90)
        ax2.set_title('Total de Informações Sensíveis Encontradas', fontsize=10)

        # Ajustamos o layout para evitar sobreposição de elementos
        plt.tight_layout()

        plt.show()

    def start_scan(self):
        directory_path = self.directory_path.get()
        if directory_path:
            # Limpa o texto antigo
            self.output_text.delete(1.0, tk.END)

            messagebox.showinfo("Scan em progresso",
                                "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar")

            results = {}
            process_directory(directory_path, results)

            sensitive_files = []

            for path, data in results.items():
                results[path] = list(set(data))
                sensitive_files.append(path)

                # Chamada para a função de detecção de informações sensíveis
                results = extract_sensitive_info_from_image(path, results)

                # Verifica se rostos foram detectados e exibe uma mensagem
                if 'Rosto' in data:
                    self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

            for path, data in results.items():
                path = os.path.normpath(path)
                self.output_text.insert(tk.END, f"Informações sensíveis encontradas em: {path}\n")

                types_found = set(info[0] for info in data)
                for info_type in types_found:
                    if info_type == 'Rosto':
                        self.output_text.insert(tk.END, f"{info_type} encontrado\n")
                    else:
                        self.output_text.insert(tk.END, f"{info_type} encontrado\n")

                self.output_text.insert(tk.END, "\n")

            self.sensitive_files = sensitive_files

            # Aumenta o tamanho da fonte
            self.output_text.configure(
                font=(
                "Times New Roman", 20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

            # Adiciona os dados do escaneamento a self.scan_reports
            current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
            for path, data in results.items():
                self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

            # Gera o relatório
            self.generate_report()

    def filtrado(self):
        # Cria uma janela modal para entrada de texto
        root = tk.Tk()
        root.withdraw()  # Esconde a janela principal

        tipos_arquivo = simpledialog.askstring("Tipos de arquivo",
                                               "Digite os tipos de arquivo que deseja filtrar (separados por vírgula) - Opções: CPF, Gênero, Religião, RG, Telefone, Rosto: ")
        messagebox.showinfo("Scan em progresso",
                            "O scan está prestes a ser realizado. Você provavelmente verá que a aplicação congelará até terminar, porém não se preocupe! O scan estará em andamento, seja paciente. (Clique em OK para começar. Caso tenha recusado a escolha, ignore esta mensagem)")  # Se o usuário cancelar a entrada, tipos_arquivo será None
        if tipos_arquivo is not None:
            # Remove espaços em branco extras e divide os tipos de arquivo
            tipos_arquivo = [tipo.strip().lower() for tipo in tipos_arquivo.split(',')]

            # Inicia o scan apenas se pelo menos um tipo de arquivo for especificado
            if tipos_arquivo:
                directory_path = self.directory_path.get()
                if directory_path:
                    # Limpa o texto antigo
                    self.output_text.delete(1.0, tk.END)

                    results = {}
                    process_directory(directory_path, results)

                    sensitive_files = []

                    for path, data in results.items():
                        results[path] = list(set(data))
                        sensitive_files.append(path)

                        # Chamada para a função de detecção de informações sensíveis
                        results = extract_sensitive_info_from_image(path, results)

                        # Verifica se rostos foram detectados e exibe uma mensagem
                        if 'rosto' in [info[0].lower() for info in data]:
                            self.output_text.insert(tk.END, f"Rosto detectado em: {path}\n")

                    # Filtra os resultados para mostrar apenas os diretórios com informações sensíveis especificadas
                    filtered_directories = []
                    for path, data in results.items():
                        tipos_encontrados = [info[0].lower() for info in data if info[0].lower() in tipos_arquivo]
                        if tipos_encontrados:
                            filtered_directories.append(path)

                    # Mostra os diretórios filtrados na saída
                    if filtered_directories:
                        self.output_text.insert(tk.END, "Diretórios com informações sensíveis filtradas:\n")
                        for directory in filtered_directories:
                            self.output_text.insert(tk.END, f"{directory}\n")
                    else:
                        self.output_text.insert(tk.END,
                                                "Nenhum diretório com informações sensíveis filtradas encontrado.\n")

                    self.sensitive_files = sensitive_files

                    # Aumenta o tamanho da fonte
                    self.output_text.configure(
                        font=(
                            "Times New Roman",
                            20))  # Substitua "Helvetica" pelo nome da fonte desejada e 12 pelo tamanho desejado

                    # Adiciona os dados do escaneamento a self.scan_reports
                    current_time = datetime.now().strftime("%d-%m-%Y %H:%M:%S")
                    for path, data in results.items():
                        self.scan_reports.append([current_time, directory_path, path, data, "Escaneado"])

                    # Gera o relatório
                    self.generate_report()
            else:
                print("Nenhum tipo de arquivo especificado para filtrar.")
        else:
            print("Entrada cancelada pelo usuário.")



    def delete_files(self):
        directory_path = self.directory_path.get()
        if directory_path:
            confirmation = messagebox.askyesno("Confirmação",
                                               "Tem certeza de que deseja excluir todos os arquivos sensíveis no diretório?")
            if confirmation:
                for sensitive_file in self.sensitive_files:
                    if os.path.isfile(sensitive_file):
                        os.remove(sensitive_file)

                messagebox.showinfo("Concluído", "Todos os arquivos sensíveis foram excluídos com sucesso.")

    def move_files(self):
        if not self.sensitive_files:
            messagebox.showwarning("Aviso", "Nenhum arquivo sensível foi encontrado.")
            return

        destination_directory = filedialog.askdirectory()
        if destination_directory:
            for sensitive_file in self.sensitive_files:
                try:
                    shutil.move(sensitive_file, destination_directory)
                except Exception as e:
                    self.output_text.insert(tk.END, f"Erro ao mover arquivo '{sensitive_file}': {str(e)}\n")

        messagebox.showinfo("Transferência concluída!", "Todos os arquivos foram transferidos com sucesso!")

    def close_program(self):
        self.destroy()  # Fecha a janela principal do aplicativo

    def choose_blacklist_directory(self):
        blacklist_directory = filedialog.askdirectory()
        if blacklist_directory:
            # Salvar o diretório escolhido na lista de diretórios de lista negra
            self.blacklist_directories.append(blacklist_directory)
            messagebox.showinfo("Concluído", "Diretório escolhido para a blacklist com sucesso.")

    def start_schedule_loop(self):
        # Loop para verificar e executar os agendamentos
        self.scan_blacklist_directories()  # Executar imediatamente antes de entrar no loop
        # Agendar a próxima execução após 5 minutos (300.000 milissegundos)
        self.after(300000, self.start_schedule_loop)

    def scan_blacklist_directories(self):
        sensitive_files = []
        for directory in self.blacklist_directories:
            results = {}
            process_directory(directory, results)
            for path, data in results.items():
                if data:  # Verifica se há informações sensíveis encontradas
                    path = os.path.normpath(path)
                    sensitive_files.append(path)

        if sensitive_files:
            sensitive_message = "\n".join(f"Informações sensíveis encontradas em: {path}" for path in sensitive_files)
            messagebox.showinfo("Aviso", sensitive_message)


    def informacoes_outlook(self):
        def limpar_nome(nome):
            # Substituir caracteres não permitidos por '_'
            nome_limpo = re.sub(r'[\\/*?:"<>|]', '_', nome)
            return nome_limpo

        root = Tk()
        root.title("Login no Outlook")

        username = ""
        password = ""

        # Função para lidar com o clique no botão de login
        def handle_login():
            nonlocal username, password
            username = username_entry.get()
            password = password_entry.get()
            if username and password:
                root.destroy()  # Fecha a janela pop-up após a conclusão
                try:
                    # Conectamos ao servidor do outlook com IMAP
                    Objeto_conexao = imaplib.IMAP4_SSL("imap.outlook.com")

                    # Efetuamos o login
                    Objeto_conexao.login(username, password)

                    # Loopar a caixa de entrada
                    Objeto_conexao.select(mailbox='inbox', readonly=True)
                    resposta, idDosEmails = Objeto_conexao.search(None, 'All')

                    for num in idDosEmails[0].split():
                        resultados, dados = Objeto_conexao.fetch(num, '(RFC822)')
                        texto_do_email = dados[0][1]
                        texto_do_email = texto_do_email.decode('utf-8')
                        texto_do_email = email.message_from_string((texto_do_email))

                        # Limpar o título do e-mail
                        titulo_email = texto_do_email['Subject']
                        nome_pasta = limpar_nome(titulo_email)

                        # Criar a pasta para salvar os anexos se ela ainda não existir
                        pasta_anexos = os.path.join("Anexos_email", nome_pasta)
                        if not os.path.exists(pasta_anexos):
                            os.makedirs(pasta_anexos)

                        for part in texto_do_email.walk():
                            # Se tiver anexo, pegar o nome do anexo
                            if part.get_content_maintype() == 'multipart':
                                continue
                            if part.get('Content-Disposition') is None:
                                continue
                            # Pegando o nome do arquivo em anexo
                            fileName = part.get_filename()
                            # Salvando o arquivo dentro da pasta "Anexos_email"
                            caminho_arquivo = os.path.join(pasta_anexos, fileName)
                            # Escrevendo o binário do anexo no arquivo
                            with open(caminho_arquivo, 'wb') as arquivo:
                                arquivo.write(part.get_payload(decode=True))
                            print(f"Anexo '{fileName}' salvo em '{pasta_anexos}'")

                except Exception as e:
                    print("Ocorreu um erro:", e)

                finally:
                    if 'Objeto_conexao' in locals():
                        # Fecha a conexão com o servidor do Outlook
                        Objeto_conexao.logout()
            else:
                # Se o login ou a senha estiverem em branco, exibe uma mensagem de erro
                messagebox.showerror("Erro", "Por favor, insira o login e a senha.")

        def handle_delete_folder():
            try:
                shutil.rmtree("Anexos_email")  # Excluir a pasta "Anexos_email" e todo o seu conteúdo
                messagebox.showinfo("Sucesso", "Pasta 'Anexos_email' excluída com sucesso!")
                root.destroy()  # Fecha a janela após a exclusão da pasta
            except Exception as e:
                messagebox.showerror("Erro", f"Ocorreu um erro ao excluir a pasta: {e}")

        def show_success_message():
            messagebox.showinfo("Sucesso",
                                f"Pasta 'Anexos_email' criada com sucesso em:\n{os.path.abspath('Anexos_email')}")

        # Labels e campos de entrada para login e senha
        Label(root, text="Login:").grid(row=0, column=0, padx=5, pady=5)
        username_entry = Entry(root)
        username_entry.grid(row=0, column=1, padx=5, pady=5)

        Label(root, text="Senha:").grid(row=1, column=0, padx=5, pady=5)
        password_entry = Entry(root, show="*")
        password_entry.grid(row=1, column=1, padx=5, pady=5)

        # Botão de login
        login_button = Button(root, text="Login", command=lambda: [handle_login(), check_and_show_message()])
        login_button.grid(row=2, column=0, columnspan=2, padx=5, pady=5)

        # Botão para excluir a pasta "Anexos_email"
        delete_folder_button = Button(root, text="Excluir pasta Anexos_email", command=handle_delete_folder)
        delete_folder_button.grid(row=3, column=0, columnspan=2, padx=5, pady=5)

        def check_and_show_message():
            pasta_anexos = "Anexos_email"
            if not os.path.exists(pasta_anexos):
                os.makedirs(pasta_anexos)
            show_success_message()
def process_directory(directory_path, results):
    # Percorre a estrutura de diretórios
    for root, dirs, files in os.walk(directory_path):
        for filename in files:
            # Verifica se o arquivo é uma imagem, PDF, TXT ou DOCX e chama a função correspondente
            if filename.endswith(('.jpg', '.png', '.bmp')):
                image_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_image(image_path, results)
            elif filename.endswith('.pdf'):
                pdf_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pdf(pdf_path, results)
            elif filename.endswith('.txt'):
                txt_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_txt(txt_path, results)
            elif filename.endswith('.docx'):
                docx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_docx(docx_path, results)
            elif filename.endswith('.pptx'):
                pptx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_pptx(pptx_path, results)
            elif filename.endswith('.xlsx'):
                xlsx_path = os.path.join(root, filename)
                results = extract_sensitive_info_from_xlsx(xlsx_path, results)


if __name__ == "__main__":
    app = MeuApp()
    app.mainloop()
